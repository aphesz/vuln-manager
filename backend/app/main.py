# backend/app/main.py

from fastapi import FastAPI, UploadFile, File, HTTPException, Depends, WebSocket, WebSocketDisconnect, Request, Body, Query, Path, status
from fastapi.middleware.cors import CORSMiddleware
from fastapi.middleware.trustedhost import TrustedHostMiddleware
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse, Response
from sqlmodel import SQLModel, Session, create_engine, select
from sqlalchemy import text, case, func, delete
from typing import Optional, List, Dict, Any
from contextlib import asynccontextmanager
from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.util import get_remote_address
from slowapi.errors import RateLimitExceeded
import os
import sys
import logging
import re

# Configure logging to output to stdout
logging.basicConfig(
    level=logging.INFO if os.getenv('ENVIRONMENT') == 'production' else logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[logging.StreamHandler(sys.stdout)]
)
logger = logging.getLogger(__name__)

# Test log message on startup
logger.debug("Logging initialized in main.py")

# Custom imports
from app.models import (
    User,
    UserRead,
    UserCreate,
    UserUpdate,
    UserUpdatePassword,
    Project,
    Finding,
    Instance,
    InstanceRead,
    RiskMapping,
    FindingBase,  # Import FindingBase instead of top-level RiskRating
    ProjectReadWithFindings,
    FindingReadWithInstances,
    Comment,
    CommentRead,
    CommentBase,
    CommentCreate,
    AuditLog,
    AuditLogRead,
    JiraSettings,
    JiraSettingsRead,
    JiraSettingsBase,
    VulnerabilityTemplate,
    VulnerabilityTemplateRead,
    VulnerabilityTemplateVersion,
    VulnerabilityTemplateVersionRead,
    VulnerabilityMatch,
    VulnerabilityMatchRead,
    ImportHistory,
    ImportHistoryRead,
    ProjectMetrics,
    SLAComplianceMetrics,
    ReviewProgressMetrics,
    FindingTrend,
    TopVulnerability,
    Tag,
    TagRead,
    TagCreate,
    TagUpdate,
    FindingTag,
    FindingArtifact,
    FindingArtifactRead,
    ReportTemplate,
    ReportTemplateRead,
    ReportTemplateCreate,
    ReportTemplateUpdate,
)

# Use the RiskRating from FindingBase
RiskRating = FindingBase.RiskRating
from app.parsers import parse_xml_content
from app import scoring
from app import owasp
from app import cwe_top25
from app.reports import generate_report_docx, generate_report_pdf, generate_executive_report_pdf
from app.report_poc_simple import render_docx_simple, render_docx_raw, build_simple_template_docx
from app.report_modular import assemble_report, list_available_modules
from app.auth import (
    get_password_hash,
    verify_password,
    authenticate_user,
    create_access_token,
    create_refresh_token,
    validate_password_strength,
    decode_token,
)
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from datetime import timedelta

# OAuth2 scheme for token authentication
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="/auth/login")
from app.sla import (
    calculate_sla_deadline,
    update_finding_sla,
    get_overdue_findings,
    get_sla_summary,
)
from app.jira import get_jira_client, encrypt_token
from app.timezone_utils import (
    get_utc_now,
    convert_to_user_timezone,
    parse_iso_datetime,
    DEFAULT_TIMEZONE,
    TIMEZONE_CHOICES,
    is_valid_timezone,
)
import json
from datetime import datetime
from pathlib import Path as FilePath
import secrets

# --- Database Setup ---

# Use SQLite as the default when no DATABASE_URL is provided.
# This avoids the need for the `psycopg2` driver in local/dev environments.
# For production you can still set DATABASE_URL to a PostgreSQL DSN (e.g.,
# "postgresql+psycopg2://user:password@db:5432/vuln_db").
DATABASE_URL = os.getenv("DATABASE_URL", "sqlite:///./test.db")
engine = create_engine(DATABASE_URL, echo=False)

def create_db_and_tables():
    """Initializes the database and creates all tables defined in SQLModel."""
    SQLModel.metadata.create_all(engine)

def get_session():
    """Dependency to get a new DB session."""
    with Session(engine) as session:
        yield session

# File uploads base directory (for POC evidence)
# Use /code/uploads inside container (writable by appuser) or override via env var
EVIDENCE_BASE_DIR = FilePath(os.getenv("UPLOAD_DIR", "/code/uploads")).resolve()
EVIDENCE_ARTIFACTS_DIR = EVIDENCE_BASE_DIR / "artifacts"
EVIDENCE_ARTIFACTS_DIR.mkdir(parents=True, exist_ok=True)


# --- Authentication Dependencies ---

def get_current_user(
    token: str = Depends(oauth2_scheme),
    session: Session = Depends(get_session)
) -> User:
    """Get current authenticated user from JWT token."""
    try:
        payload = decode_token(token)
        user_id: str = payload.get("sub")
        token_type: str = payload.get("type")
        
        if user_id is None or token_type != "access":
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="Could not validate credentials",
                headers={"WWW-Authenticate": "Bearer"},
            )
    except Exception:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Could not validate credentials",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    user = session.get(User, int(user_id))
    if user is None:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="User not found",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    return user


def get_current_active_user(current_user: User = Depends(get_current_user)) -> User:
    """Ensure current user is active."""
    if not current_user.is_active:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Inactive user"
        )
    return current_user


def require_role(*allowed_roles: str):
    """Dependency factory to require specific user roles."""
    def role_checker(current_user: User = Depends(get_current_active_user)) -> User:
        if current_user.role not in allowed_roles:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail=f"Insufficient permissions. Required role: {', '.join(allowed_roles)}"
            )
        return current_user
    return role_checker

# --- Input Validation Utilities ---

def validate_string_length(value: str, field_name: str, max_length: int = 500, allow_empty: bool = False) -> str:
    """Validate string length and basic content."""
    if not value and not allow_empty:
        raise HTTPException(status_code=400, detail=f"{field_name} cannot be empty")
    if len(value) > max_length:
        raise HTTPException(status_code=400, detail=f"{field_name} exceeds maximum length of {max_length} characters")
    return value.strip()

def validate_url(url: str, field_name: str = "URL") -> str:
    """Validate URL format."""
    url_pattern = re.compile(
        r'^https?://'  # http:// or https://
        r'(?:(?:[A-Z0-9](?:[A-Z0-9-]{0,61}[A-Z0-9])?\.)+[A-Z]{2,6}\.?|'  # domain
        r'localhost|'  # localhost
        r'\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3})'  # or IP
        r'(?::\d+)?'  # optional port
        r'(?:/?|[/?]\S+)$', re.IGNORECASE)
    if url and not url_pattern.match(url):
        raise HTTPException(status_code=400, detail=f"Invalid {field_name} format")
    return url

def sanitize_html_input(value: str) -> str:
    """Remove potentially dangerous HTML/script tags from user input."""
    if not value:
        return value
    # Remove script tags and their content
    value = re.sub(r'<script[^>]*>.*?</script>', '', value, flags=re.DOTALL | re.IGNORECASE)
    # Remove common XSS vectors
    dangerous_patterns = [
        r'javascript:',
        r'on\w+\s*=',  # onclick=, onerror=, etc.
        r'<iframe',
        r'<embed',
        r'<object',
    ]
    for pattern in dangerous_patterns:
        value = re.sub(pattern, '', value, flags=re.IGNORECASE)
    return value

# --- Rate Limiting Setup ---

# Disable rate limiting during tests
import sys
is_testing = "pytest" in sys.modules

limiter = Limiter(key_func=get_remote_address, enabled=not is_testing)

# --- Request Models (for API endpoints) ---
from pydantic import BaseModel

class ReviewStatusUpdate(BaseModel):
    """Request model for updating finding review status."""
    status: str
    reviewer_name: Optional[str] = None

class JiraSettingsCreate(BaseModel):
    """Request model for creating/updating Jira settings."""
    project_id: int
    jira_url: str
    project_key: str
    api_token: str
    is_active: bool = True

class JiraConnectionTest(BaseModel):
    """Request model for testing Jira connection."""
    jira_url: str
    email: str  # Jira user email for authentication
    api_token: str

class JiraIssueCreate(BaseModel):
    """Request model for creating Jira issue from finding."""
    user: str = "system"

class RemediationUpdate(BaseModel):
    """Request model for updating remediation deadline and owner."""
    remediation_deadline: Optional[str] = None
    remediation_owner: Optional[str] = None
    user: str = "system"

class IssueStatusUpdate(BaseModel):
    """Request model for updating finding issue status."""
    issue_status: str  # "Open", "Partially Closed", or "Closed"
    issue_status_comment: Optional[str] = None
    user: str = "system"

# --- Application Lifespan ---

@asynccontextmanager
async def lifespan(app: FastAPI):
    """Lifespan context manager for FastAPI application startup and shutdown."""
    # Startup
    create_db_and_tables()
    
    # Fix the database using direct connection first
    from app.db import fix_risk_ratings_direct
    fix_risk_ratings_direct()
    
    yield
    
    # Shutdown (if needed in future)

# --- FastAPI App Initialization ---

from app.websocket import WebSocketManager

app = FastAPI(
    title="VulnManager API",
    description="""
## VulnManager - Vulnerability Assessment Management Platform

A comprehensive API for managing security vulnerability assessments, findings, and reports.

### Key Features

- **Project Management**: Create and manage vulnerability assessment projects
- **Finding Tracking**: Track findings with risk ratings, instances, and remediation status
- **Quick Add**: Rapidly create findings from vulnerability templates
- **Vulnerability Repository**: Searchable template library with 100+ common vulnerabilities
- **Report Generation**: Export findings to DOCX, PDF, Excel, CSV
- **CVSS & OWASP Scoring**: Built-in risk calculators
- **SLA Tracking**: Monitor remediation deadlines and compliance
- **Peer Review**: Comment system and review workflow
- **Jira Integration**: Sync findings to Jira issues

### Security

- Rate limiting on all write endpoints
- Input validation and sanitization
- XSS protection
- XXE-safe XML parsing
- Security headers (CSP, X-Frame-Options, etc.)

### Version

Current version: 0.10.0
    """,
    version="0.10.0",
    lifespan=lifespan,
    docs_url="/docs",
    redoc_url="/redoc",
    contact={
        "name": "VulnManager Support",
        "url": "https://github.com/aphesz/vuln-manager",
    },
    license_info={
        "name": "MIT",
    },
)

# Add rate limiter to app state
app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)

# Add CORS middleware to allow frontend requests
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins (can be restricted in production)
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Add security headers middleware (strict HTTP security)
@app.middleware("http")
async def add_security_headers(request: Request, call_next):
    """Add security headers to all responses."""
    response = await call_next(request)
    
    # Prevent clickjacking attacks
    response.headers["X-Frame-Options"] = "DENY"
    
    # Prevent MIME type sniffing
    response.headers["X-Content-Type-Options"] = "nosniff"
    
    # Enable XSS protection (legacy, but good for older browsers)
    response.headers["X-XSS-Protection"] = "1; mode=block"
    
    # Referrer policy - only send referrer for same-origin requests
    response.headers["Referrer-Policy"] = "strict-origin-when-cross-origin"
    
    # Permissions policy - disable sensitive features not in use
    response.headers["Permissions-Policy"] = "geolocation=(), microphone=(), camera=()"
    
    # Content Security Policy - restrict resource loading
    # Allow Swagger UI CDN resources (cdn.jsdelivr.net) for API documentation
    # Allow: same-origin scripts/styles, unsafe-inline for Material-UI and Swagger UI
    response.headers["Content-Security-Policy"] = (
        "default-src 'self'; "
        "script-src 'self' 'unsafe-inline' 'unsafe-eval' https://cdn.jsdelivr.net; "
        "style-src 'self' 'unsafe-inline' https://cdn.jsdelivr.net; "
        "img-src 'self' data: https:; "
        "font-src 'self' https://cdn.jsdelivr.net; "
        "connect-src 'self' ws: wss:; "
        "frame-ancestors 'none'; "
        "base-uri 'self'; "
        "form-action 'self'"
    )
    
    return response

# Initialize WebSocket manager
ws_manager = WebSocketManager()

@app.websocket("/ws/{project_id}")
async def websocket_endpoint(websocket: WebSocket, project_id: int):
    await ws_manager.connect(websocket, project_id)
    try:
        while True:
            await websocket.receive_text()
    except WebSocketDisconnect:
        ws_manager.disconnect(websocket, project_id)

# --- Utility Functions ---

def fix_risk_ratings_in_db(session: Session):
    """
    Fix any incorrectly cased risk ratings in the database.
    """
    logger.info("Checking for risk ratings that need fixing...")
    
    try:
        # First, get a list of all findings
        findings = session.exec(select(Finding)).all()
        logger.info(f"Found {len(findings)} total findings")
        
        fixed_count = 0
        for finding in findings:
            current = finding.risk_rating
            logger.debug(f"Checking finding {finding.id} with risk_rating '{current}'")
            
            # Map the current rating to the correct enum value
            if current == 'MEDIUM':
                finding.risk_rating = RiskRating.Medium.value
                fixed_count += 1
            elif current == 'LOW':
                finding.risk_rating = RiskRating.Low.value
                fixed_count += 1
            elif current == 'HIGH':
                finding.risk_rating = RiskRating.High.value
                fixed_count += 1
            elif current == 'CRITICAL':
                finding.risk_rating = RiskRating.Critical.value
                fixed_count += 1
            elif current in ('INFORMATION', 'INFORMATIONAL'):
                finding.risk_rating = RiskRating.Informational.value
                fixed_count += 1
        
        if fixed_count > 0:
            logger.info(f"Fixed {fixed_count} risk ratings")
            session.commit()
        else:
            logger.info("No risk ratings needed fixing")
            
    except Exception as e:
        logger.error(f"Error fixing risk ratings: {str(e)}", exc_info=True)
        session.rollback()
        raise

def get_risk_rating(raw_rating: str) -> RiskRating:
    """Map raw scanner severity to the :class:`RiskRating` enum.

    The function returns a ``RiskRating`` member, which SQLModel stores as the
    corresponding string value in the PostgreSQL ``risk_rating`` ENUM column.
    
    Valid values are: Critical, High, Medium, Low, Informational
    
    Args:
        raw_rating: The raw risk rating from a scanner (can be numeric or text)
        
    Returns:
        RiskRating: A valid enum member with proper case
    """
    logger.debug(f"Converting raw risk rating: {raw_rating!r}")
    
    # Handle None/empty values
    if not raw_rating:
        logger.warning("Empty risk rating, defaulting to Low")
        return RiskRating.Low
    
    raw_str = str(raw_rating).strip()
    
    # First try exact match (handles already correct values)
    try:
        return RiskRating(raw_str)
    except ValueError:
        pass
    
    # Then try proper case version of the string
    try:
        proper_case = raw_str.lower().capitalize()
        if proper_case == 'Information':  # Special case
            proper_case = 'Informational'
        return RiskRating(proper_case)
    except ValueError:
        pass
    
    # Finally try mapping from known scanner values
    mapping = {
        # Numeric ratings (Nessus)
        '0': RiskRating.Informational,
        '1': RiskRating.Low,
        '2': RiskRating.Medium,
        '3': RiskRating.High,
        '4': RiskRating.Critical,
        
        # Common text variations
        'INFO': RiskRating.Informational,
        'INFORMATIONAL': RiskRating.Informational,
        'INFORMATION': RiskRating.Informational,
        'FALSE POSITIVE': RiskRating.Informational,
    }
    
    result = mapping.get(raw_str.upper(), RiskRating.Low)
    logger.debug(f"Mapped {raw_rating!r} to {result.value}")
    return result

async def process_and_save_issue(session: Session, project_id: int, issue_data: Dict[str, Any]):
    """
    Checks for an existing Finding, creates or updates it, and adds a new Instance.
    Also auto-creates or links to vulnerability templates.
    
    Args:
        session: Database session
        project_id: ID of the project to add findings to
        issue_data: Dictionary containing finding data from scanner
        
    Returns:
        bool: True if successful
        
    Raises:
        ValueError: If required data is missing or invalid
    """
    import re
    
    # Validate required fields
    if not issue_data.get('title'):
        raise ValueError("Finding title is required")
        
    title = issue_data['title']
    raw_risk = issue_data.get('risk_rating_raw')
    description = issue_data.get('description', 'No description provided.')
    remediation = issue_data.get('remediation', 'No remediation provided.')
    
    logger.info(f"Processing finding: {title!r}")
    logger.debug(f"Raw risk rating: {raw_risk!r}")
    
    try:
        standard_risk = get_risk_rating(raw_risk)
        logger.debug(f"Normalized risk rating: {standard_risk.value}")
    except Exception as e:
        logger.error(f"Error converting risk rating {raw_risk!r}: {e}")
        standard_risk = RiskRating.Low
    
    # Extract CWE and CVE IDs from description
    cwe_match = re.search(r'CWE-(\d+)', description, re.IGNORECASE)
    cve_match = re.search(r'CVE-\d{4}-\d{4,}', description, re.IGNORECASE)
    
    cwe_id = f"CWE-{cwe_match.group(1)}" if cwe_match else None
    cve_id = cve_match.group(0) if cve_match else None
    
    if cwe_id:
        logger.debug(f"Extracted CWE: {cwe_id}")
    if cve_id:
        logger.debug(f"Extracted CVE: {cve_id}")
    
    # 1. Check for existing Finding (Deduplication Logic)
    existing_finding = session.exec(
        select(Finding)
        .where(Finding.project_id == project_id)
        .where(Finding.title == title)
    ).first()
    
    # 2. Find or create vulnerability template
    template = None
    template_id = None
    
    # Try to find existing template by title, CWE, or CVE
    if cwe_id or cve_id:
        statement = select(VulnerabilityTemplate)
        if cwe_id:
            statement = statement.where(VulnerabilityTemplate.cwe_id == cwe_id)
        elif cve_id:
            statement = statement.where(VulnerabilityTemplate.cve_id == cve_id)
        template = session.exec(statement).first()
    
    # If no template found by CWE/CVE, try matching by title
    if not template:
        template = session.exec(
            select(VulnerabilityTemplate).where(VulnerabilityTemplate.title == title)
        ).first()
    
    # Create new template if none exists
    if not template:
        logger.info(f"Creating new vulnerability template from scan: {title}")
        template = VulnerabilityTemplate(
            title=title,
            description=description,
            cwe_id=cwe_id,
            cve_id=cve_id,
            default_risk_rating=standard_risk.value,
            vulnerability_type=_extract_vulnerability_type(title),
            remediation_summary=remediation[:500] if remediation else None,  # First 500 chars
            remediation_steps=remediation,
            source=issue_data.get('scanner_type', 'scan'),  # 'burp', 'nessus', etc.
            is_verified=False,  # Auto-created templates need review
            usage_count=0,
            created_at=get_utc_now(),
            updated_at=get_utc_now()
        )
        session.add(template)
        session.flush()  # Get template.id
    
    template_id = template.id
    
    # Update template usage count
    template.usage_count += 1
    template.last_used = get_utc_now()
    session.add(template)
    
    if existing_finding:
        finding = existing_finding
        # Update template_id if not set
        if not finding.template_id:
            finding.template_id = template_id
            session.add(finding)
    else:
        # 3. Create a new Finding with template link and timeline tracking
        
        # Auto-detect OWASP category (v0.8.3)
        cwe_int = owasp.extract_cwe_from_text(description) if description else None
        vulnerability_type = _extract_vulnerability_type(title)
        owasp_category = owasp.detect_owasp_category(
            title=title,
            description=description,
            cwe_id=cwe_int,
            vulnerability_type=vulnerability_type
        )
        
        if owasp_category:
            logger.debug(f"Auto-detected OWASP category: {owasp_category}")
        
        finding = Finding(
            project_id=project_id,
            title=title,
            risk_rating=standard_risk,
            description=description,
            remediation=remediation,
            template_id=template_id,
            discovered_at=get_utc_now(),  # Track when finding was first detected (v0.8.1)
            owasp_category=owasp_category  # Auto-detect OWASP category (v0.8.3)
        )
        session.add(finding)
        session.flush() # Flushes to get the finding.id for the instance

    # 4. Create a new Instance
    instance = Instance(
        finding_id=finding.id,
        location=issue_data.get('location', 'N/A'),
        details=issue_data.get('details', 'N/A'),
        status='New - Unvalidated',  # Default status
        created_at=get_utc_now()  # Set creation timestamp
    )
    session.add(instance)
    
    session.commit()
    
    # Send WebSocket notification
    await ws_manager.broadcast_finding_update(
        project_id,
        finding.id,
        'update' if existing_finding else 'create'
    )
    
    logger.info(f"Finding processed successfully. Template ID: {template_id}, Finding ID: {finding.id}")
    
    return True


def _extract_vulnerability_type(title: str) -> str:
    """
    Extract vulnerability type from title for categorization.
    
    Args:
        title: Finding title
        
    Returns:
        str: Vulnerability type (XSS, SQLi, CSRF, etc.)
    """
    title_lower = title.lower()
    
    # Common vulnerability type patterns
    if 'xss' in title_lower or 'cross-site scripting' in title_lower or 'cross site scripting' in title_lower:
        return 'XSS'
    elif 'sql' in title_lower and 'injection' in title_lower:
        return 'SQLi'
    elif 'csrf' in title_lower or 'cross-site request forgery' in title_lower:
        return 'CSRF'
    elif 'ssrf' in title_lower or 'server-side request forgery' in title_lower:
        return 'SSRF'
    elif 'rce' in title_lower or 'remote code execution' in title_lower:
        return 'RCE'
    elif 'lfi' in title_lower or 'local file inclusion' in title_lower:
        return 'LFI'
    elif 'rfi' in title_lower or 'remote file inclusion' in title_lower:
        return 'RFI'
    elif 'xxe' in title_lower or 'xml external entity' in title_lower:
        return 'XXE'
    elif 'idor' in title_lower or 'insecure direct object' in title_lower:
        return 'IDOR'
    elif 'path traversal' in title_lower or 'directory traversal' in title_lower:
        return 'Path Traversal'
    elif 'authentication' in title_lower or 'auth' in title_lower:
        return 'Authentication'
    elif 'authorization' in title_lower:
        return 'Authorization'
    elif 'session' in title_lower:
        return 'Session Management'
    elif 'encryption' in title_lower or 'crypto' in title_lower:
        return 'Cryptography'
    elif 'information disclosure' in title_lower or 'info leak' in title_lower:
        return 'Information Disclosure'
    elif 'dos' in title_lower or 'denial of service' in title_lower:
        return 'DoS'
    elif 'buffer overflow' in title_lower:
        return 'Buffer Overflow'
    elif 'insecure deserialization' in title_lower:
        return 'Insecure Deserialization'
    else:
        return 'Other'


# --- Health Check ---

@app.get("/health")
def health_check(session: Session = Depends(get_session)):
    """
    Health check endpoint for monitoring and orchestrators.
    Returns 200 if service and database are healthy.
    """
    try:
        # Test database connection with a simple query
        session.exec(text("SELECT 1")).first()
        return {
            "status": "healthy",
            "service": "vuln-manager-api",
            "database": "connected",
            "version": "0.3.0"
        }
    except Exception as e:
        logger.error(f"Health check failed: {e}")
        return JSONResponse(
            status_code=503,
            content={
                "status": "unhealthy",
                "service": "vuln-manager-api",
                "database": "disconnected",
                "error": str(e)
            }
        )

@app.get("/ready")
def readiness_check(session: Session = Depends(get_session)):
    """
    Readiness check for Kubernetes/orchestrators.
    Returns 200 only if service is ready to accept requests.
    """
    try:
        # Test database connectivity
        session.exec(text("SELECT 1")).first()
        # Check that at least tables exist
        session.exec(select(Project)).first()
        return {
            "ready": True,
            "service": "vuln-manager-api"
        }
    except Exception as e:
        logger.warning(f"Readiness check failed: {e}")
        return JSONResponse(
            status_code=503,
            content={
                "ready": False,
                "service": "vuln-manager-api",
                "reason": "Database not initialized"
            }
        )

# --- Authentication & User Management ---

@app.post("/auth/register", response_model=UserRead, status_code=201)
@limiter.limit("5/minute")
def register_user(
    request: Request,
    user_data: UserCreate,
    session: Session = Depends(get_session)
):
    """
    Register a new user account.
    
    - **email**: Valid email address (unique)
    - **username**: Username (unique, 3-30 characters)
    - **password**: Password (min 8 chars, uppercase, lowercase, digit)
    - **full_name**: Optional full name
    
    Rate limited to 5 registrations per minute per IP.
    """
    # Validate password strength
    is_valid, error_msg = validate_password_strength(user_data.password)
    if not is_valid:
        raise HTTPException(status_code=400, detail=error_msg)
    
    # Check if email already exists
    existing_email = session.exec(select(User).where(User.email == user_data.email)).first()
    if existing_email:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    # Check if username already exists
    existing_username = session.exec(select(User).where(User.username == user_data.username)).first()
    if existing_username:
        raise HTTPException(status_code=400, detail="Username already taken")
    
    # Create new user
    hashed_password = get_password_hash(user_data.password)
    db_user = User(
        email=user_data.email,
        username=user_data.username,
        full_name=user_data.full_name,
        hashed_password=hashed_password,
        role=user_data.role if user_data.role else "viewer",  # Default to viewer role
        is_active=True,
        created_at=get_utc_now()
    )
    
    session.add(db_user)
    session.commit()
    session.refresh(db_user)
    
    logger.info(f"New user registered: {db_user.email} (ID: {db_user.id})")
    
    return db_user


@app.post("/auth/login")
@limiter.limit("10/minute")
def login(
    request: Request,
    form_data: OAuth2PasswordRequestForm = Depends(),
    session: Session = Depends(get_session)
):
    """
    Login with email and password to receive JWT tokens.
    
    Returns access token (30min expiry) and refresh token (7 days expiry).
    Rate limited to 10 login attempts per minute per IP.
    """
    user = authenticate_user(session, form_data.username, form_data.password)
    if not user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect email or password",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    if not user.is_active:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Account is disabled"
        )
    
    # Update last login timestamp
    user.last_login = get_utc_now()
    session.add(user)
    session.commit()
    
    # Create JWT tokens
    access_token = create_access_token(data={"sub": str(user.id)})
    refresh_token = create_refresh_token(data={"sub": str(user.id)})
    
    logger.info(f"User logged in: {user.email} (ID: {user.id})")
    
    return {
        "access_token": access_token,
        "refresh_token": refresh_token,
        "token_type": "bearer",
        "user": UserRead.from_orm(user)
    }


@app.post("/auth/refresh")
@limiter.limit("20/minute")
def refresh_access_token(
    request: Request,
    refresh_token: str = Body(..., embed=True),
    session: Session = Depends(get_session)
):
    """
    Refresh an access token using a valid refresh token.
    
    Rate limited to 20 refreshes per minute per IP.
    """
    try:
        payload = decode_token(refresh_token)
        token_type = payload.get("type")
        user_id = payload.get("sub")
        
        if token_type != "refresh" or not user_id:
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="Invalid refresh token"
            )
        
        # Verify user still exists and is active
        user = session.get(User, int(user_id))
        if not user or not user.is_active:
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="User not found or inactive"
            )
        
        # Create new access token
        new_access_token = create_access_token(data={"sub": str(user.id)})
        
        return {
            "access_token": new_access_token,
            "token_type": "bearer"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Token refresh error: {e}")
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Could not validate refresh token"
        )


@app.post("/auth/logout")
def logout(current_user: User = Depends(get_current_active_user)):
    """
    Logout current user.
    
    Note: JWT tokens are stateless. Client should discard tokens.
    This endpoint exists for logging/auditing purposes.
    """
    logger.info(f"User logged out: {current_user.email} (ID: {current_user.id})")
    
    return {"message": "Logged out successfully"}


@app.get("/auth/me", response_model=UserRead)
def get_current_user_profile(current_user: User = Depends(get_current_active_user)):
    """
    Get current authenticated user's profile.
    
    Requires valid JWT access token in Authorization header.
    """
    return current_user


@app.put("/auth/me", response_model=UserRead)
@limiter.limit("10/minute")
def update_current_user_profile(
    request: Request,
    user_update: UserUpdate,
    current_user: User = Depends(get_current_active_user),
    session: Session = Depends(get_session)
):
    """
    Update current user's profile information.
    
    Can update: full_name, avatar_url
    Cannot update: email, username, role, is_active (use admin endpoints)
    
    Rate limited to 10 updates per minute.
    """
    # Update allowed fields
    if user_update.full_name is not None:
        current_user.full_name = user_update.full_name
    
    if user_update.avatar_url is not None:
        current_user.avatar_url = user_update.avatar_url
    
    session.add(current_user)
    session.commit()
    session.refresh(current_user)
    
    logger.info(f"User profile updated: {current_user.email} (ID: {current_user.id})")
    
    return current_user


@app.put("/auth/me/password")
@limiter.limit("5/minute")
def change_password(
    request: Request,
    password_data: UserUpdatePassword,
    current_user: User = Depends(get_current_active_user),
    session: Session = Depends(get_session)
):
    """
    Change current user's password.
    
    Requires: current_password and new_password
    Rate limited to 5 changes per minute.
    """
    # Verify current password
    if not verify_password(password_data.current_password, current_user.hashed_password):
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Current password is incorrect"
        )
    
    # Validate new password strength
    is_valid, error_msg = validate_password_strength(password_data.new_password)
    if not is_valid:
        raise HTTPException(status_code=400, detail=error_msg)
    
    # Update password
    current_user.hashed_password = get_password_hash(password_data.new_password)
    session.add(current_user)
    session.commit()
    
    logger.info(f"Password changed: {current_user.email} (ID: {current_user.id})")
    
    return {"message": "Password updated successfully"}


# Admin-only user management endpoints

@app.get("/users", response_model=List[UserRead])
def list_users(
    skip: int = 0,
    limit: int = 100,
    current_user: User = Depends(require_role("admin")),
    session: Session = Depends(get_session)
):
    """
    List all users (admin only).
    
    Requires admin role.
    """
    statement = select(User).offset(skip).limit(limit)
    users = session.exec(statement).all()
    return users


@app.get("/users/{user_id}", response_model=UserRead)
def get_user(
    user_id: int,
    current_user: User = Depends(require_role("admin")),
    session: Session = Depends(get_session)
):
    """
    Get specific user by ID (admin only).
    
    Requires admin role.
    """
    user = session.get(User, user_id)
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    return user


@app.put("/users/{user_id}", response_model=UserRead)
@limiter.limit("20/minute")
def update_user(
    request: Request,
    user_id: int,
    user_update: UserUpdate,
    current_user: User = Depends(require_role("admin")),
    session: Session = Depends(get_session)
):
    """
    Update user information (admin only).
    
    Can update: email, username, full_name, role, is_active, avatar_url
    Requires admin role.
    """
    user = session.get(User, user_id)
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Update fields
    if user_update.email is not None:
        # Check email uniqueness
        existing = session.exec(select(User).where(User.email == user_update.email, User.id != user_id)).first()
        if existing:
            raise HTTPException(status_code=400, detail="Email already in use")
        user.email = user_update.email
    
    if user_update.username is not None:
        # Check username uniqueness
        existing = session.exec(select(User).where(User.username == user_update.username, User.id != user_id)).first()
        if existing:
            raise HTTPException(status_code=400, detail="Username already in use")
        user.username = user_update.username
    
    if user_update.full_name is not None:
        user.full_name = user_update.full_name
    
    if user_update.role is not None:
        if user_update.role not in ["admin", "analyst", "viewer"]:
            raise HTTPException(status_code=400, detail="Invalid role")
        user.role = user_update.role
    
    if user_update.is_active is not None:
        user.is_active = user_update.is_active
    
    if user_update.avatar_url is not None:
        user.avatar_url = user_update.avatar_url
    
    session.add(user)
    session.commit()
    session.refresh(user)
    
    logger.info(f"User {user.id} updated by admin {current_user.id}")
    
    return user


@app.delete("/users/{user_id}")
@limiter.limit("10/minute")
def delete_user(
    request: Request,
    user_id: int,
    current_user: User = Depends(require_role("admin")),
    session: Session = Depends(get_session)
):
    """
    Delete user (admin only).
    
    Cannot delete yourself. Requires admin role.
    """
    if user_id == current_user.id:
        raise HTTPException(status_code=400, detail="Cannot delete yourself")
    
    user = session.get(User, user_id)
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    session.delete(user)
    session.commit()
    
    logger.info(f"User {user_id} deleted by admin {current_user.id}")
    
    return {"message": "User deleted successfully"}


# --- Endpoint: Project Management ---

@app.post("/projects/", response_model=Project)
@limiter.limit("30/hour")  # Max 30 projects per hour per IP
def create_project(
    request: Request,
    project: Project,
    session: Session = Depends(get_session)
):
    """Creates a new project entry."""
    # Validate and sanitize inputs
    project.name = validate_string_length(sanitize_html_input(project.name), "project name", max_length=200)
    if project.consultant_name:
        project.consultant_name = validate_string_length(
            sanitize_html_input(project.consultant_name),
            "consultant name",
            max_length=100,
            allow_empty=True
        )
    
    session.add(project)
    session.commit()
    session.refresh(project)
    return project

@app.get("/projects/", response_model=List[Project])
def read_projects(session: Session = Depends(get_session)):
    """Returns a list of all projects."""
    projects = session.exec(select(Project)).all()
    return projects

@app.get("/projects/stats/all")
def get_all_projects_with_stats(session: Session = Depends(get_session)):
    """Returns all projects with their statistics (findings count, risk summary)."""
    projects = session.exec(select(Project)).all()
    
    projects_with_stats = []
    for project in projects:
        # Get findings count
        findings = session.exec(
            select(Finding).where(Finding.project_id == project.id)
        ).all()
        
        # Calculate risk summary
        risk_summary = {
            'Critical': 0,
            'High': 0,
            'Medium': 0,
            'Low': 0,
            'Informational': 0
        }
        
        for finding in findings:
            risk_level = finding.risk_rating
            if risk_level in risk_summary:
                risk_summary[risk_level] += 1
        
        # Get last upload date from most recent instance
        last_upload = None
        for finding in findings:
            instances = session.exec(
                select(Instance).where(Instance.finding_id == finding.id)
            ).all()
            for instance in instances:
                if instance.created_at and (not last_upload or instance.created_at > last_upload):
                    last_upload = instance.created_at
        
        projects_with_stats.append({
            'id': project.id,
            'name': project.name,
            'consultant_name': project.consultant_name,
            'is_archived': project.is_archived,
            'archived_at': project.archived_at,
            'total_findings': len(findings),
            'critical_count': risk_summary['Critical'],
            'high_count': risk_summary['High'],
            'medium_count': risk_summary['Medium'],
            'low_count': risk_summary['Low'],
            'last_upload_date': last_upload,
            'risk_summary': risk_summary,
            'critical_high_count': risk_summary['Critical'] + risk_summary['High']
        })
    
    return projects_with_stats

@app.get("/projects/stats")
def get_portfolio_stats(session: Session = Depends(get_session)):
    """Returns aggregated statistics across all projects for portfolio dashboard."""
    # Get all projects
    total_projects = session.exec(
        select(func.count(Project.id))
    ).one()
    
    active_projects = session.exec(
        select(func.count(Project.id)).where(Project.is_archived == False)
    ).one()
    
    archived_projects = total_projects - active_projects
    
    # Get all findings with risk ratings
    findings = session.exec(select(Finding)).all()
    
    # Count by risk level
    critical_findings = sum(1 for f in findings if f.risk_rating == 'Critical')
    high_findings = sum(1 for f in findings if f.risk_rating == 'High')
    medium_findings = sum(1 for f in findings if f.risk_rating == 'Medium')
    low_findings = sum(1 for f in findings if f.risk_rating == 'Low')
    informational_findings = sum(1 for f in findings if f.risk_rating == 'Informational')
    
    total_findings = len(findings)
    avg_findings_per_project = total_findings / total_projects if total_projects > 0 else 0
    
    # Count projects with critical findings
    projects_with_critical = session.exec(
        select(func.count(func.distinct(Finding.project_id)))
        .where(Finding.risk_rating == 'Critical')
    ).one()
    
    # Get most recent upload date
    most_recent_instance = session.exec(
        select(Instance)
        .order_by(Instance.created_at.desc())
        .limit(1)
    ).first()
    
    most_recent_upload = most_recent_instance.created_at if most_recent_instance else None
    
    return {
        'total_projects': total_projects,
        'active_projects': active_projects,
        'archived_projects': archived_projects,
        'total_findings': total_findings,
        'critical_findings': critical_findings,
        'high_findings': high_findings,
        'medium_findings': medium_findings,
        'low_findings': low_findings,
        'informational_findings': informational_findings,
        'avg_findings_per_project': round(avg_findings_per_project, 2),
        'projects_with_critical': projects_with_critical,
        'most_recent_upload': most_recent_upload
    }

@app.get("/projects/{project_id}", response_model=ProjectReadWithFindings) # Using the corrected model name
def read_project(project_id: int, session: Session = Depends(get_session)):
    """Returns details for a specific project, including all findings and instances."""
    project = session.exec(
        select(Project)
        .where(Project.id == project_id)
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Build the response with tags and instances
    findings_with_tags = []
    for finding in project.findings:
        # Load tags for this finding
        finding_tags = session.exec(
            select(Tag)
            .join(FindingTag, Tag.id == FindingTag.tag_id)
            .where(FindingTag.finding_id == finding.id)
            .order_by(Tag.name)
        ).all()
        
        # Convert to dict and add tags, instances, and artifacts
        finding_dict = finding.model_dump()
        finding_dict['tags'] = [TagRead.model_validate(tag) for tag in finding_tags]
        # Explicitly include instances from the relationship
        finding_dict['instances'] = [InstanceRead.model_validate(inst) for inst in finding.instances]
        # Include artifacts (POC evidence)
        try:
            finding_dict['artifacts'] = [FindingArtifactRead.model_validate(a) for a in getattr(finding, 'artifacts', [])]
        except Exception:
            finding_dict['artifacts'] = []
        findings_with_tags.append(FindingReadWithInstances(**finding_dict))
    
    # Build project response
    project_dict = project.model_dump(exclude={'findings'})
    project_dict['findings'] = findings_with_tags
    
    return ProjectReadWithFindings(**project_dict)

@app.put("/projects/{project_id}", response_model=Project)
def update_project(project_id: int, project_update: Project, session: Session = Depends(get_session)):
    """Updates an existing project (name, consultant, archive status)."""
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Update fields
    project.name = project_update.name
    project.consultant_name = project_update.consultant_name
    project.is_archived = project_update.is_archived
    project.archived_at = project_update.archived_at
    
    session.add(project)
    session.commit()
    session.refresh(project)
    return project

@app.delete("/projects/{project_id}", status_code=204)
def delete_project(project_id: int, session: Session = Depends(get_session)):
    """Deletes a project and all associated findings and instances."""
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Delete all instances for this project's findings
    findings = session.exec(
        select(Finding).where(Finding.project_id == project_id)
    ).all()
    
    for finding in findings:
        instances = session.exec(
            select(Instance).where(Instance.finding_id == finding.id)
        ).all()
        for instance in instances:
            session.delete(instance)
        session.delete(finding)
    
    # Delete the project
    session.delete(project)
    session.commit()
    
    logger.info(f"Project {project_id} and all associated findings deleted")

# --- Endpoint: Report Upload and Parsing ---

async def _process_upload(
    project_id: int,
    scanner_type: str,
    xml_content: bytes,
    session: Session
) -> Dict[str, Any]:
    """
    Shared logic for processing uploaded reports.
    """
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    try:
        logger.info(f"Processing report as {scanner_type} for project {project_id}")
        
        # 1. Parse the XML content using the robust utility function
        logger.debug("Parsing XML content...")
        issues_data = parse_xml_content(xml_content, scanner_type)
        logger.debug(f"Found {len(issues_data)} issues in report")
        
    except ValueError as e:
        logger.error(f"Parsing Error: {e}")
        # Handles DTD security block, invalid XML, or unknown scanner type errors
        raise HTTPException(status_code=400, detail=f"Parsing Error: {e}")
    except Exception as e:
        logger.error(f"Unexpected error during parsing: {e}", exc_info=True)
        # Catch all other unexpected errors during file processing
        raise HTTPException(status_code=500, detail=f"Internal Server Error during file processing: {e}")

    # 2. Process and Save all issues to the database (Deduplication occurs here)
    new_instances_count = 0
    for issue in issues_data:
        # Add scanner_type to issue data for template creation
        issue['scanner_type'] = scanner_type
        await process_and_save_issue(session, project_id, issue)
        new_instances_count += 1

    return {
        "message": f"Successfully processed report for Project {project_id}.",
        "new_instances_count": new_instances_count
    }

@app.post("/projects/{project_id}/upload/auto", status_code=201)
@limiter.limit("10/minute")  # Max 10 uploads per minute per IP
async def upload_report_auto(
    request: Request,
    project_id: int,
    file: UploadFile = File(...),
    session: Session = Depends(get_session)
):
    """
    Auto-detects scanner type from XML content and processes accordingly.
    Supports: Burp Suite XML, Nessus v2 XML
    
    Security:
    - File size limit: 10 MiB (enforced in parsers.py)
    - Content-type validation: XML only
    - XXE protection: defusedxml parsing
    """
    # Validate content type
    if file.content_type and not file.content_type.startswith('application/xml') and not file.content_type.startswith('text/xml'):
        raise HTTPException(
            status_code=400,
            detail=f"Invalid file type: {file.content_type}. Expected XML file."
        )
    
    # Validate filename
    if not file.filename or not file.filename.lower().endswith(('.xml', '.nessus')):
        raise HTTPException(
            status_code=400,
            detail="Invalid filename. Must be .xml or .nessus file."
        )
    
    xml_content = await file.read()
    
    # Auto-detect scanner type
    xml_str = xml_content.decode('utf-8', errors='ignore')
    
    if '<NessusClientData' in xml_str:
        scanner_type = 'nessus'
    elif '<issues burpversion' in xml_str.lower() or '<issues' in xml_str.lower():
        scanner_type = 'burp'
    else:
        raise HTTPException(
            status_code=400,
            detail="Could not auto-detect scanner type. Supported formats: Burp Suite XML, Nessus v2 XML"
        )
    
    logger.info(f"Auto-detected scanner type: {scanner_type}")
    return await _process_upload(project_id, scanner_type, xml_content, session)

@app.post("/projects/{project_id}/upload/{scanner_type}", status_code=201)
@limiter.limit("10/minute")  # Max 10 uploads per minute per IP
async def upload_report(
    request: Request,
    project_id: int,
    scanner_type: str,
    file: UploadFile = File(...),
    session: Session = Depends(get_session)
):
    """
    Uploads an XML report, parses it, and saves findings/instances to the database.
    Supported types: 'burp', 'nessus'.
    
    Security:
    - Scanner type whitelist validation
    - Content-type validation: XML only
    - File size limit: 10 MiB (enforced in parsers.py)
    """
    # Whitelist valid scanner types
    valid_scanners = {'burp', 'nessus'}
    scanner_type_lower = scanner_type.lower().strip()
    
    if scanner_type_lower not in valid_scanners:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid scanner type: {scanner_type}. Supported: {', '.join(valid_scanners)}"
        )
    
    # Validate content type
    if file.content_type and not file.content_type.startswith('application/xml') and not file.content_type.startswith('text/xml'):
        raise HTTPException(
            status_code=400,
            detail=f"Invalid file type: {file.content_type}. Expected XML file."
        )
    
    # Validate filename
    if not file.filename or not file.filename.lower().endswith(('.xml', '.nessus')):
        raise HTTPException(
            status_code=400,
            detail="Invalid filename. Must be .xml or .nessus file."
        )
    
    xml_content = await file.read()
    return await _process_upload(project_id, scanner_type_lower, xml_content, session)

@app.get("/projects/{project_id}/risk_summary")
def get_risk_summary(project_id: int, session: Session = Depends(get_session)):
    """
    Returns a summary of findings grouped by risk level.
    Used for risk visualization charts.
    """
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Initialize risk summary
    risk_summary = {
        'Critical': 0,
        'High': 0,
        'Medium': 0,
        'Low': 0,
        'Informational': 0
    }
    
    # Count findings by risk level
    findings = session.exec(
        select(Finding).where(Finding.project_id == project_id)
    ).all()
    
    for finding in findings:
        # finding.risk_rating is already a string (enum value)
        risk_level = finding.risk_rating
        if risk_level in risk_summary:
            risk_summary[risk_level] += 1
    
    return risk_summary

@app.get("/projects/{project_id}/metrics", response_model=ProjectMetrics)
def get_project_metrics(project_id: int, session: Session = Depends(get_session)):
    """
    Returns comprehensive dashboard metrics for a project.
    Includes SLA compliance, review progress, trends, and top vulnerabilities.
    """
    from datetime import timedelta
    from collections import defaultdict
    
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Fetch all findings for this project
    findings = session.exec(
        select(Finding).where(Finding.project_id == project_id)
    ).all()
    
    # 1. SLA Compliance Metrics
    sla_on_track = sum(1 for f in findings if f.sla_status == "On Track")
    sla_at_risk = sum(1 for f in findings if f.sla_status == "At Risk")
    sla_overdue = sum(1 for f in findings if f.sla_status == "Overdue")
    sla_total = len(findings)
    sla_compliance_rate = (sla_on_track / sla_total * 100) if sla_total > 0 else 0.0
    
    sla_compliance = SLAComplianceMetrics(
        on_track=sla_on_track,
        at_risk=sla_at_risk,
        overdue=sla_overdue,
        total=sla_total,
        compliance_rate=round(sla_compliance_rate, 2)
    )
    
    # 2. Review Progress Metrics
    review_pending = sum(1 for f in findings if f.review_status == "Pending")
    review_in_review = sum(1 for f in findings if f.review_status == "In Review")
    review_approved = sum(1 for f in findings if f.review_status == "Approved")
    review_rejected = sum(1 for f in findings if f.review_status == "Rejected")
    review_total = len(findings)
    approval_rate = (review_approved / review_total * 100) if review_total > 0 else 0.0
    
    review_progress = ReviewProgressMetrics(
        pending=review_pending,
        in_review=review_in_review,
        approved=review_approved,
        rejected=review_rejected,
        total=review_total,
        approval_rate=round(approval_rate, 2)
    )
    
    # 3. Finding Trends (last 30 days or since project creation)
    # Group findings by date (we'll use a simplified approach - just show overall counts)
    # For a real implementation, you'd track when findings were created
    now = get_utc_now()
    trends = []
    
    # Generate trend for last 30 days
    for i in range(30, -1, -1):
        date = now - timedelta(days=i)
        date_str = date.strftime("%Y-%m-%d")
        
        # Count findings by status (simplified - in reality you'd filter by creation date)
        open_count = sum(1 for f in findings if f.issue_status in ["Open", "Partially Closed"])
        closed_count = sum(1 for f in findings if f.issue_status == "Closed")
        
        trends.append(FindingTrend(
            date=date_str,
            total_findings=len(findings),
            open_findings=open_count,
            closed_findings=closed_count
        ))
    
    # 4. Top Vulnerabilities (top 5 by instance count)
    finding_instance_counts = []
    for finding in findings:
        instance_count = len(session.exec(
            select(Instance).where(Instance.finding_id == finding.id)
        ).all())
        finding_instance_counts.append((finding, instance_count))
    
    # Sort by instance count descending
    finding_instance_counts.sort(key=lambda x: x[1], reverse=True)
    
    top_vulnerabilities = [
        TopVulnerability(
            title=finding.title,
            risk_rating=finding.risk_rating,
            instance_count=count,
            finding_id=finding.id
        )
        for finding, count in finding_instance_counts[:5]
    ]
    
    # 5. Key Metrics
    total_instances = sum(
        len(session.exec(select(Instance).where(Instance.finding_id == f.id)).all())
        for f in findings
    )
    
    # Calculate average CVSS score (if available)
    cvss_scores = []
    for finding in findings:
        # Check if finding has CVSS data (you'd need to parse from description or store separately)
        # For now, we'll skip this or use a placeholder
        pass
    average_cvss = None  # Would calculate from actual CVSS scores
    
    # Count findings with Jira tickets
    findings_with_jira = sum(1 for f in findings if f.jira_issue_key)
    jira_sync_rate = (findings_with_jira / len(findings) * 100) if findings else 0.0
    
    # Average time to approval (would need timestamps on review_status changes)
    average_time_to_approval = None  # Would calculate from audit logs
    
    return ProjectMetrics(
        sla_compliance=sla_compliance,
        review_progress=review_progress,
        finding_trends=trends,
        top_vulnerabilities=top_vulnerabilities,
        total_findings=len(findings),
        total_instances=total_instances,
        average_cvss_score=average_cvss,
        findings_with_jira=findings_with_jira,
        jira_sync_rate=round(jira_sync_rate, 2),
        average_time_to_approval=average_time_to_approval
    )

# --- Trend Analysis Endpoints (v0.8.1) ---

@app.get("/projects/{project_id}/trends/findings")
@limiter.limit("60/minute")
def get_findings_trend(
    request: Request,
    project_id: int,
    start_date: Optional[str] = Query(None, description="Start date (ISO format: YYYY-MM-DD)"),
    end_date: Optional[str] = Query(None, description="End date (ISO format: YYYY-MM-DD)"),
    granularity: str = Query("daily", regex="^(daily|weekly|monthly)$"),
    session: Session = Depends(get_session)
):
    """
    Get time-series data of finding counts by risk rating.
    
    Returns trend data for visualizing finding discovery over time.
    
    Parameters:
    - start_date: Start of date range (default: 30 days ago)
    - end_date: End of date range (default: now)
    - granularity: Time grouping - 'daily', 'weekly', or 'monthly' (default: daily)
    
    Returns:
    - labels: Array of date labels
    - datasets: Finding counts by risk rating
    - totals: Total counts by risk rating
    """
    from app.trends import get_findings_timeline
    from datetime import datetime
    
    # Parse dates
    start = datetime.fromisoformat(start_date) if start_date else None
    end = datetime.fromisoformat(end_date) if end_date else None
    
    try:
        result = get_findings_timeline(
            session=session,
            project_id=project_id,
            start_date=start,
            end_date=end,
            granularity=granularity  # type: ignore
        )
        return result
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))


@app.get("/projects/{project_id}/trends/remediation")
@limiter.limit("60/minute")
def get_remediation_trend(
    request: Request,
    project_id: int,
    start_date: Optional[str] = Query(None, description="Start date (ISO format: YYYY-MM-DD)"),
    end_date: Optional[str] = Query(None, description="End date (ISO format: YYYY-MM-DD)"),
    granularity: str = Query("daily", regex="^(daily|weekly|monthly)$"),
    session: Session = Depends(get_session)
):
    """
    Get remediation progress metrics over time.
    
    Returns trend data showing open vs. closed findings, remediation velocity,
    and mean time to remediate (MTTR) by risk level.
    
    Parameters:
    - start_date: Start of date range (default: 30 days ago)
    - end_date: End of date range (default: now)
    - granularity: Time grouping - 'daily', 'weekly', or 'monthly' (default: daily)
    
    Returns:
    - labels: Array of date labels
    - open_findings: Count of open findings at each point
    - closed_findings: Count of closed findings at each point
    - remediation_velocity: Findings closed per week
    - mean_time_to_remediate: MTTR in days by risk level
    - by_risk: Current open/closed counts by risk level
    """
    from app.trends import get_remediation_progress
    from datetime import datetime
    
    # Parse dates
    start = datetime.fromisoformat(start_date) if start_date else None
    end = datetime.fromisoformat(end_date) if end_date else None
    
    try:
        result = get_remediation_progress(
            session=session,
            project_id=project_id,
            start_date=start,
            end_date=end,
            granularity=granularity  # type: ignore
        )
        return result
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))


@app.get("/projects/{project_id}/trends/risk-score")
@limiter.limit("60/minute")
def get_risk_score_trend(
    request: Request,
    project_id: int,
    start_date: Optional[str] = Query(None, description="Start date (ISO format: YYYY-MM-DD)"),
    end_date: Optional[str] = Query(None, description="End date (ISO format: YYYY-MM-DD)"),
    granularity: str = Query("daily", regex="^(daily|weekly|monthly)$"),
    session: Session = Depends(get_session)
):
    """
    Get aggregate risk score evolution over time.
    
    Risk score is calculated as weighted sum:
    - Critical = 10 points
    - High = 5 points  
    - Medium = 3 points
    - Low = 1 point
    - Informational = 0 points
    
    Parameters:
    - start_date: Start of date range (default: 30 days ago)
    - end_date: End of date range (default: now)
    - granularity: Time grouping - 'daily', 'weekly', or 'monthly' (default: daily)
    
    Returns:
    - labels: Array of date labels
    - risk_scores: Risk score at each point
    - trend: Overall trend ('improving', 'stable', 'worsening')
    - change_percent: Percentage change from start to end
    - current_score: Most recent risk score
    - start_score: Initial risk score
    """
    from app.trends import get_risk_score_trend
    from datetime import datetime
    
    # Parse dates
    start = datetime.fromisoformat(start_date) if start_date else None
    end = datetime.fromisoformat(end_date) if end_date else None
    
    try:
        result = get_risk_score_trend(
            session=session,
            project_id=project_id,
            start_date=start,
            end_date=end,
            granularity=granularity  # type: ignore
        )
        return result
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))


@app.get("/projects/{project_id}/trends/uploads")
@limiter.limit("60/minute")
def get_upload_history_trend(
    request: Request,
    project_id: int,
    start_date: Optional[str] = Query(None, description="Start date (ISO format: YYYY-MM-DD)"),
    end_date: Optional[str] = Query(None, description="End date (ISO format: YYYY-MM-DD)"),
    session: Session = Depends(get_session)
):
    """
    Get timeline of scan uploads with findings discovered per upload.
    
    Groups findings by discovery time to approximate upload events.
    Findings discovered within 1 hour are considered part of the same upload.
    
    Parameters:
    - start_date: Start of date range (default: 90 days ago)
    - end_date: End of date range (default: now)
    
    Returns:
    - uploads: Array of upload events with finding counts
    - total_uploads: Number of upload events
    - average_findings_per_upload: Average findings discovered per upload
    """
    from app.trends import get_upload_history
    from datetime import datetime
    
    # Parse dates
    start = datetime.fromisoformat(start_date) if start_date else None
    end = datetime.fromisoformat(end_date) if end_date else None
    
    try:
        result = get_upload_history(
            session=session,
            project_id=project_id,
            start_date=start,
            end_date=end
        )
        return result
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))

# --- Predictive Analytics Endpoints (v0.8.5) ---

@app.get("/projects/{project_id}/predict/remediation-time")
@limiter.limit("60/minute")
def predict_remediation_time(
    request: Request,
    project_id: int,
    risk_level: Optional[str] = Query(None, description="Filter by risk level (Critical, High, Medium, Low, Informational)"),
    session: Session = Depends(get_session)
):
    """
    Estimate remediation time based on historical data.
    
    Uses median time-to-remediate from resolved findings to predict
    how long it will take to fix current open findings.
    
    Parameters:
    - risk_level: Optional filter for specific risk level
    
    Returns:
    - List of RemediationTimeEstimate per risk level
    - estimated_days: Median days to remediate
    - confidence_interval_low/high: 95% confidence bounds
    - sample_size: Number of historical findings used
    
    Example:
        GET /api/projects/1/predict/remediation-time
        GET /api/projects/1/predict/remediation-time?risk_level=Critical
    """
    from app.predict import estimate_remediation_time
    from app.models import RemediationTimeEstimate
    
    # Verify project exists
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    try:
        estimates = estimate_remediation_time(session, project_id, risk_level)
        return {
            "project_id": project_id,
            "project_name": project.name,
            "estimates": [e.model_dump() for e in estimates]
        }
    except Exception as e:
        logger.error(f"Error estimating remediation time: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Prediction error: {str(e)}")


@app.get("/projects/{project_id}/predict/risk-forecast")
@limiter.limit("60/minute")
def forecast_project_risk(
    request: Request,
    project_id: int,
    session: Session = Depends(get_session)
):
    """
    Forecast risk score 30/60/90 days ahead using historical trend analysis.
    
    Uses simple linear regression on historical risk scores to predict
    future trajectory. Helps identify if security posture is improving
    or worsening.
    
    Returns:
    - current_risk_score: Current aggregate risk
    - forecast_30_days/60_days/90_days: Predicted risk scores
    - trend: "improving", "stable", or "worsening"
    - confidence: 0.0-1.0 prediction confidence
    
    Example:
        GET /api/projects/1/predict/risk-forecast
        
    Response:
        {
          "current_risk_score": 45.0,
          "forecast_30_days": {"predicted_risk_score": 38.5, ...},
          "trend": "improving",
          "confidence": 0.85
        }
    """
    from app.predict import forecast_risk_score
    
    # Verify project exists
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    try:
        forecast = forecast_risk_score(session, project_id)
        return {
            "project_id": project_id,
            "project_name": project.name,
            **forecast.model_dump()
        }
    except Exception as e:
        logger.error(f"Error forecasting risk: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Prediction error: {str(e)}")


@app.get("/projects/{project_id}/predict/anomalies")
@limiter.limit("60/minute")
def detect_security_anomalies(
    request: Request,
    project_id: int,
    session: Session = Depends(get_session)
):
    """
    Detect anomalies in security metrics using statistical analysis.
    
    Automatically identifies:
    - Sudden spike in findings (possible security incident)
    - Remediation slowdown (team capacity issues)
    - Regression (resolved findings reopened)
    
    Returns:
    - List of detected anomalies with severity and recommendations
    - Empty list if no anomalies detected
    
    Example:
        GET /api/projects/1/predict/anomalies
        
    Response:
        {
          "project_id": 1,
          "anomalies": [
            {
              "anomaly_type": "spike_in_findings",
              "severity": "high",
              "description": "15 findings in last 7 days (avg: 3/week)",
              "recommendation": "Review recent scans for false positives..."
            }
          ]
        }
    """
    from app.predict import detect_anomalies
    
    # Verify project exists
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    try:
        anomalies = detect_anomalies(session, project_id)
        return {
            "project_id": project_id,
            "project_name": project.name,
            "anomaly_count": len(anomalies),
            "anomalies": [a.model_dump() for a in anomalies]
        }
    except Exception as e:
        logger.error(f"Error detecting anomalies: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Prediction error: {str(e)}")


@app.get("/projects/{project_id}/predict/recommendations")
@limiter.limit("60/minute")
def generate_security_recommendations(
    request: Request,
    project_id: int,
    session: Session = Depends(get_session)
):
    """
    Generate actionable recommendations based on project analysis.
    
    Automatically identifies:
    - Quick wins (low-effort, high-impact fixes)
    - Stale findings (>90 days open)
    - SLA at-risk findings
    - Resource allocation issues (overloaded owners)
    
    Returns:
    - List of prioritized recommendations with estimated effort
    - Recommendations sorted by priority (critical → high → medium → low)
    
    Example:
        GET /api/projects/1/predict/recommendations
        
    Response:
        {
          "project_id": 1,
          "recommendation_count": 3,
          "recommendations": [
            {
              "priority": "high",
              "category": "quick_wins",
              "title": "5 low/medium findings with 5+ instances each",
              "description": "Fix root cause to remediate 45 instances",
              "estimated_effort": "1-2 days"
            }
          ]
        }
    """
    from app.predict import generate_recommendations
    
    # Verify project exists
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    try:
        recommendations = generate_recommendations(session, project_id)
        return {
            "project_id": project_id,
            "project_name": project.name,
            "recommendation_count": len(recommendations),
            "recommendations": [r.model_dump() for r in recommendations]
        }
    except Exception as e:
        logger.error(f"Error generating recommendations: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Prediction error: {str(e)}")

# --- Compliance Endpoints ---

@app.get("/projects/{project_id}/compliance/owasp-top-10")
@limiter.limit("60/minute")
def get_owasp_top_10_coverage(
    request: Request,
    project_id: int,
    session: Session = Depends(get_session)
):
    """
    Get OWASP Top 10 2021 coverage statistics for a project.
    
    Returns finding counts per OWASP category, coverage percentage,
    and category details.
    
    Returns:
    - categories: Dict mapping category ID (A01-A10) to finding count and details
    - statistics: Overall coverage statistics
    """
    # Verify project exists
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Get all findings for the project
    findings = session.exec(
        select(Finding).where(Finding.project_id == project_id)
    ).all()
    
    # Count findings per OWASP category
    findings_by_category: Dict[str, int] = {
        category_id: 0 for category_id in owasp.get_owasp_categories().keys()
    }
    
    unmapped_findings = 0
    
    for finding in findings:
        if finding.owasp_category and finding.owasp_category in findings_by_category:
            findings_by_category[finding.owasp_category] += 1
        else:
            unmapped_findings += 1
    
    # Calculate coverage statistics
    statistics = owasp.calculate_coverage_statistics(findings_by_category)
    statistics["unmapped_findings"] = unmapped_findings
    statistics["total_findings_in_project"] = len(findings)
    
    # Build category details
    categories = {}
    for category_id, count in findings_by_category.items():
        category_info = owasp.get_category_description(category_id)
        categories[category_id] = {
            "name": category_info["name"],
            "description": category_info["description"],
            "finding_count": count,
            "has_findings": count > 0
        }
    
    return {
        "categories": categories,
        "statistics": statistics
    }

@app.get("/projects/{project_id}/compliance/cwe-top-25")
@limiter.limit("60/minute")
def get_cwe_top_25_coverage(
    request: Request,
    project_id: int,
    session: Session = Depends(get_session)
):
    """
    Get CWE Top 25 2024 coverage statistics for a project.
    
    Returns finding counts per CWE weakness, coverage percentage,
    and weakness details with severity.
    
    Returns:
    - weaknesses: List of top 10 CWE entries sorted by finding count
    - all_weaknesses: Dict mapping CWE ID to finding count
    - statistics: Overall coverage statistics
    """
    # Verify project exists
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Get all findings for the project
    findings = session.exec(
        select(Finding).where(Finding.project_id == project_id)
    ).all()
    
    # Count findings per CWE ID (extract from description)
    findings_by_cwe: Dict[int, int] = {}
    
    for finding in findings:
        # Try to extract CWE ID from description
        cwe_id = owasp.extract_cwe_from_text(finding.description) if finding.description else None
        
        if cwe_id and cwe_top25.is_in_top_25(cwe_id):
            findings_by_cwe[cwe_id] = findings_by_cwe.get(cwe_id, 0) + 1
    
    # Calculate coverage statistics
    statistics = cwe_top25.calculate_top25_statistics(findings_by_cwe)
    
    # Get top 10 weaknesses by finding count
    top_weaknesses = cwe_top25.get_top_10_by_findings(findings_by_cwe)
    
    # Get all CWE Top 25 with finding counts
    all_weaknesses = {}
    for rank, cwe_data in cwe_top25.get_cwe_top_25().items():
        cwe_id = cwe_data["cwe_id"]
        all_weaknesses[cwe_id] = {
            "rank": cwe_data["rank"],
            "name": cwe_data["name"],
            "severity": cwe_data["severity"],
            "score": cwe_data["score"],
            "finding_count": findings_by_cwe.get(cwe_id, 0),
            "has_findings": cwe_id in findings_by_cwe
        }
    
    return {
        "weaknesses": top_weaknesses,
        "all_weaknesses": all_weaknesses,
        "statistics": statistics
    }

# --- Endpoint: Report Generation ---

@app.get("/projects/{project_id}/report.docx", response_class=FileResponse)
def get_docx_report(project_id: int, session: Session = Depends(get_session)):
    """Generates and returns the assessment report in DOCX format."""
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Generate the report file path
    file_path = f"/tmp/report_{project_id}.docx"
    generate_report_docx(project, file_path)
    
    return FileResponse(
        file_path, 
        media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        filename=f"{project.name.replace(' ', '_')}_Report.docx"
    )

@app.get("/reports/poc/template.docx", response_class=FileResponse)
def download_poc_template():
    """Download a minimal sample DOCX template for the PoC renderer."""
    tmp_path = "/tmp/poc_template.docx"
    try:
        with open(tmp_path, "wb") as f:
            f.write(build_simple_template_docx())
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to build template: {e}")
    return FileResponse(
        tmp_path,
        media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        filename="vulnmanager_poc_template.docx"
    )

@app.post("/projects/{project_id}/report/poc", response_class=FileResponse)
async def generate_poc_report(
    project_id: int,
    apply_style: bool = Query(True, description="Apply left-border styling post-render"),
    donut_size_cm: float | None = Query(None, description="Donut image size in centimeters (optional)"),
    donut_dpi: int | None = Query(None, description="Donut image DPI (optional)"),
    template_file: UploadFile = File(..., description="DOCX template with Jinja2 placeholders"),
    session: Session = Depends(get_session)
):
    """Proof-of-concept templated DOCX report generation using docxtpl + post-processing.

    Upload a .docx file containing a finding loop:
        {% for f in findings %}
        <table> ... {{ f.risk_rating }} {{ f.donut_img }} ... {% endfor %}
        {% endfor %}

    Returns a rendered DOCX with risk-colored donut images and left colored border per finding table.
    """
    if not template_file.filename.lower().endswith('.docx'):
        raise HTTPException(status_code=400, detail="Template must be a .docx file")

    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    # Build read model (reuse existing response model builder logic if available)
    # Minimal reconstruction: load findings + instances
    findings = project.findings  # Relationship already loaded lazily; may trigger queries
    # Wrap project into a pseudo read model object with needed attributes
    class _ProjectReadShim:
        pass
    _proj = _ProjectReadShim()
    _proj.name = project.name
    _proj.findings = findings

    template_bytes = await template_file.read()
    try:
        # Use simpler renderer to avoid corruption; allow disabling styling
        rendered = render_docx_simple(
            template_bytes,
            _proj,
            apply_style=apply_style,
            donut_size_cm=donut_size_cm,
            donut_dpi=donut_dpi,
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Template render failed: {e}")

    out_path = f"/tmp/poc_report_{project_id}.docx"
    with open(out_path, 'wb') as f:
        f.write(rendered)

    return FileResponse(
        out_path,
        media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        filename=f"{project.name.replace(' ', '_')}_POC_Report.docx"
    )

@app.post("/projects/{project_id}/report/poc/raw", response_class=FileResponse)
async def generate_poc_report_raw(
    project_id: int,
    template_file: UploadFile = File(..., description="DOCX template with Jinja2 placeholders"),
    session: Session = Depends(get_session)
):
    """Raw variant: renders template WITHOUT images and WITHOUT any post-processing.

    Use this to diagnose Word open errors. If this opens successfully while the
    styled version fails, the issue is with images or border styling.
    """
    if not template_file.filename.lower().endswith('.docx'):
        raise HTTPException(status_code=400, detail="Template must be a .docx file")

    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    class _ProjectReadShim:
        pass
    shim = _ProjectReadShim()
    shim.name = project.name
    shim.findings = project.findings

    template_bytes = await template_file.read()
    try:
        rendered = render_docx_raw(template_bytes, shim)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Raw template render failed: {e}")

    out_path = f"/tmp/poc_report_raw_{project_id}.docx"
    with open(out_path, 'wb') as f:
        f.write(rendered)

    return FileResponse(
        out_path,
        media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        filename=f"{project.name.replace(' ', '_')}_POC_Report_RAW.docx"
    )

@app.get("/projects/{project_id}/report.pdf", response_class=FileResponse)
def get_pdf_report(project_id: int, session: Session = Depends(get_session)):
    """Generates and returns the assessment report in PDF format."""
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    # Build a risk summary for the project.
    # Use the RiskRating enum values to ensure we cover all possible ratings.
    risk_levels = [r.value for r in RiskRating]
    summary: dict[str, int] = {level: 0 for level in risk_levels}

    findings = session.exec(
        select(Finding.risk_rating).where(Finding.project_id == project_id)
    ).all()

    for risk in findings:
        # ``risk`` is stored as the enum string value.
        if risk in summary:
            summary[risk] += 1

    # Generate a placeholder PDF (or real PDF in the future).
    file_path = f"/tmp/report_{project_id}.pdf"
    generate_report_pdf(project, file_path)

    return FileResponse(
        file_path,
        media_type="application/pdf",
        filename=f"{project.name.replace(' ', '_')}_Report.pdf",
    )


@app.get("/projects/{project_id}/reports/executive", response_class=FileResponse)
def get_executive_report(
    project_id: int,
    date_from: Optional[str] = None,
    date_to: Optional[str] = None,
    include_charts: bool = True,
    company_name: Optional[str] = None,
    custom_header: Optional[str] = None,
    custom_footer: Optional[str] = None,
    session: Session = Depends(get_session)
):
    """
    Generates and returns an executive summary report in PDF format.
    
    This is a polished, stakeholder-friendly report with charts and visualizations,
    designed for executives and non-technical audiences.
    
    Query Parameters:
    - date_from: Optional start date for filtering findings (YYYY-MM-DD)
    - date_to: Optional end date for filtering findings (YYYY-MM-DD)
    - include_charts: Whether to include charts in the report (default: true)
    - company_name: Optional company name for branding
    - custom_header: Optional custom header text
    - custom_footer: Optional custom footer text
    """
    from datetime import date as date_type
    
    # Fetch project with all findings and instances
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Parse date parameters
    parsed_date_from = None
    parsed_date_to = None
    
    if date_from:
        try:
            parsed_date_from = date_type.fromisoformat(date_from)
        except ValueError:
            raise HTTPException(status_code=400, detail="Invalid date_from format. Use YYYY-MM-DD")
    
    if date_to:
        try:
            parsed_date_to = date_type.fromisoformat(date_to)
        except ValueError:
            raise HTTPException(status_code=400, detail="Invalid date_to format. Use YYYY-MM-DD")
    
    # Generate filename
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    file_path = f"/tmp/executive_report_{project_id}_{timestamp}.pdf"
    
    # Generate the executive report
    try:
        generate_executive_report_pdf(
            project=project,
            file_path=file_path,
            date_from=parsed_date_from,
            date_to=parsed_date_to,
            include_charts=include_charts,
            logo_url=None,  # TODO(future): Add logo upload feature in branding settings
            company_name=company_name,
            custom_header=custom_header,
            custom_footer=custom_footer
        )
    except Exception as e:
        print(f"Error generating executive report: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate executive report: {str(e)}")
    
    # Return the PDF file
    return FileResponse(
        file_path,
        media_type="application/pdf",
        filename=f"{project.name.replace(' ', '_')}_Executive_Report.pdf",
    )


@app.post("/projects/{project_id}/reports/from-template", response_class=FileResponse)
def generate_report_from_template(
    project_id: int,
    template_id: int = Body(..., embed=True),
    variables: Optional[dict] = Body(default=None, embed=True),
    session: Session = Depends(get_session)
):
    """
    Generate a report from a template.
    
    POST body:
    {
        "template_id": 1,
        "variables": {
            "company_name": "Acme Corp",
            "include_charts": true,
            "max_findings": 10
        }
    }
    """
    # Fetch project
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Fetch template
    template = session.get(ReportTemplate, template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    
    # Generate filename
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    file_path = f"/tmp/report_{project_id}_{template_id}_{timestamp}.pdf"

@app.post("/projects/{project_id}/report/assemble", response_class=FileResponse)
async def assemble_modular_report(
    project_id: int,
    modules: List[str] = Body(..., description="List of module names to include (in order)"),
    variables: Optional[Dict[str, Any]] = Body(default=None, description="Optional template variables"),
    session: Session = Depends(get_session)
):
    """
    Assemble a modular report from selected template modules.
    
    This endpoint allows you to compose a custom report by selecting and ordering
    reusable template modules (title_page, executive_summary, detailed_findings, etc.).
    
    POST body:
    ```json
    {
        "modules": [
            "title_page",
            "executive_summary",
            "risk_charts",
            "detailed_findings",
            "recommendations"
        ],
        "variables": {
            "company_name": "Acme Corporation",
            "assessment_period": "Q4 2024",
            "include_charts": true
        }
    }
    ```
    
    Available modules:
    - title_page: Project title, metadata, and company branding
    - executive_summary: High-level overview and key metrics
    - risk_charts: Visual risk distribution and trends
    - top_findings: Top N critical findings summary
    - detailed_findings: Full finding details with all fields
    - recommendations: Remediation recommendations and action items
    - appendix: Additional technical details
    - sla_status: SLA tracking and deadline summary
    - compliance_owasp: OWASP Top 10 compliance mapping
    - compliance_cwe: CWE Top 25 compliance mapping
    - jira_integration: Jira ticket status and linking
    
    Returns:
        Assembled DOCX report with selected modules merged into a single document
    """
    # Fetch project with all findings
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Validate modules list
    if not modules:
        raise HTTPException(status_code=400, detail="Must specify at least one module")
    
    # Generate filename
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    file_path = f"/tmp/modular_report_{project_id}_{timestamp}.docx"
    
    try:
        # Assemble the modular report
        report_bytes = assemble_report(
            project=project,
            modules=modules,
            variables=variables or {}
        )
        
        # Write to temporary file
        with open(file_path, 'wb') as f:
            f.write(report_bytes)
        
        # Generate descriptive filename
        modules_str = "_".join(modules[:3])  # First 3 modules
        if len(modules) > 3:
            modules_str += f"_plus{len(modules) - 3}"
        
        filename = f"{project.name.replace(' ', '_')}_Modular_{modules_str}_Report.docx"
        
        return FileResponse(
            file_path,
            media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            filename=filename
        )
        
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        print(f"Error assembling modular report: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to assemble report: {str(e)}")


@app.get("/report/modules")
def get_available_modules():
    """
    List all available report modules with metadata.
    
    Returns information about each module including whether it exists,
    its path, and a brief description.
    
    Use this endpoint to discover which modules can be used with the
    /projects/{id}/report/assemble endpoint.
    """
    try:
        modules_info = list_available_modules()
        return {
            "modules": modules_info,
            "total": len(modules_info),
            "available": sum(1 for m in modules_info if m["exists"]),
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to list modules: {str(e)}")


@app.get("/report/modules/generate-defaults")
def generate_default_module_templates():
    """
    Generate default module templates programmatically.
    
    This endpoint creates default DOCX templates for all available modules
    if they don't already exist. Useful for initial setup or regenerating
    templates after updates.
    
    Note: This will overwrite existing templates with the same names.
    """
    try:
        from app.report_modules.generate_templates import main as generate_main
        
        # Run the template generator
        generate_main()
        
        return {
            "message": "Successfully generated default module templates",
            "modules": list_available_modules()
        }
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate default templates: {str(e)}"
        )
    
    # Import render_template function
    from app.reports import render_template
    from app.timezone_utils import get_utc_now
    
    # Generate the report from template
    try:
        render_template(
            template=template,
            project=project,
            file_path=file_path,
            variables=variables
        )
        
        # Track template usage
        template.usage_count += 1
        template.last_used_at = get_utc_now()
        session.add(template)
        session.commit()
        
    except Exception as e:
        print(f"Error generating report from template: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate report: {str(e)}")
    
    # Return the PDF file
    return FileResponse(
        file_path,
        media_type="application/pdf",
        filename=f"{project.name.replace(' ', '_')}_{template.name.replace(' ', '_')}.pdf",
    )


# --- Export Endpoints ---

@app.get("/projects/{project_id}/export")
def export_findings(
    project_id: int,
    format: str = "excel",  # 'excel', 'csv', 'json', or 'markdown'
    columns: Optional[str] = None,  # Comma-separated column names
    risk_filter: Optional[str] = None,  # Comma-separated risk levels
    status_filter: Optional[str] = None,  # Comma-separated issue statuses
    review_filter: Optional[str] = None,  # Comma-separated review statuses
    session: Session = Depends(get_session)
):
    """
    Export findings with customizable columns and filters.
    
    Query Parameters:
    - format: 'excel', 'csv', 'json', or 'markdown' (default: 'excel')
    - columns: Comma-separated list of columns to include (default: all)
              Available: title, risk_rating, description, remediation, instance_count,
                        review_status, reviewer_name, jira_issue_key, jira_status,
                        remediation_deadline, sla_status, remediation_owner, issue_status
    - risk_filter: Comma-separated risk levels (Critical, High, Medium, Low, Informational)
    - status_filter: Comma-separated issue statuses (Open, Partially Closed, Closed)
    - review_filter: Comma-separated review statuses (Pending, In Review, Approved, Rejected)
    """
    import csv
    import io
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
    
    # Fetch project
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Fetch findings with instances
    findings = session.exec(
        select(Finding).where(Finding.project_id == project_id)
    ).all()
    
    # Apply filters
    filtered_findings = []
    for finding in findings:
        # Risk filter
        if risk_filter:
            risk_levels = [r.strip() for r in risk_filter.split(',')]
            if finding.risk_rating not in risk_levels:
                continue
        
        # Issue status filter
        if status_filter:
            statuses = [s.strip() for s in status_filter.split(',')]
            if finding.issue_status not in statuses:
                continue
        
        # Review status filter
        if review_filter:
            reviews = [r.strip() for r in review_filter.split(',')]
            if finding.review_status not in reviews:
                continue
        
        filtered_findings.append(finding)
    
    # Define all available columns
    all_columns = {
        'title': 'Title',
        'risk_rating': 'Risk Rating',
        'description': 'Description',
        'remediation': 'Remediation',
        'instance_count': 'Instance Count',
        'review_status': 'Review Status',
        'reviewer_name': 'Reviewer',
        'jira_issue_key': 'Jira Issue',
        'jira_status': 'Jira Status',
        'remediation_deadline': 'Deadline',
        'sla_status': 'SLA Status',
        'remediation_owner': 'Owner',
        'issue_status': 'Issue Status',
    }
    
    # Parse selected columns
    if columns:
        selected_cols = [c.strip() for c in columns.split(',')]
        # Validate columns
        invalid_cols = [c for c in selected_cols if c not in all_columns]
        if invalid_cols:
            raise HTTPException(
                status_code=400, 
                detail=f"Invalid columns: {', '.join(invalid_cols)}"
            )
        column_map = {k: v for k, v in all_columns.items() if k in selected_cols}
    else:
        column_map = all_columns
    
    # Build data rows
    rows = []
    for finding in filtered_findings:
        row = {}
        for col_key, col_name in column_map.items():
            if col_key == 'instance_count':
                # Count instances
                instance_count = session.exec(
                    select(Instance).where(Instance.finding_id == finding.id)
                ).all()
                row[col_name] = len(instance_count)
            else:
                # Get attribute value
                value = getattr(finding, col_key, '')
                row[col_name] = value if value is not None else ''
        rows.append(row)
    
    # Generate Excel file
    if format.lower() == 'excel':
        wb = Workbook()
        ws = wb.active
        ws.title = "Findings"
        
        # Add headers
        headers = list(column_map.values())
        ws.append(headers)
        
        # Style header row
        header_fill = PatternFill(start_color="E0E0E0", end_color="E0E0E0", fill_type="solid")
        header_font = Font(bold=True)
        for cell in ws[1]:
            cell.fill = header_fill
            cell.font = header_font
        
        # Add data rows
        for row in rows:
            ws.append([row.get(h, '') for h in headers])
        
        # Auto-adjust column widths
        for column in ws.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 60)
            ws.column_dimensions[column_letter].width = adjusted_width
        
        # Save to bytes
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        
        return StreamingResponse(
            output,
            media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            headers={"Content-Disposition": f"attachment; filename={project.name.replace(' ', '_')}_findings.xlsx"}
        )
    
    # Generate CSV file
    elif format.lower() == 'csv':
        output = io.StringIO()
        headers = list(column_map.values())
        writer = csv.DictWriter(output, fieldnames=headers)
        writer.writeheader()
        writer.writerows(rows)
        
        # Convert to bytes
        output.seek(0)
        
        return StreamingResponse(
            iter([output.getvalue()]),
            media_type='text/csv',
            headers={"Content-Disposition": f"attachment; filename={project.name.replace(' ', '_')}_findings.csv"}
        )
    
    # Generate JSON file
    elif format.lower() == 'json':
        # Build full JSON structure with project metadata
        export_data = {
            "project": {
                "id": project.id,
                "name": project.name,
                "consultant_name": project.consultant_name,
            },
            "export_metadata": {
                "exported_at": get_utc_now().isoformat(),
                "total_findings": len(filtered_findings),
                "columns_included": list(column_map.values()),
                "filters_applied": {
                    "risk_levels": risk_filter.split(',') if risk_filter else None,
                    "issue_statuses": status_filter.split(',') if status_filter else None,
                    "review_statuses": review_filter.split(',') if review_filter else None,
                }
            },
            "findings": []
        }
        
        # Add full finding data
        for finding in filtered_findings:
            finding_data = {}
            for col_key in column_map.keys():
                if col_key == 'instance_count':
                    instance_count = session.exec(
                        select(Instance).where(Instance.finding_id == finding.id)
                    ).all()
                    finding_data[col_key] = len(instance_count)
                else:
                    value = getattr(finding, col_key, None)
                    # Convert datetime to ISO format
                    if hasattr(value, 'isoformat'):
                        finding_data[col_key] = value.isoformat()
                    else:
                        finding_data[col_key] = value
            
            # Add instances if needed
            if 'instance_count' in column_map:
                instances = session.exec(
                    select(Instance).where(Instance.finding_id == finding.id)
                ).all()
                finding_data['instances'] = [
                    {
                        "location": inst.location,
                        "details": inst.details,
                        "status": inst.status,
                        "created_at": inst.created_at.isoformat() if inst.created_at else None
                    } for inst in instances
                ]
            
            export_data["findings"].append(finding_data)
        
        return JSONResponse(
            content=export_data,
            headers={"Content-Disposition": f"attachment; filename={project.name.replace(' ', '_')}_findings.json"}
        )
    
    # Generate Markdown file
    elif format.lower() == 'markdown':
        output = io.StringIO()
        
        # Write header
        output.write(f"# Vulnerability Assessment Report: {project.name}\n\n")
        output.write(f"**Consultant:** {project.consultant_name}\n\n")
        output.write(f"**Generated:** {get_utc_now().strftime('%B %d, %Y at %H:%M UTC')}\n\n")
        output.write(f"**Total Findings:** {len(filtered_findings)}\n\n")
        
        # Write summary table
        if filtered_findings:
            output.write("## Summary\n\n")
            output.write("| Risk Level | Count |\n")
            output.write("|------------|-------|\n")
            risk_counts = {}
            for finding in filtered_findings:
                risk = finding.risk_rating
                risk_counts[risk] = risk_counts.get(risk, 0) + 1
            for risk in ['Critical', 'High', 'Medium', 'Low', 'Informational']:
                if risk in risk_counts:
                    output.write(f"| {risk} | {risk_counts[risk]} |\n")
            output.write("\n")
        
        # Write findings
        output.write("## Findings\n\n")
        for idx, finding in enumerate(filtered_findings, 1):
            output.write(f"### {idx}. {finding.title}\n\n")
            
            # Risk rating badge
            risk_emoji = {
                'Critical': '🔴',
                'High': '🟠',
                'Medium': '🟡',
                'Low': '🟢',
                'Informational': '🔵'
            }
            emoji = risk_emoji.get(finding.risk_rating, '⚪')
            output.write(f"**Risk Rating:** {emoji} {finding.risk_rating}\n\n")
            
            # Add selected fields
            if 'description' in column_map:
                output.write(f"**Description:**\n\n{finding.description}\n\n")
            
            if 'remediation' in column_map:
                output.write(f"**Remediation:**\n\n{finding.remediation}\n\n")
            
            if 'instance_count' in column_map:
                instances = session.exec(
                    select(Instance).where(Instance.finding_id == finding.id)
                ).all()
                output.write(f"**Instances Found:** {len(instances)}\n\n")
                if instances:
                    output.write("**Instance Details:**\n\n")
                    for i, inst in enumerate(instances, 1):
                        output.write(f"{i}. **Location:** `{inst.location}`\n")
                        output.write(f"   - **Details:** {inst.details}\n")
                        output.write(f"   - **Status:** {inst.status}\n")
                    output.write("\n")
            
            if 'review_status' in column_map:
                output.write(f"**Review Status:** {finding.review_status}\n\n")
            
            if 'issue_status' in column_map:
                output.write(f"**Issue Status:** {finding.issue_status}\n\n")
            
            if 'jira_issue_key' in column_map and finding.jira_issue_key:
                output.write(f"**Jira Issue:** {finding.jira_issue_key}\n\n")
            
            output.write("---\n\n")
        
        output.seek(0)
        
        return StreamingResponse(
            iter([output.getvalue()]),
            media_type='text/markdown',
            headers={"Content-Disposition": f"attachment; filename={project.name.replace(' ', '_')}_findings.md"}
        )
    
    else:
        raise HTTPException(status_code=400, detail="Invalid format. Use 'excel', 'csv', 'json', or 'markdown'")


# --- Enhanced Export Formats (v1.1.0 Phase 4) ---

@app.get("/projects/{project_id}/export/sarif")
def export_sarif(
    project_id: int,
    session: Session = Depends(get_session)
):
    """
    Export findings in SARIF 2.1.0 format for CI/CD integration.
    SARIF (Static Analysis Results Interchange Format) is widely supported by:
    - GitHub Security, GitLab, Azure DevOps
    - SonarQube, Snyk, Checkmarx
    - VS Code, IntelliJ IDEA
    """
    import json
    from datetime import datetime
    
    # Fetch project
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Fetch findings with instances
    findings = session.exec(
        select(Finding).where(Finding.project_id == project_id)
    ).all()
    
    # Map risk ratings to SARIF severity levels
    risk_to_severity = {
        'Critical': 'error',
        'High': 'error',
        'Medium': 'warning',
        'Low': 'note',
        'Informational': 'none'
    }
    
    # Build SARIF document
    sarif_doc = {
        "$schema": "https://raw.githubusercontent.com/oasis-tcs/sarif-spec/master/Schemata/sarif-schema-2.1.0.json",
        "version": "2.1.0",
        "runs": [
            {
                "tool": {
                    "driver": {
                        "name": "VulnManager",
                        "informationUri": "https://github.com/aphesz/vuln-manager",
                        "version": "1.1.0",
                        "semanticVersion": "1.1.0",
                        "rules": []
                    }
                },
                "results": [],
                "properties": {
                    "projectName": project.name,
                    "projectId": project.id,
                    "consultantName": project.consultant_name,
                    "exportedAt": get_utc_now().isoformat()
                }
            }
        ]
    }
    
    run = sarif_doc["runs"][0]
    
    # Add rules (one per unique finding)
    for finding in findings:
        rule = {
            "id": f"VULN-{finding.id}",
            "name": finding.title.replace(" ", ""),
            "shortDescription": {
                "text": finding.title
            },
            "fullDescription": {
                "text": finding.description or finding.title
            },
            "help": {
                "text": finding.remediation or "No remediation guidance available",
                "markdown": finding.remediation or "No remediation guidance available"
            },
            "defaultConfiguration": {
                "level": risk_to_severity.get(finding.risk_rating, 'warning')
            },
            "properties": {
                "tags": [finding.risk_rating, finding.issue_status],
                "precision": "high"
            }
        }
        run["tool"]["driver"]["rules"].append(rule)
        
        # Add result instances
        instances = session.exec(
            select(Instance).where(Instance.finding_id == finding.id)
        ).all()
        
        for instance in instances:
            result = {
                "ruleId": f"VULN-{finding.id}",
                "level": risk_to_severity.get(finding.risk_rating, 'warning'),
                "message": {
                    "text": f"{finding.title}: {instance.details or 'See description'}"
                },
                "locations": [
                    {
                        "physicalLocation": {
                            "artifactLocation": {
                                "uri": instance.location or "unknown"
                            }
                        }
                    }
                ],
                "properties": {
                    "instanceId": instance.id,
                    "status": instance.status,
                    "riskRating": finding.risk_rating,
                    "reviewStatus": finding.review_status,
                    "issueStatus": finding.issue_status
                }
            }
            run["results"].append(result)
    
    # Return JSON
    return Response(
        content=json.dumps(sarif_doc, indent=2),
        media_type="application/json",
        headers={"Content-Disposition": f"attachment; filename={project.name.replace(' ', '_')}_findings.sarif"}
    )


@app.get("/projects/{project_id}/export/html")
def export_html_interactive(
    project_id: int,
    session: Session = Depends(get_session)
):
    """
    Export interactive HTML report with sortable/filterable findings table.
    Includes JavaScript for client-side filtering and sorting.
    """
    # Fetch project
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Fetch findings with instances
    findings = session.exec(
        select(Finding).where(Finding.project_id == project_id)
    ).all()
    
    # Calculate risk distribution
    risk_counts = {'Critical': 0, 'High': 0, 'Medium': 0, 'Low': 0, 'Informational': 0}
    for finding in findings:
        if finding.risk_rating in risk_counts:
            risk_counts[finding.risk_rating] += 1
    
    # Build HTML
    html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{project.name} - Security Assessment Report</title>
    <style>
        * {{ box-sizing: border-box; margin: 0; padding: 0; }}
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Arial, sans-serif; line-height: 1.6; color: #333; background: #f5f5f5; }}
        .container {{ max-width: 1400px; margin: 0 auto; padding: 20px; }}
        .header {{ background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 40px; border-radius: 12px; margin-bottom: 30px; box-shadow: 0 4px 6px rgba(0,0,0,0.1); }}
        .header h1 {{ font-size: 2.5rem; margin-bottom: 10px; }}
        .header .subtitle {{ opacity: 0.9; font-size: 1.1rem; }}
        .stats {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 20px; margin-bottom: 30px; }}
        .stat-card {{ background: white; padding: 20px; border-radius: 8px; box-shadow: 0 2px 4px rgba(0,0,0,0.1); text-align: center; }}
        .stat-card .number {{ font-size: 2.5rem; font-weight: bold; margin-bottom: 5px; }}
        .stat-card .label {{ color: #666; font-size: 0.9rem; text-transform: uppercase; letter-spacing: 0.5px; }}
        .critical {{ color: #d32f2f; }}
        .high {{ color: #f57c00; }}
        .medium {{ color: #fbc02d; }}
        .low {{ color: #388e3c; }}
        .info {{ color: #1976d2; }}
        .controls {{ background: white; padding: 20px; border-radius: 8px; margin-bottom: 20px; box-shadow: 0 2px 4px rgba(0,0,0,0.1); }}
        .controls label {{ display: inline-block; margin-right: 15px; font-weight: 500; }}
        .controls input, .controls select {{ padding: 8px 12px; border: 1px solid #ddd; border-radius: 4px; margin-right: 10px; }}
        .table-container {{ background: white; border-radius: 8px; box-shadow: 0 2px 4px rgba(0,0,0,0.1); overflow: hidden; }}
        table {{ width: 100%; border-collapse: collapse; }}
        th {{ background: #667eea; color: white; padding: 15px; text-align: left; font-weight: 600; cursor: pointer; user-select: none; }}
        th:hover {{ background: #5568d3; }}
        th::after {{ content: ' ⇅'; opacity: 0.5; }}
        td {{ padding: 12px 15px; border-bottom: 1px solid #eee; }}
        tr:hover {{ background: #f8f9fa; }}
        .risk-badge {{ padding: 4px 12px; border-radius: 20px; font-size: 0.85rem; font-weight: 600; text-transform: uppercase; }}
        .risk-critical {{ background: #ffebee; color: #d32f2f; }}
        .risk-high {{ background: #fff3e0; color: #f57c00; }}
        .risk-medium {{ background: #fffde7; color: #f9a825; }}
        .risk-low {{ background: #e8f5e9; color: #388e3c; }}
        .risk-info {{ background: #e3f2fd; color: #1976d2; }}
        .status-badge {{ padding: 4px 10px; border-radius: 4px; font-size: 0.8rem; }}
        .status-open {{ background: #ffcdd2; color: #c62828; }}
        .status-partial {{ background: #fff9c4; color: #f57f17; }}
        .status-closed {{ background: #c8e6c9; color: #2e7d32; }}
        .expandable {{ cursor: pointer; color: #667eea; text-decoration: underline; }}
        .details {{ padding: 20px; background: #f8f9fa; margin: 10px 0; border-left: 4px solid #667eea; border-radius: 4px; }}
        .details h3 {{ color: #667eea; margin-top: 15px; margin-bottom: 10px; font-size: 1.1rem; }}
        .details h3:first-child {{ margin-top: 0; }}
        .details p {{ line-height: 1.8; color: #444; margin-bottom: 15px; }}
        .details ul {{ margin-left: 20px; }}
        .details li {{ margin-bottom: 8px; line-height: 1.6; }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>🔒 {project.name}</h1>
            <p class="subtitle">Security Assessment Report • Generated {get_utc_now().strftime('%B %d, %Y')}</p>
        </div>
        
        <div class="stats">
            <div class="stat-card">
                <div class="number critical">{risk_counts['Critical']}</div>
                <div class="label">Critical</div>
            </div>
            <div class="stat-card">
                <div class="number high">{risk_counts['High']}</div>
                <div class="label">High</div>
            </div>
            <div class="stat-card">
                <div class="number medium">{risk_counts['Medium']}</div>
                <div class="label">Medium</div>
            </div>
            <div class="stat-card">
                <div class="number low">{risk_counts['Low']}</div>
                <div class="label">Low</div>
            </div>
            <div class="stat-card">
                <div class="number info">{risk_counts['Informational']}</div>
                <div class="label">Informational</div>
            </div>
        </div>
        
        <div class="controls">
            <label>Filter: <input type="text" id="searchInput" placeholder="Search findings..." /></label>
            <label>Risk: 
                <select id="riskFilter">
                    <option value="">All</option>
                    <option value="Critical">Critical</option>
                    <option value="High">High</option>
                    <option value="Medium">Medium</option>
                    <option value="Low">Low</option>
                    <option value="Informational">Informational</option>
                </select>
            </label>
            <label>Status: 
                <select id="statusFilter">
                    <option value="">All</option>
                    <option value="Open">Open</option>
                    <option value="Partially Closed">Partially Closed</option>
                    <option value="Closed">Closed</option>
                </select>
            </label>
        </div>
        
        <div class="table-container">
            <table id="findingsTable">
                <thead>
                    <tr>
                        <th onclick="sortTable(0)">ID</th>
                        <th onclick="sortTable(1)">Title</th>
                        <th onclick="sortTable(2)">Risk</th>
                        <th onclick="sortTable(3)">Status</th>
                        <th onclick="sortTable(4)">Instances</th>
                        <th>Actions</th>
                    </tr>
                </thead>
                <tbody>"""
    
    # Add findings
    for finding in findings:
        instances = session.exec(
            select(Instance).where(Instance.finding_id == finding.id)
        ).all()
        
        risk_class = finding.risk_rating.lower()
        status_class = finding.issue_status.lower().replace(' ', '-')
        
        # Use content directly (already contains HTML from rich text editor)
        title = finding.title or 'Untitled'
        description = finding.description or 'No description provided'
        remediation = finding.remediation or 'No remediation guidance provided'
        
        html_content += f"""
                    <tr data-risk="{finding.risk_rating}" data-status="{finding.issue_status}">
                        <td>{finding.id}</td>
                        <td><strong>{title}</strong></td>
                        <td><span class="risk-badge risk-{risk_class}">{finding.risk_rating}</span></td>
                        <td><span class="status-badge status-{status_class}">{finding.issue_status}</span></td>
                        <td>{len(instances)}</td>
                        <td><span class="expandable" onclick="toggleDetails({finding.id})">View Details</span></td>
                    </tr>
                    <tr id="details-{finding.id}" style="display: none;">
                        <td colspan="6">
                            <div class="details">
                                <h3>📝 Description</h3>
                                <div style="white-space: pre-wrap;">{description}</div>
                                <h3>🔧 Remediation</h3>
                                <div style="white-space: pre-wrap;">{remediation}</div>
                                <h3>📍 Instances ({len(instances)})</h3>
                                <ul>"""
        
        for instance in instances:
            location = instance.location or 'Unknown location'
            details = instance.details or 'No details'
            html_content += f"""
                                    <li><strong>{location}</strong> - {details} (Status: {instance.status})</li>"""
        
        html_content += """
                                </ul>
                            </div>
                        </td>
                    </tr>"""
    
    html_content += """
                </tbody>
            </table>
        </div>
    </div>
    
    <script>
        // Search filter
        document.getElementById('searchInput').addEventListener('keyup', filterTable);
        document.getElementById('riskFilter').addEventListener('change', filterTable);
        document.getElementById('statusFilter').addEventListener('change', filterTable);
        
        function filterTable() {
            const searchValue = document.getElementById('searchInput').value.toLowerCase();
            const riskValue = document.getElementById('riskFilter').value;
            const statusValue = document.getElementById('statusFilter').value;
            const table = document.getElementById('findingsTable');
            const rows = table.getElementsByTagName('tr');
            
            for (let i = 1; i < rows.length; i += 2) {
                const row = rows[i];
                const title = row.cells[1].textContent.toLowerCase();
                const risk = row.getAttribute('data-risk');
                const status = row.getAttribute('data-status');
                
                const matchesSearch = title.includes(searchValue);
                const matchesRisk = !riskValue || risk === riskValue;
                const matchesStatus = !statusValue || status === statusValue;
                
                if (matchesSearch && matchesRisk && matchesStatus) {
                    row.style.display = '';
                    if (i + 1 < rows.length) rows[i + 1].style.display = '';
                } else {
                    row.style.display = 'none';
                    if (i + 1 < rows.length) rows[i + 1].style.display = 'none';
                }
            }
        }
        
        function toggleDetails(id) {
            const detailsRow = document.getElementById('details-' + id);
            if (detailsRow.style.display === 'none' || detailsRow.style.display === '') {
                detailsRow.style.display = 'table-row';
            } else {
                detailsRow.style.display = 'none';
            }
        }
        
        function sortTable(n) {
            const table = document.getElementById('findingsTable');
            let switching = true;
            let dir = 'asc';
            let switchcount = 0;
            
            while (switching) {
                switching = false;
                const rows = table.rows;
                
                for (let i = 1; i < (rows.length - 1); i += 2) {
                    let shouldSwitch = false;
                    const x = rows[i].getElementsByTagName('TD')[n];
                    const y = rows[i + 2].getElementsByTagName('TD')[n];
                    
                    if (dir === 'asc') {
                        if (x.innerHTML.toLowerCase() > y.innerHTML.toLowerCase()) {
                            shouldSwitch = true;
                            break;
                        }
                    } else if (dir === 'desc') {
                        if (x.innerHTML.toLowerCase() < y.innerHTML.toLowerCase()) {
                            shouldSwitch = true;
                            break;
                        }
                    }
                }
                
                if (shouldSwitch) {
                    rows[i].parentNode.insertBefore(rows[i + 2], rows[i]);
                    rows[i].parentNode.insertBefore(rows[i + 3], rows[i + 1]);
                    switching = true;
                    switchcount++;
                } else {
                    if (switchcount === 0 && dir === 'asc') {
                        dir = 'desc';
                        switching = true;
                    }
                }
            }
        }
    </script>
</body>
</html>"""
    
    return Response(
        content=html_content,
        media_type="text/html",
        headers={"Content-Disposition": f"attachment; filename={project.name.replace(' ', '_')}_interactive.html"}
    )


@app.get("/projects/{project_id}/export/pptx")
def export_powerpoint(
    project_id: int,
    session: Session = Depends(get_session)
):
    """
    Export PowerPoint presentation with executive summary and findings.
    Uses python-pptx library to generate professional slide deck.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    import io
    
    # Fetch project
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Fetch findings
    findings = session.exec(
        select(Finding).where(Finding.project_id == project_id)
    ).all()
    
    # Calculate statistics
    risk_counts = {'Critical': 0, 'High': 0, 'Medium': 0, 'Low': 0, 'Informational': 0}
    status_counts = {'Open': 0, 'Partially Closed': 0, 'Closed': 0}
    
    for finding in findings:
        if finding.risk_rating in risk_counts:
            risk_counts[finding.risk_rating] += 1
        if finding.issue_status in status_counts:
            status_counts[finding.issue_status] += 1
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    color_critical = RGBColor(211, 47, 47)
    color_high = RGBColor(245, 124, 0)
    color_medium = RGBColor(251, 192, 45)
    color_low = RGBColor(56, 142, 60)
    color_info = RGBColor(25, 118, 210)
    color_primary = RGBColor(102, 126, 234)
    
    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"{project.name}"
    subtitle.text = f"Security Assessment Report\\n{get_utc_now().strftime('%B %d, %Y')}"
    
    # Slide 2: Executive Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Executive Summary"
    
    tf = body_shape.text_frame
    tf.text = f"Total Findings: {len(findings)}"
    
    p = tf.add_paragraph()
    p.text = f"Critical: {risk_counts['Critical']}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"High: {risk_counts['High']}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"Medium: {risk_counts['Medium']}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"Low: {risk_counts['Low']}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"Informational: {risk_counts['Informational']}"
    p.level = 1
    
    # Slide 3: Risk Distribution
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.75)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Risk Distribution"
    
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color_primary
    
    # Add risk breakdown
    y_offset = 1.5
    for risk, count in risk_counts.items():
        if count > 0:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(y_offset), Inches(3), Inches(0.5))
            tf = txBox.text_frame
            tf.text = f"{risk}:"
            p = tf.paragraphs[0]
            p.font.size = Pt(18)
            p.font.bold = True
            
            # Add count with color
            txBox = slide.shapes.add_textbox(Inches(4), Inches(y_offset), Inches(2), Inches(0.5))
            tf = txBox.text_frame
            tf.text = str(count)
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            
            if risk == 'Critical':
                p.font.color.rgb = color_critical
            elif risk == 'High':
                p.font.color.rgb = color_high
            elif risk == 'Medium':
                p.font.color.rgb = color_medium
            elif risk == 'Low':
                p.font.color.rgb = color_low
            else:
                p.font.color.rgb = color_info
            
            y_offset += 0.7
    
    # Add slides for top findings (Critical and High only)
    critical_high_findings = [f for f in findings if f.risk_rating in ['Critical', 'High']]
    
    for finding in critical_high_findings[:10]:  # Limit to top 10
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        # Title with risk badge
        title_shape.text = f"{finding.title}"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        
        # Body
        tf = body_shape.text_frame
        tf.clear()
        
        # Risk rating
        p = tf.add_paragraph()
        p.text = f"Risk: {finding.risk_rating}"
        p.font.size = Pt(16)
        p.font.bold = True
        if finding.risk_rating == 'Critical':
            p.font.color.rgb = color_critical
        else:
            p.font.color.rgb = color_high
        
        # Description
        p = tf.add_paragraph()
        p.text = "Description:"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        desc_text = finding.description[:200] + "..." if finding.description and len(finding.description) > 200 else (finding.description or "No description")
        p.text = desc_text
        p.font.size = Pt(12)
        p.level = 1
        
        # Remediation
        p = tf.add_paragraph()
        p.text = "Remediation:"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        rem_text = finding.remediation[:150] + "..." if finding.remediation and len(finding.remediation) > 150 else (finding.remediation or "See full report")
        p.text = rem_text
        p.font.size = Pt(12)
        p.level = 1
    
    # Final slide: Next Steps
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Next Steps"
    
    tf = body_shape.text_frame
    tf.text = "Review and prioritize critical and high-risk findings"
    
    p = tf.add_paragraph()
    p.text = "Assign remediation owners and deadlines"
    
    p = tf.add_paragraph()
    p.text = "Track progress in VulnManager dashboard"
    
    p = tf.add_paragraph()
    p.text = "Schedule follow-up assessment"
    
    # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return StreamingResponse(
        iter([output.getvalue()]),
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        headers={"Content-Disposition": f"attachment; filename={project.name.replace(' ', '_')}_presentation.pptx"}
    )


@app.post("/export/bulk")
def export_bulk(
    project_ids: List[int],
    format: str = "pdf",
    session: Session = Depends(get_session)
):
    """
    Export multiple projects as a ZIP archive.
    Generates individual reports for each project and packages them together.
    """
    import io
    import zipfile
    from datetime import datetime
    
    if not project_ids:
        raise HTTPException(status_code=400, detail="No project IDs provided")
    
    # Validate format
    valid_formats = ['pdf', 'docx', 'html', 'json', 'sarif']
    if format not in valid_formats:
        raise HTTPException(status_code=400, detail=f"Invalid format. Use one of: {', '.join(valid_formats)}")
    
    # Create ZIP file in memory
    zip_buffer = io.BytesIO()
    
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        for project_id in project_ids:
            # Fetch project
            project = session.exec(select(Project).where(Project.id == project_id)).first()
            if not project:
                continue  # Skip missing projects
            
            # Generate filename
            safe_name = project.name.replace(' ', '_').replace('/', '_')
            
            # Call appropriate export endpoint based on format
            try:
                if format == 'html':
                    # Generate HTML content
                    response = export_html_interactive(project_id, session)
                    content = response.body
                    filename = f"{safe_name}_report.html"
                
                elif format == 'sarif':
                    # Generate SARIF content
                    response = export_sarif(project_id, session)
                    content = response.body
                    filename = f"{safe_name}_findings.sarif"
                
                elif format == 'json':
                    # Export as JSON
                    findings = session.exec(
                        select(Finding).where(Finding.project_id == project_id)
                    ).all()
                    
                    export_data = {
                        "project": {
                            "id": project.id,
                            "name": project.name,
                            "consultant_name": project.consultant_name
                        },
                        "findings": []
                    }
                    
                    for finding in findings:
                        instances = session.exec(
                            select(Instance).where(Instance.finding_id == finding.id)
                        ).all()
                        
                        export_data["findings"].append({
                            "id": finding.id,
                            "title": finding.title,
                            "risk_rating": finding.risk_rating,
                            "description": finding.description,
                            "remediation": finding.remediation,
                            "issue_status": finding.issue_status,
                            "review_status": finding.review_status,
                            "instances": [
                                {
                                    "location": inst.location,
                                    "details": inst.details,
                                    "status": inst.status
                                }
                                for inst in instances
                            ]
                        })
                    
                    import json
                    content = json.dumps(export_data, indent=2).encode('utf-8')
                    filename = f"{safe_name}_findings.json"
                
                else:  # pdf or docx - use existing report generation
                    # For now, skip PDF/DOCX in bulk export (would need to refactor report generation)
                    continue
                
                # Add to ZIP
                zip_file.writestr(filename, content)
                
            except Exception as e:
                logger.error(f"Error exporting project {project_id}: {str(e)}")
                continue
    
    zip_buffer.seek(0)
    
    # Generate ZIP filename with timestamp
    timestamp = datetime.utcnow().strftime('%Y%m%d_%H%M%S')
    zip_filename = f"bulk_export_{format}_{timestamp}.zip"
    
    return StreamingResponse(
        iter([zip_buffer.getvalue()]),
        media_type='application/zip',
        headers={"Content-Disposition": f"attachment; filename={zip_filename}"}
    )


# --- Peer Review Workflow Endpoints ---

@app.patch("/findings/{finding_id}/review")
def update_finding_review_status(
    finding_id: int,
    data: ReviewStatusUpdate,
    session: Session = Depends(get_session)
):
    """
    Update the review status of a finding.
    
    Args:
        finding_id: ID of the finding
        data: Review status update data (status, reviewer_name)
    """
    finding = session.get(Finding, finding_id)
    if not finding:
        raise HTTPException(status_code=404, detail="Finding not found")
    
    # Map 'status' from request to 'review_status' for validation
    review_status = data.status
    reviewer_name = data.reviewer_name
    user = reviewer_name if reviewer_name else "system"
    
    # Validate review status
    valid_statuses = ['Pending', 'In Review', 'Approved', 'Rejected']
    if review_status not in valid_statuses:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid review status. Must be one of: {', '.join(valid_statuses)}"
        )
    
    old_status = finding.review_status.value if finding.review_status else None
    old_reviewer = finding.reviewer_name
    
    finding.review_status = FindingBase.ReviewStatus(review_status)
    finding.reviewer_name = reviewer_name
    
    # Create audit log entry
    changes = {
        "review_status": {
            "old_value": old_status,
            "new_value": review_status
        }
    }
    if old_reviewer != reviewer_name:
        changes["reviewer_name"] = {
            "old_value": old_reviewer,
            "new_value": reviewer_name
        }
    
    audit_entry = AuditLog(
        entity_type="finding",
        entity_id=finding_id,
        action="review_status_changed",
        user=user,
        timestamp=get_utc_now(),
        changes_json=json.dumps(changes)
    )
    
    session.add(finding)
    session.add(audit_entry)
    session.commit()
    session.refresh(finding)
    
    logger.info(f"Finding {finding_id} review status updated: {old_status} -> {review_status} by {user}")
    
    return {
        "id": finding.id,
        "review_status": finding.review_status.value,
        "reviewer_name": finding.reviewer_name,
        "updated_by": user
    }


@app.post("/projects/{project_id}/findings", response_model=FindingReadWithInstances, status_code=201)
@limiter.limit("20/minute")  # Max 20 manual findings per minute per IP
async def create_finding_manually(
    request: Request,
    project_id: int,
    title: str = Body(...),
    description: str = Body(...),
    remediation: str = Body(...),
    risk_rating: str = Body(...),
    template_id: Optional[int] = Body(None),
    instances: List[Dict[str, str]] = Body(..., description="List of instance objects with 'location' and 'details'"),
    issue_status: Optional[str] = Body("Open"),
    session: Session = Depends(get_session)
):
    """
    Manually create a new finding with instances (Quick Add feature).
    
    This endpoint allows creating findings from vulnerability templates
    or from scratch with multiple instances.
    
    Request body:
    - title: Finding title
    - description: Detailed description
    - remediation: Remediation guidance
    - risk_rating: One of: Critical, High, Medium, Low, Informational
    - template_id: Optional link to vulnerability template
    - instances: Array of {location, details} objects
    - issue_status: One of: Open, Partially Closed, Closed (default: Open)
    """
    from app.models import VulnerabilityTemplate
    
    # Input validation
    title = validate_string_length(sanitize_html_input(title), "title", max_length=200)
    description = validate_string_length(sanitize_html_input(description), "description", max_length=5000)
    remediation = validate_string_length(sanitize_html_input(remediation), "remediation", max_length=5000)
    
    # Validate project exists
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Validate risk rating
    try:
        risk = FindingBase.RiskRating(risk_rating)
    except ValueError:
        valid_ratings = [r.value for r in FindingBase.RiskRating]
        raise HTTPException(
            status_code=400,
            detail=f"Invalid risk_rating. Must be one of: {', '.join(valid_ratings)}"
        )
    
    # Validate issue status
    try:
        status = FindingBase.IssueStatus(issue_status) if issue_status else FindingBase.IssueStatus.Open
    except ValueError:
        valid_statuses = [s.value for s in FindingBase.IssueStatus]
        raise HTTPException(
            status_code=400,
            detail=f"Invalid issue_status. Must be one of: {', '.join(valid_statuses)}"
        )
    
    # Validate instances
    if not instances or len(instances) == 0:
        raise HTTPException(
            status_code=400,
            detail="At least one instance is required"
        )
    
    if len(instances) > 100:  # Limit to prevent abuse
        raise HTTPException(
            status_code=400,
            detail="Maximum 100 instances allowed per request"
        )
    
    for idx, inst in enumerate(instances):
        if 'location' not in inst or 'details' not in inst:
            raise HTTPException(
                status_code=400,
                detail=f"Instance {idx} must have 'location' and 'details' fields"
            )
        # Validate and sanitize instance data
        inst['location'] = validate_string_length(sanitize_html_input(inst['location']), f"Instance {idx} location", max_length=500)
        inst['details'] = validate_string_length(sanitize_html_input(inst['details']), f"Instance {idx} details", max_length=2000)
    
    # If template_id provided, validate it exists and update usage
    template = None
    if template_id:
        template = session.get(VulnerabilityTemplate, template_id)
        if not template:
            raise HTTPException(status_code=404, detail=f"Vulnerability template {template_id} not found")
        
        # Update template usage
        template.usage_count += 1
        template.last_used = get_utc_now()
        template.updated_at = get_utc_now()
        session.add(template)
    
    # Check for existing finding with same title in project
    existing_finding = session.exec(
        select(Finding).where(
            Finding.project_id == project_id,
            Finding.title == title
        )
    ).first()
    
    if existing_finding:
        # Add instances to existing finding
        finding = existing_finding
        logger.info(f"Adding instances to existing finding {finding.id}: {title}")
    else:
        # Create new finding
        finding = Finding(
            project_id=project_id,
            title=title,
            risk_rating=risk,
            description=description,
            remediation=remediation,
            template_id=template_id,
            issue_status=status,
            discovered_at=get_utc_now(),  # Track discovery time (v0.8.1)
            resolved_at=get_utc_now() if status == FindingBase.IssueStatus.Closed else None  # Track resolution (v0.8.1)
        )
        session.add(finding)
        session.flush()  # Get finding.id
        logger.info(f"Created new finding {finding.id}: {title}")
    
    # Create instances
    created_instances = []
    for inst_data in instances:
        instance = Instance(
            finding_id=finding.id,
            location=inst_data['location'],
            details=inst_data['details'],
            status='New - Unvalidated',
            created_at=get_utc_now()
        )
        session.add(instance)
        created_instances.append(instance)
    
    session.commit()
    session.refresh(finding)
    
    # Send WebSocket notification
    await ws_manager.broadcast_finding_update(
        project_id,
        finding.id,
        'update' if existing_finding else 'create'
    )
    
    logger.info(f"Successfully created finding with {len(created_instances)} instances")
    
    return finding


@app.post("/findings/{finding_id}/comments", response_model=CommentRead)
@limiter.limit("60/minute")  # Max 60 comments per minute per IP
def create_comment(
    request: Request,
    finding_id: int,
    comment_data: CommentCreate,
    session: Session = Depends(get_session)
):
    """
    Add a comment to a finding.
    
    Args:
        finding_id: ID of the finding
        comment_data: Comment data (text, user)
    """
    finding = session.get(Finding, finding_id)
    if not finding:
        raise HTTPException(status_code=404, detail="Finding not found")
    
    # Validate and sanitize comment
    comment_text = validate_string_length(sanitize_html_input(comment_data.text), "comment text", max_length=5000)
    
    comment = Comment(
        text=comment_text,
        user=comment_data.user,
        created_at=get_utc_now(),
        finding_id=finding_id
    )
    
    # Create audit log entry
    audit_entry = AuditLog(
        entity_type="comment",
        entity_id=0,  # Will be updated after commit
        action="created",
        user=comment_data.user,
        timestamp=get_utc_now(),
        changes_json=json.dumps({
            "finding_id": finding_id,
            "text_preview": comment_data.text[:100] + "..." if len(comment_data.text) > 100 else comment_data.text
        })
    )
    
    session.add(comment)
    session.commit()
    session.refresh(comment)
    
    # Update audit entry with comment ID
    audit_entry.entity_id = comment.id
    session.add(audit_entry)
    session.commit()
    
    logger.info(f"Comment {comment.id} added to finding {finding_id} by {comment_data.user}")
    
    return comment

@app.get("/findings/{finding_id}/comments", response_model=List[CommentRead])
def get_finding_comments(
    finding_id: int,
    session: Session = Depends(get_session)
):
    """Get all comments for a finding."""
    finding = session.get(Finding, finding_id)
    if not finding:
        raise HTTPException(status_code=404, detail="Finding not found")
    
    comments = session.exec(
        select(Comment)
        .where(Comment.finding_id == finding_id)
        .order_by(Comment.created_at)  # Chronological order (oldest first)
    ).all()
    
    return comments

@app.get("/audit-log", response_model=List[AuditLogRead])
def get_audit_log(
    entity_type: Optional[str] = None,
    entity_id: Optional[int] = None,
    user: Optional[str] = None,
    limit: int = 100,
    session: Session = Depends(get_session)
):
    """
    Get audit log entries with optional filtering.
    
    Args:
        entity_type: Filter by entity type (finding, project, comment, etc.)
        entity_id: Filter by entity ID
        user: Filter by user
        limit: Maximum number of entries to return (default 100, max 1000)
    """
    if limit > 1000:
        limit = 1000
    
    query = select(AuditLog)
    
    if entity_type:
        query = query.where(AuditLog.entity_type == entity_type)
    if entity_id:
        query = query.where(AuditLog.entity_id == entity_id)
    if user:
        query = query.where(AuditLog.user == user)
    
    query = query.order_by(AuditLog.timestamp.desc()).limit(limit)  # Newest first for better UX
    
    entries = session.exec(query).all()
    return entries

# --- Jira Integration Endpoints ---

@app.post("/jira/settings", response_model=JiraSettingsRead)
def create_jira_settings(
    data: JiraSettingsCreate,
    session: Session = Depends(get_session)
):
    """
    Create or update Jira integration settings.
    
    Security: API token is encrypted before storage.
    """
    # Encrypt the API token
    encrypted_token = encrypt_token(data.api_token)
    
    settings = JiraSettings(
        jira_url=data.jira_url,
        project_key=data.project_key,
        api_token_encrypted=encrypted_token,
        is_active=data.is_active,
        project_id=data.project_id
    )
    
    session.add(settings)
    session.commit()
    session.refresh(settings)
    
    logger.info(f"Jira settings created for project {data.project_id}")
    
    return settings

@app.get("/jira/settings/{project_id}", response_model=JiraSettingsRead)
def get_jira_settings(
    project_id: int,
    session: Session = Depends(get_session)
):
    """
    Retrieve Jira integration settings for a specific project.
    Returns 404 if settings not found.
    """
    statement = select(JiraSettings).where(JiraSettings.project_id == project_id)
    settings = session.exec(statement).first()
    
    if not settings:
        raise HTTPException(status_code=404, detail=f"Jira settings not found for project {project_id}")
    
    return settings

@app.post("/jira/test-connection")
async def test_jira_connection(
    data: JiraConnectionTest
):
    """Test Jira connection without saving settings."""
    from app.jira import JiraClient
    
    try:
        client = JiraClient(data.jira_url, data.email, data.api_token)
        result = await client.test_connection()
        
        if not result.get("success"):
            raise HTTPException(status_code=400, detail=result.get("message", "Connection failed"))
        
        return result
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Connection failed: {str(e)}")

@app.post("/findings/{finding_id}/create-jira-issue")
async def create_jira_issue_for_finding(
    finding_id: int,
    data: JiraIssueCreate,
    session: Session = Depends(get_session)
):
    """
    Create a Jira issue from a finding.
    
    Args:
        finding_id: ID of the finding
        data: Request data containing user information
    """
    finding = session.get(Finding, finding_id)
    if not finding:
        raise HTTPException(status_code=404, detail="Finding not found")
    
    if finding.jira_issue_key:
        raise HTTPException(
            status_code=400,
            detail=f"Finding already has Jira issue: {finding.jira_issue_key}"
        )
    
    # Get Jira client for the project
    client = get_jira_client(session, finding.project_id)
    if not client:
        raise HTTPException(
            status_code=400,
            detail="Jira integration not configured for this project"
        )
    
    # Get project key from settings
    settings_query = select(JiraSettings).where(
        JiraSettings.project_id == finding.project_id,
        JiraSettings.is_active == True
    )
    settings = session.exec(settings_query).first()
    if not settings:
        raise HTTPException(status_code=400, detail="No active Jira settings found")
    
    project_key = settings.project_key
    
    # Create the issue
    issue_data = await client.create_issue(project_key, finding)
    
    if not issue_data:
        raise HTTPException(status_code=500, detail="Failed to create Jira issue")
    
    # Update finding with Jira issue key
    finding.jira_issue_key = issue_data["key"]
    finding.jira_status = "Open"  # Default status
    
    # Create audit log entry
    audit_entry = AuditLog(
        entity_type="finding",
        entity_id=finding_id,
        action="jira_issue_created",
        user=data.user,
        timestamp=get_utc_now(),
        changes_json=json.dumps({
            "jira_issue_key": issue_data["key"],
            "jira_url": issue_data.get("self", "")
        })
    )
    
    session.add(finding)
    session.add(audit_entry)
    session.commit()
    session.refresh(finding)
    
    
    logger.info(f"Jira issue {issue_data['key']} created for finding {finding_id} by {data.user}")
    
    return {
        "jira_issue_key": finding.jira_issue_key,
        "jira_url": issue_data.get("self", ""),
        "finding_id": finding.id
    }

# ===============================
@app.post("/webhooks/jira")
async def jira_webhook(
    request: Request,
    session: Session = Depends(get_session)
):
    """
    Webhook endpoint for receiving Jira updates.
    
    Security: In production, verify webhook signature/token.
    """
    try:
        payload = await request.json()
        
        # Extract relevant data from Jira webhook
        webhook_event = payload.get("webhookEvent")
        issue_key = payload.get("issue", {}).get("key")
        new_status = payload.get("issue", {}).get("fields", {}).get("status", {}).get("name")
        
        logger.info(f"Received Jira webhook: {webhook_event} for {issue_key}")
        
        if not issue_key:
            return {"status": "ignored", "reason": "No issue key in payload"}
        
        # Find finding with this Jira issue key
        finding = session.exec(
            select(Finding).where(Finding.jira_issue_key == issue_key)
        ).first()
        
        if not finding:
            return {"status": "ignored", "reason": f"No finding found for issue {issue_key}"}
        
        # Update finding's Jira status
        old_status = finding.jira_status
        finding.jira_status = new_status
        
        # Create audit log entry
        audit_entry = AuditLog(
            entity_type="finding",
            entity_id=finding.id,
            action="jira_status_updated",
            user="jira_webhook",
            timestamp=get_utc_now(),
            changes_json=json.dumps({
                "jira_issue_key": issue_key,
                "old_status": old_status,
                "new_status": new_status,
                "webhook_event": webhook_event
            })
        )
        
        session.add(finding)
        session.add(audit_entry)
        session.commit()
        
        logger.info(f"Updated finding {finding.id} Jira status: {old_status} -> {new_status}")
        
        return {
            "status": "success",
            "finding_id": finding.id,
            "old_status": old_status,
            "new_status": new_status
        }
        
    except Exception as e:
        logger.error(f"Error processing Jira webhook: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail="Error processing webhook")

# --- SLA & Remediation Tracking Endpoints ---

@app.get("/findings/sla")
def get_all_findings_with_sla(session: Session = Depends(get_session)):
    """Get all findings with SLA tracking information."""
    statement = select(Finding)
    findings = session.exec(statement).all()
    return [FindingReadWithInstances.model_validate(f) for f in findings]

@app.get("/findings/overdue")
def get_overdue_findings_endpoint(session: Session = Depends(get_session)):
    """Get all findings that are past their SLA deadline."""
    findings = get_overdue_findings(session)
    return [FindingReadWithInstances.model_validate(f) for f in findings]

@app.patch("/findings/{finding_id}/remediation")
def update_finding_remediation(
    finding_id: int,
    data: RemediationUpdate,
    session: Session = Depends(get_session)
):
    """
    Update remediation tracking for a finding.
    
    Args:
        finding_id: ID of the finding
        data: Remediation update data (deadline, owner, user)
    """
    finding = session.get(Finding, finding_id)
    if not finding:
        raise HTTPException(status_code=404, detail="Finding not found")
    
    changes = {}
    
    if data.remediation_deadline:
        try:
            new_deadline = datetime.fromisoformat(data.remediation_deadline.replace('Z', '+00:00'))
            old_deadline = finding.remediation_deadline
            finding.remediation_deadline = new_deadline
            changes["remediation_deadline"] = {
                "old": old_deadline.isoformat() if old_deadline else None,
                "new": new_deadline.isoformat()
            }
        except ValueError:
            raise HTTPException(status_code=400, detail="Invalid datetime format. Use ISO format.")
    
    if data.remediation_owner is not None:
        old_owner = finding.remediation_owner
        finding.remediation_owner = data.remediation_owner
        changes["remediation_owner"] = {
            "old": old_owner,
            "new": data.remediation_owner
        }
    
    # Recalculate SLA status
    finding = update_finding_sla(finding, session)
    
    # Create audit log entry
    if changes:
        audit_entry = AuditLog(
            entity_type="finding",
            entity_id=finding_id,
            action="remediation_updated",
            user=data.user,
            timestamp=get_utc_now(),
            changes_json=json.dumps(changes)
        )
        session.add(audit_entry)
        session.commit()
    
    logger.info(f"Finding {finding_id} remediation updated by {data.user}")
    
    return {
        "id": finding.id,
        "remediation_deadline": finding.remediation_deadline.isoformat() if finding.remediation_deadline else None,
        "remediation_owner": finding.remediation_owner,
        "sla_status": finding.sla_status.value if finding.sla_status else None
    }

@app.patch("/findings/{finding_id}/issue-status")
def update_finding_issue_status(
    finding_id: int,
    data: IssueStatusUpdate,
    session: Session = Depends(get_session)
):
    """
    Update issue status for a finding.
    
    Args:
        finding_id: ID of the finding
        data: Issue status update data (status, comment, user)
    """
    finding = session.get(Finding, finding_id)
    if not finding:
        raise HTTPException(status_code=404, detail="Finding not found")
    
    # Validate issue status
    valid_statuses = ["Open", "Partially Closed", "Closed"]
    if data.issue_status not in valid_statuses:
        raise HTTPException(
            status_code=400, 
            detail=f"Invalid issue status. Must be one of: {', '.join(valid_statuses)}"
        )
    
    changes = {}
    old_status = finding.issue_status.value if finding.issue_status else None
    old_comment = finding.issue_status_comment
    
    # Update status
    finding.issue_status = FindingBase.IssueStatus(data.issue_status)
    changes["issue_status"] = {
        "old": old_status,
        "new": data.issue_status
    }
    
    # Update comment if provided
    if data.issue_status_comment is not None:
        finding.issue_status_comment = data.issue_status_comment
        changes["issue_status_comment"] = {
            "old": old_comment,
            "new": data.issue_status_comment
        }
    
    session.add(finding)
    
    # Create audit log entry
    audit_entry = AuditLog(
        entity_type="finding",
        entity_id=finding_id,
        action="issue_status_updated",
        user=data.user,
        timestamp=get_utc_now(),
        changes_json=json.dumps(changes)
    )
    session.add(audit_entry)
    session.commit()
    session.refresh(finding)
    
    logger.info(f"Finding {finding_id} issue status updated to {data.issue_status} by {data.user}")
    
    return {
        "id": finding.id,
        "issue_status": finding.issue_status.value,
        "issue_status_comment": finding.issue_status_comment
    }

@app.patch("/findings/{finding_id}")
def update_finding(
    finding_id: int,
    data: dict,
    session: Session = Depends(get_session)
):
    """
    Update finding fields.
    
    Supported fields:
    - title, risk_rating, description, remediation
    - impact, references_url, poc_description
    - cwe_id, cve_id, cvss_vector, cvss_score
    - owasp_likelihood, owasp_impact, owasp_risk_rating
    - review_status, sla_status, issue_status
    
    Args:
        finding_id: ID of the finding
        data: Dictionary with fields to update
    """
    finding = session.get(Finding, finding_id)
    if not finding:
        raise HTTPException(status_code=404, detail="Finding not found")
    
    changes = {}
    
    # Update title if provided
    if "title" in data:
        old_title = finding.title
        finding.title = data["title"]
        changes["title"] = {
            "old": old_title,
            "new": data["title"]
        }
    
    # Update risk_rating if provided
    if "risk_rating" in data:
        old_risk = finding.risk_rating.value if finding.risk_rating else None
        new_risk = data["risk_rating"]
        
        # Validate risk rating
        valid_risks = ["Critical", "High", "Medium", "Low", "Informational"]
        if new_risk not in valid_risks:
            raise HTTPException(
                status_code=400,
                detail=f"Invalid risk rating. Must be one of: {', '.join(valid_risks)}"
            )
        
        finding.risk_rating = RiskRating(new_risk)
        
        # Update SLA deadline when risk changes
        from app.sla import calculate_sla_deadline
        finding.remediation_deadline = calculate_sla_deadline(finding.risk_rating.value)
        
        changes["risk_rating"] = {
            "old": old_risk,
            "new": new_risk
        }
    
    # Update description if provided
    if "description" in data:
        old_desc = finding.description
        finding.description = data["description"]
        changes["description"] = {
            "old": old_desc[:100] if old_desc else None,  # Truncate for audit log
            "new": data["description"][:100] if data["description"] else None
        }
    
    # Update impact if provided
    if "impact" in data:
        old_val = finding.impact
        finding.impact = sanitize_html_input(data["impact"]) if data["impact"] is not None else None
        changes["impact"] = {
            "old": (old_val[:100] if old_val else None),
            "new": (finding.impact[:100] if finding.impact else None)
        }
    
    # Update references_url if provided
    if "references_url" in data:
        new_url = data["references_url"] or None
        if new_url:
            try:
                validate_url(new_url, "references_url")
            except HTTPException as e:
                raise e
        old_url = finding.references_url
        finding.references_url = new_url
        changes["references_url"] = {"old": old_url, "new": new_url}
    
    # Update POC description
    if "poc_description" in data:
        old_val = finding.poc_description
        finding.poc_description = sanitize_html_input(data["poc_description"]) if data["poc_description"] is not None else None
        changes["poc_description"] = {
            "old": (old_val[:100] if old_val else None),
            "new": (finding.poc_description[:100] if finding.poc_description else None)
        }
    
    # Update CWE ID if provided
    if "cwe_id" in data:
        old_val = finding.cwe_id
        new_val = data["cwe_id"] or None
        finding.cwe_id = new_val
        changes["cwe_id"] = {"old": old_val, "new": new_val}
    
    # Update CVE ID if provided
    if "cve_id" in data:
        old_val = finding.cve_id
        new_val = data["cve_id"] or None
        finding.cve_id = new_val
        changes["cve_id"] = {"old": old_val, "new": new_val}
    
    # Update CVSS vector if provided
    if "cvss_vector" in data:
        old_val = finding.cvss_vector
        new_val = data["cvss_vector"] or None
        finding.cvss_vector = new_val
        changes["cvss_vector"] = {"old": old_val, "new": new_val}
    
    # Update CVSS score if provided
    if "cvss_score" in data:
        old_val = finding.cvss_score
        new_val = data["cvss_score"]
        if new_val is not None:
            try:
                new_val = float(new_val)
                if new_val < 0.0 or new_val > 10.0:
                    raise ValueError("CVSS score must be between 0.0 and 10.0")
            except (ValueError, TypeError) as e:
                raise HTTPException(status_code=400, detail=f"Invalid CVSS score: {str(e)}")
        finding.cvss_score = new_val
        changes["cvss_score"] = {"old": old_val, "new": new_val}
    
    # Update OWASP likelihood if provided
    if "owasp_likelihood" in data:
        old_val = finding.owasp_likelihood
        new_val = data["owasp_likelihood"]
        if new_val is not None:
            try:
                new_val = int(new_val)
                if new_val < 1 or new_val > 9:
                    raise ValueError("OWASP likelihood must be between 1 and 9")
            except (ValueError, TypeError) as e:
                raise HTTPException(status_code=400, detail=f"Invalid OWASP likelihood: {str(e)}")
        finding.owasp_likelihood = new_val
        changes["owasp_likelihood"] = {"old": old_val, "new": new_val}
    
    # Update OWASP impact if provided
    if "owasp_impact" in data:
        old_val = finding.owasp_impact
        new_val = data["owasp_impact"]
        if new_val is not None:
            try:
                new_val = int(new_val)
                if new_val < 1 or new_val > 9:
                    raise ValueError("OWASP impact must be between 1 and 9")
            except (ValueError, TypeError) as e:
                raise HTTPException(status_code=400, detail=f"Invalid OWASP impact: {str(e)}")
        finding.owasp_impact = new_val
        changes["owasp_impact"] = {"old": old_val, "new": new_val}
    
    # Update OWASP risk rating if provided
    if "owasp_risk_rating" in data:
        old_val = finding.owasp_risk_rating
        new_val = data["owasp_risk_rating"] or None
        finding.owasp_risk_rating = new_val
        changes["owasp_risk_rating"] = {"old": old_val, "new": new_val}
    
    # Update review_status if provided
    if "review_status" in data:
        new_status = data["review_status"]
        valid_statuses = ['Pending', 'In Review', 'Approved', 'Rejected']
        if new_status not in valid_statuses:
            raise HTTPException(status_code=400, detail=f"Invalid review status. Must be one of: {', '.join(valid_statuses)}")
        old_status = finding.review_status.value if finding.review_status else None
        finding.review_status = FindingBase.ReviewStatus(new_status)
        changes["review_status"] = {"old": old_status, "new": new_status}
    
    # Update sla_status if provided
    if "sla_status" in data:
        new_sla = data["sla_status"]
        valid_sla = ['On Track', 'At Risk', 'Overdue', None]
        if new_sla not in ['On Track', 'At Risk', 'Overdue', None, '']:
            raise HTTPException(status_code=400, detail=f"Invalid SLA status. Must be one of: On Track, At Risk, Overdue or null")
        old_sla = finding.sla_status.value if finding.sla_status else None
        finding.sla_status = FindingBase.SLAStatus(new_sla) if new_sla else None
        changes["sla_status"] = {"old": old_sla, "new": new_sla}
    
    # Update issue_status if provided (v0.8.1 - Track resolved_at)
    if "issue_status" in data:
        old_status = finding.issue_status.value if finding.issue_status else None
        new_status = data["issue_status"]
        
        # Validate issue status
        valid_statuses = ["Open", "Partially Closed", "Closed"]
        if new_status not in valid_statuses:
            raise HTTPException(
                status_code=400,
                detail=f"Invalid issue status. Must be one of: {', '.join(valid_statuses)}"
            )
        
        finding.issue_status = FindingBase.IssueStatus(new_status)
        
        # Set resolved_at when status changes to Closed
        if new_status == "Closed" and finding.resolved_at is None:
            finding.resolved_at = get_utc_now()
            changes["resolved_at"] = {
                "old": None,
                "new": "Set to current time"
            }
        # Clear resolved_at if reopened
        elif new_status != "Closed" and finding.resolved_at is not None:
            finding.resolved_at = None
            changes["resolved_at"] = {
                "old": "Previously resolved",
                "new": None
            }
        
        changes["issue_status"] = {
            "old": old_status,
            "new": new_status
        }
    
    session.add(finding)
    
    # Create audit log entry if changes were made
    if changes:
        audit_entry = AuditLog(
            entity_type="finding",
            entity_id=finding_id,
            action="finding_updated",
            user=data.get("user", "system"),
            timestamp=get_utc_now(),
            changes_json=json.dumps(changes)
        )
        session.add(audit_entry)
    
    session.commit()
    session.refresh(finding)
    
    logger.info(f"Finding {finding_id} updated with changes: {list(changes.keys())}")
    
    return {
        "id": finding.id,
        "title": finding.title,
        "risk_rating": finding.risk_rating.value,
        "description": finding.description,
        "impact": finding.impact,
        "references_url": finding.references_url,
        "poc_description": finding.poc_description,
        "remediation_deadline": finding.remediation_deadline.isoformat() if finding.remediation_deadline else None
    }

# --- Instance Management Endpoints ---

@app.patch("/instances/{instance_id}")
@limiter.limit("60/minute")
def update_instance(
    request: Request,
    instance_id: int,
    data: dict,
    session: Session = Depends(get_session)
):
    """
    Update an instance's location, details, or status.
    
    Args:
        instance_id: ID of the instance to update
        data: Dict with fields to update (location, details, status)
    
    Returns:
        Updated instance data
    """
    instance = session.get(Instance, instance_id)
    if not instance:
        raise HTTPException(status_code=404, detail="Instance not found")
    
    # Update location if provided
    if "location" in data and data["location"] is not None:
        instance.location = sanitize_html_input(str(data["location"]))
    
    # Update details if provided
    if "details" in data and data["details"] is not None:
        instance.details = sanitize_html_input(str(data["details"]))
    
    # Update status if provided
    if "status" in data and data["status"] is not None:
        instance.status = sanitize_html_input(str(data["status"]))
    
    session.add(instance)
    session.commit()
    session.refresh(instance)
    
    logger.info(f"Instance {instance_id} updated (finding {instance.finding_id})")
    
    return {
        "id": instance.id,
        "finding_id": instance.finding_id,
        "location": instance.location,
        "details": instance.details,
        "status": instance.status
    }


@app.delete("/instances/{instance_id}", status_code=204)
@limiter.limit("60/minute")
def delete_instance(
    request: Request,
    instance_id: int,
    session: Session = Depends(get_session)
):
    """
    Delete an instance.
    
    Note: Deleting the last instance of a finding will NOT delete the finding itself.
    """
    instance = session.get(Instance, instance_id)
    if not instance:
        raise HTTPException(status_code=404, detail="Instance not found")
    
    finding_id = instance.finding_id
    session.delete(instance)
    session.commit()
    
    logger.info(f"Instance {instance_id} deleted from finding {finding_id}")
    return


@app.post("/findings/{finding_id}/instances", status_code=201)
@limiter.limit("60/minute")
def add_instance(
    request: Request,
    finding_id: int,
    data: dict,
    session: Session = Depends(get_session)
):
    """
    Add a new instance to an existing finding.
    
    Args:
        finding_id: ID of the finding
        data: Dict with location, details, and optional status
    
    Returns:
        Created instance data
    """
    finding = session.get(Finding, finding_id)
    if not finding:
        raise HTTPException(status_code=404, detail="Finding not found")
    
    # Validate required fields
    if "location" not in data or not data["location"]:
        raise HTTPException(status_code=400, detail="location is required")
    if "details" not in data or not data["details"]:
        raise HTTPException(status_code=400, detail="details is required")
    
    # Create new instance
    instance = Instance(
        finding_id=finding_id,
        location=sanitize_html_input(str(data["location"])),
        details=sanitize_html_input(str(data["details"])),
        status=data.get("status", "New - Unvalidated"),
        created_at=get_utc_now()
    )
    
    session.add(instance)
    session.commit()
    session.refresh(instance)
    
    logger.info(f"New instance {instance.id} added to finding {finding_id}")
    
    return {
        "id": instance.id,
        "finding_id": instance.finding_id,
        "location": instance.location,
        "details": instance.details,
        "status": instance.status,
        "created_at": instance.created_at.isoformat() if instance.created_at else None
    }

# --- Finding Artifact (POC Evidence) Endpoints ---

@app.post("/findings/{finding_id}/artifacts", response_model=FindingArtifactRead, status_code=201)
@limiter.limit("30/minute")
async def upload_finding_artifact(
    request: Request,
    finding_id: int,
    file: UploadFile = File(...),
    description: Optional[str] = Body(None),
    session: Session = Depends(get_session)
):
    """Upload a POC evidence image for a finding. Accepts JPEG/PNG up to 5 MiB."""
    finding = session.get(Finding, finding_id)
    if not finding:
        raise HTTPException(status_code=404, detail="Finding not found")
    
    # Validate content type
    allowed_types = {"image/jpeg", "image/png"}
    if file.content_type not in allowed_types:
        raise HTTPException(status_code=400, detail="Only JPEG and PNG images are allowed")
    
    # Read and size-check
    content = await file.read()
    max_bytes = 5 * 1024 * 1024  # 5 MiB
    if len(content) > max_bytes:
        raise HTTPException(status_code=400, detail="File too large (max 5 MiB)")
    
    # Build path: artifacts/<finding_id>/<random>_<sanitized_name>
    safe_name = re.sub(r"[^A-Za-z0-9._-]", "_", file.filename or "artifact")
    random_prefix = secrets.token_hex(8)
    artifact_dir = EVIDENCE_ARTIFACTS_DIR / str(finding_id)
    artifact_dir.mkdir(parents=True, exist_ok=True)
    final_path = artifact_dir / f"{random_prefix}_{safe_name}"
    
    # Write to disk
    with open(final_path, "wb") as f:
        f.write(content)
    
    art = FindingArtifact(
        finding_id=finding_id,
        file_name=file.filename or "artifact",
        file_path=str(final_path.relative_to(EVIDENCE_BASE_DIR)),
        mime_type=file.content_type or "application/octet-stream",
        size_bytes=len(content),
        description=sanitize_html_input(description) if description else None,
        created_at=get_utc_now(),
    )
    session.add(art)
    session.commit()
    session.refresh(art)
    
    return FindingArtifactRead.model_validate(art)


@app.get("/findings/{finding_id}/artifacts", response_model=List[FindingArtifactRead])
def list_finding_artifacts(
    finding_id: int,
    session: Session = Depends(get_session)
):
    finding = session.get(Finding, finding_id)
    if not finding:
        raise HTTPException(status_code=404, detail="Finding not found")
    artifacts = session.exec(select(FindingArtifact).where(FindingArtifact.finding_id == finding_id)).all()
    return [FindingArtifactRead.model_validate(a) for a in artifacts]


@app.get("/artifacts/{artifact_id}/download")
def download_artifact(
    artifact_id: int,
    session: Session = Depends(get_session)
):
    art = session.get(FindingArtifact, artifact_id)
    if not art:
        raise HTTPException(status_code=404, detail="Artifact not found")
    abs_path = (EVIDENCE_BASE_DIR / art.file_path).resolve()
    if not abs_path.exists():
        raise HTTPException(status_code=404, detail="File not found on server")
    return FileResponse(
        path=str(abs_path),
        media_type=art.mime_type,
        filename=art.file_name
    )


@app.delete("/artifacts/{artifact_id}", status_code=204)
def delete_artifact(
    artifact_id: int,
    session: Session = Depends(get_session)
):
    art = session.get(FindingArtifact, artifact_id)
    if not art:
        raise HTTPException(status_code=404, detail="Artifact not found")
    # Try deleting file from disk
    try:
        abs_path = (EVIDENCE_BASE_DIR / art.file_path).resolve()
        if abs_path.exists():
            abs_path.unlink()
    except Exception:
        # Non-fatal: continue to remove DB row
        pass
    session.delete(art)
    session.commit()
    return Response(status_code=204)

@app.get("/sla-summary")
def get_sla_summary_endpoint(session: Session = Depends(get_session)):
    """Get a summary of findings by SLA status."""
    return get_sla_summary(session)

# --- User Preferences & Timezone Endpoints ---

@app.get("/timezones")
def get_available_timezones():
    """Get list of available timezones for user selection."""
    return {
        "default": DEFAULT_TIMEZONE,
        "timezones": TIMEZONE_CHOICES
    }

@app.get("/user-preferences/{user_email}", response_model=Optional[Dict[str, Any]])
def get_user_preferences(
    user_email: str,
    session: Session = Depends(get_session)
):
    """
    Get user preferences including timezone settings.
    
    Args:
        user_email: User's email address
    """
    from app.models import UserPreferences
    
    statement = select(UserPreferences).where(UserPreferences.user_email == user_email)
    prefs = session.exec(statement).first()
    
    if not prefs:
        # Return defaults if no preferences exist
        return {
            "user_email": user_email,
            "timezone": DEFAULT_TIMEZONE,
            "date_format": "%Y-%m-%d %H:%M:%S %Z",
            "locale": "en_MY"
        }
    
    return {
        "id": prefs.id,
        "user_email": prefs.user_email,
        "timezone": prefs.timezone,
        "date_format": prefs.date_format,
        "locale": prefs.locale
    }

@app.post("/user-preferences")
def create_or_update_user_preferences(
    user_email: str = Body(...),
    timezone: str = Body(DEFAULT_TIMEZONE),
    date_format: str = Body("%Y-%m-%d %H:%M:%S %Z"),
    locale: str = Body("en_MY"),
    session: Session = Depends(get_session)
):
    """
    Create or update user preferences.
    
    Args:
        user_email: User's email address
        timezone: User's preferred timezone (IANA format)
        date_format: Preferred date format string
        locale: User's locale preference
    """
    from app.models import UserPreferences
    
    # Validate timezone
    if not is_valid_timezone(timezone):
        raise HTTPException(
            status_code=400,
            detail=f"Invalid timezone: {timezone}. Use IANA timezone names (e.g., 'Asia/Kuala_Lumpur')"
        )
    
    # Check if preferences exist
    statement = select(UserPreferences).where(UserPreferences.user_email == user_email)
    prefs = session.exec(statement).first()
    
    if prefs:
        # Update existing preferences
        prefs.timezone = timezone
        prefs.date_format = date_format
        prefs.locale = locale
    else:
        # Create new preferences
        prefs = UserPreferences(
            user_email=user_email,
            timezone=timezone,
            date_format=date_format,
            locale=locale
        )
        session.add(prefs)
    
    session.commit()
    session.refresh(prefs)
    
    logger.info(f"User preferences updated for {user_email}: timezone={timezone}")
    
    return {
        "id": prefs.id,
        "user_email": prefs.user_email,
        "timezone": prefs.timezone,
        "date_format": prefs.date_format,
        "locale": prefs.locale
    }

# =====================================================
# VULNERABILITY REPOSITORY ENDPOINTS
# =====================================================

@app.get("/vulnerability-templates", response_model=List[VulnerabilityTemplateRead])
def get_vulnerability_templates(
    skip: int = 0,
    limit: int = 100,
    search: Optional[str] = None,
    source: Optional[str] = None,
    risk_rating: Optional[str] = None,
    is_verified: Optional[bool] = None,
    cwe_id: Optional[str] = None,
    cve_id: Optional[str] = None,
    session: Session = Depends(get_session)
):
    """
    Get all vulnerability templates with optional filtering and pagination.
    
    Query parameters:
    - skip: Number of records to skip (default: 0)
    - limit: Maximum number of records to return (default: 100, max: 1000)
    - search: Search in title and description
    - source: Filter by source (manual, burp, nessus, nvd, cwe)
    - risk_rating: Filter by default risk rating
    - is_verified: Filter by verification status
    - cwe_id: Filter by CWE ID
    - cve_id: Filter by CVE ID
    """
    from app.models import VulnerabilityTemplate
    
    # Build query
    statement = select(VulnerabilityTemplate)
    
    # Apply filters
    if search:
        search_pattern = f"%{search}%"
        statement = statement.where(
            (VulnerabilityTemplate.title.ilike(search_pattern)) |
            (VulnerabilityTemplate.description.ilike(search_pattern))
        )
    
    if source:
        statement = statement.where(VulnerabilityTemplate.source == source)
    
    if risk_rating:
        statement = statement.where(VulnerabilityTemplate.default_risk_rating == risk_rating)
    
    if is_verified is not None:
        statement = statement.where(VulnerabilityTemplate.is_verified == is_verified)
    
    if cwe_id:
        statement = statement.where(VulnerabilityTemplate.cwe_id == cwe_id)
    
    if cve_id:
        statement = statement.where(VulnerabilityTemplate.cve_id == cve_id)
    
    # Apply pagination and ordering
    statement = statement.order_by(VulnerabilityTemplate.created_at.desc())
    statement = statement.offset(skip).limit(min(limit, 1000))
    
    templates = session.exec(statement).all()
    return templates


@app.post("/vulnerability-templates", response_model=VulnerabilityTemplateRead, status_code=201)
@limiter.limit("30/hour")  # Max 30 templates per hour per IP
async def create_vulnerability_template(
    request: Request,
    title: str = Body(...),
    description: str = Body(...),
    cwe_id: Optional[str] = Body(None),
    cve_id: Optional[str] = Body(None),
    cvss_vector: Optional[str] = Body(None),
    cvss_score: Optional[float] = Body(None),
    owasp_likelihood: Optional[int] = Body(None),
    owasp_impact: Optional[int] = Body(None),
    owasp_risk_rating: Optional[str] = Body(None),
    default_risk_rating: Optional[str] = Body(None),
    vulnerability_type: Optional[str] = Body(None),
    remediation_summary: Optional[str] = Body(None),
    remediation_steps: Optional[str] = Body(None),
    references: Optional[str] = Body(None),
    is_verified: bool = Body(True),
    session: Session = Depends(get_session)
):
    """
    Create a new vulnerability template.
    
    All templates created manually are marked with source='manual' and is_verified=True by default.
    """
    from app.models import VulnerabilityTemplate
    
    # Validate and sanitize inputs
    title = validate_string_length(sanitize_html_input(title), "title", max_length=200)
    description = validate_string_length(sanitize_html_input(description), "description", max_length=5000)
    if remediation_summary:
        remediation_summary = validate_string_length(sanitize_html_input(remediation_summary), "remediation summary", max_length=2000, allow_empty=True)
    if remediation_steps:
        remediation_steps = validate_string_length(sanitize_html_input(remediation_steps), "remediation steps", max_length=5000, allow_empty=True)
    
    # Validate CVSS score range
    if cvss_score is not None and (cvss_score < 0.0 or cvss_score > 10.0):
        raise HTTPException(status_code=400, detail="CVSS score must be between 0.0 and 10.0")
    
    # Validate OWASP ranges
    if owasp_likelihood is not None and (owasp_likelihood < 1 or owasp_likelihood > 9):
        raise HTTPException(status_code=400, detail="OWASP likelihood must be between 1 and 9")
    
    if owasp_impact is not None and (owasp_impact < 1 or owasp_impact > 9):
        raise HTTPException(status_code=400, detail="OWASP impact must be between 1 and 9")
    
    # Check for duplicate templates (same title + CWE/CVE combination)
    duplicate_check = select(VulnerabilityTemplate).where(
        VulnerabilityTemplate.title == title
    )
    
    # If CWE or CVE provided, use them for stricter duplicate detection
    if cwe_id:
        duplicate_check = duplicate_check.where(VulnerabilityTemplate.cwe_id == cwe_id)
    if cve_id:
        duplicate_check = duplicate_check.where(VulnerabilityTemplate.cve_id == cve_id)
    
    existing_template = session.exec(duplicate_check).first()
    
    if existing_template:
        raise HTTPException(
            status_code=409,
            detail=f"A template with this title{' and CWE ID' if cwe_id else ''}{' and CVE ID' if cve_id else ''} already exists (ID: {existing_template.id})"
        )
    
    # Create template
    template = VulnerabilityTemplate(
        title=title,
        description=description,
        cwe_id=cwe_id,
        cve_id=cve_id,
        cvss_vector=cvss_vector,
        cvss_score=cvss_score,
        owasp_likelihood=owasp_likelihood,
        owasp_impact=owasp_impact,
        owasp_risk_rating=owasp_risk_rating,
        default_risk_rating=default_risk_rating,
        vulnerability_type=vulnerability_type,
        remediation_summary=remediation_summary,
        remediation_steps=remediation_steps,
        references=references,
        source="manual",
        is_verified=is_verified,
        usage_count=0,
        created_at=get_utc_now(),
        updated_at=get_utc_now()
    )
    
    session.add(template)
    session.commit()
    session.refresh(template)
    
    logger.info(f"Created vulnerability template: {template.id} - {template.title}")
    
    return template


@app.get("/vulnerability-templates/analytics")
def get_template_analytics(
    session: Session = Depends(get_session)
):
    """
    Get comprehensive analytics about the vulnerability template repository.
    
    Returns statistics on:
    - Total template count
    - Distribution by source (manual/burp/nessus/nvd/cwe)
    - Distribution by risk rating
    - Top 10 most-used templates
    - Data quality metrics (CVSS coverage, CWE coverage, verification rate)
    - ATT&CK technique usage statistics
    """
    from collections import defaultdict
    import json
    
    # Total templates
    total_templates = session.exec(
        select(func.count(VulnerabilityTemplate.id))
    ).one()
    
    # Templates by source
    by_source_raw = session.exec(
        select(
            VulnerabilityTemplate.source,
            func.count(VulnerabilityTemplate.id)
        ).group_by(VulnerabilityTemplate.source)
    ).all()
    by_source = {source or "unknown": count for source, count in by_source_raw}
    
    # Templates by risk rating
    by_risk_raw = session.exec(
        select(
            VulnerabilityTemplate.default_risk_rating,
            func.count(VulnerabilityTemplate.id)
        ).group_by(VulnerabilityTemplate.default_risk_rating)
    ).all()
    by_risk = {risk or "None": count for risk, count in by_risk_raw}
    
    # Most used templates (top 10)
    most_used_templates = session.exec(
        select(VulnerabilityTemplate)
        .order_by(VulnerabilityTemplate.usage_count.desc())
        .limit(10)
    ).all()
    
    most_used = [
        {
            "template_id": t.id,
            "title": t.title,
            "usage_count": t.usage_count,
            "risk_rating": t.default_risk_rating
        }
        for t in most_used_templates
    ]
    
    # Quality metrics
    with_cvss = session.exec(
        select(func.count(VulnerabilityTemplate.id))
        .where(VulnerabilityTemplate.cvss_score.isnot(None))
    ).one()
    
    with_cwe = session.exec(
        select(func.count(VulnerabilityTemplate.id))
        .where(VulnerabilityTemplate.cwe_id.isnot(None))
    ).one()
    
    verified = session.exec(
        select(func.count(VulnerabilityTemplate.id))
        .where(VulnerabilityTemplate.is_verified == True)
    ).one()
    
    # Calculate percentages
    cvss_pct = (with_cvss / total_templates * 100) if total_templates > 0 else 0
    cwe_pct = (with_cwe / total_templates * 100) if total_templates > 0 else 0
    verified_pct = (verified / total_templates * 100) if total_templates > 0 else 0
    
    # ATT&CK technique statistics
    templates_with_attack = session.exec(
        select(func.count(VulnerabilityTemplate.id))
        .where(VulnerabilityTemplate.attack_techniques.isnot(None))
        .where(VulnerabilityTemplate.attack_techniques != "null")
        .where(VulnerabilityTemplate.attack_techniques != "[]")
    ).one()
    
    # Count total techniques and tactic distribution
    all_templates_with_attack = session.exec(
        select(VulnerabilityTemplate.attack_techniques)
        .where(VulnerabilityTemplate.attack_techniques.isnot(None))
        .where(VulnerabilityTemplate.attack_techniques != "null")
        .where(VulnerabilityTemplate.attack_techniques != "[]")
    ).all()
    
    total_techniques_mapped = 0
    tactic_counts = defaultdict(int)
    
    for attack_json in all_templates_with_attack:
        try:
            techniques = json.loads(attack_json)
            if isinstance(techniques, list):
                total_techniques_mapped += len(techniques)
                for tech in techniques:
                    if isinstance(tech, dict) and "tactic" in tech:
                        tactic_counts[tech["tactic"]] += 1
        except (json.JSONDecodeError, TypeError):
            continue
    
    # Sort tactics by count
    most_common_tactics = [
        {"tactic": tactic, "count": count}
        for tactic, count in sorted(tactic_counts.items(), key=lambda x: x[1], reverse=True)
    ]
    
    return {
        "total_templates": total_templates,
        "by_source": by_source,
        "by_risk_rating": by_risk,
        "most_used": most_used,
        "quality_metrics": {
            "with_cvss": with_cvss,
            "with_cwe": with_cwe,
            "verified": verified,
            "cvss_coverage_pct": round(cvss_pct, 1),
            "cwe_coverage_pct": round(cwe_pct, 1),
            "verification_rate_pct": round(verified_pct, 1)
        },
        "attack_techniques": {
            "templates_with_attack": templates_with_attack,
            "total_techniques_mapped": total_techniques_mapped,
            "most_common_tactics": most_common_tactics[:5]  # Top 5 tactics
        }
    }


@app.get("/vulnerability-templates/{template_id}", response_model=VulnerabilityTemplateRead)
def get_vulnerability_template(
    template_id: int,
    session: Session = Depends(get_session)
):
    """Get a single vulnerability template by ID."""
    from app.models import VulnerabilityTemplate
    
    template = session.get(VulnerabilityTemplate, template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Vulnerability template not found")
    
    return template


@app.patch("/vulnerability-templates/{template_id}", response_model=VulnerabilityTemplateRead)
async def update_vulnerability_template(
    template_id: int,
    title: Optional[str] = Body(None),
    description: Optional[str] = Body(None),
    cwe_id: Optional[str] = Body(None),
    cve_id: Optional[str] = Body(None),
    cvss_vector: Optional[str] = Body(None),
    cvss_score: Optional[float] = Body(None),
    owasp_likelihood: Optional[int] = Body(None),
    owasp_impact: Optional[int] = Body(None),
    owasp_risk_rating: Optional[str] = Body(None),
    default_risk_rating: Optional[str] = Body(None),
    vulnerability_type: Optional[str] = Body(None),
    remediation_summary: Optional[str] = Body(None),
    remediation_steps: Optional[str] = Body(None),
    references: Optional[str] = Body(None),
    is_verified: Optional[bool] = Body(None),
    changed_by: Optional[str] = Body(None),  # Who is making this change
    change_reason: Optional[str] = Body(None),  # Why is this change being made
    session: Session = Depends(get_session)
):
    """
    Update an existing vulnerability template.
    Only provided fields will be updated.
    Automatically creates a version snapshot before updating.
    """
    from app.models import VulnerabilityTemplate, VulnerabilityTemplateVersion
    from sqlmodel import func, select
    
    template = session.get(VulnerabilityTemplate, template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Vulnerability template not found")
    
    # CREATE VERSION SNAPSHOT BEFORE UPDATING
    # Get current max version number for this template
    max_version = session.exec(
        select(func.max(VulnerabilityTemplateVersion.version_number))
        .where(VulnerabilityTemplateVersion.template_id == template_id)
    ).first() or 0
    
    # Create version snapshot with current state
    version_snapshot = VulnerabilityTemplateVersion(
        template_id=template.id,
        version_number=max_version + 1,
        title=template.title,
        description=template.description,
        cwe_id=template.cwe_id,
        cve_id=template.cve_id,
        cvss_vector=template.cvss_vector,
        cvss_score=template.cvss_score,
        owasp_likelihood=template.owasp_likelihood,
        owasp_impact=template.owasp_impact,
        owasp_risk_rating=template.owasp_risk_rating,
        default_risk_rating=template.default_risk_rating,
        vulnerability_type=template.vulnerability_type,
        remediation_summary=template.remediation_summary,
        remediation_steps=template.remediation_steps,
        references=template.references,
        attack_techniques=template.attack_techniques,
        source=template.source,
        is_verified=template.is_verified,
        changed_by=changed_by or "system",
        change_reason=change_reason or "Template updated",
        created_at=get_utc_now()
    )
    session.add(version_snapshot)
    logger.info(f"Created version snapshot {max_version + 1} for template {template_id}")
    
    # NOW UPDATE THE TEMPLATE
    # Update fields if provided
    if title is not None:
        template.title = title
    if description is not None:
        template.description = description
    if cwe_id is not None:
        template.cwe_id = cwe_id
    if cve_id is not None:
        template.cve_id = cve_id
    if cvss_vector is not None:
        template.cvss_vector = cvss_vector
    if cvss_score is not None:
        if cvss_score < 0.0 or cvss_score > 10.0:
            raise HTTPException(status_code=400, detail="CVSS score must be between 0.0 and 10.0")
        template.cvss_score = cvss_score
    if owasp_likelihood is not None:
        if owasp_likelihood < 1 or owasp_likelihood > 9:
            raise HTTPException(status_code=400, detail="OWASP likelihood must be between 1 and 9")
        template.owasp_likelihood = owasp_likelihood
    if owasp_impact is not None:
        if owasp_impact < 1 or owasp_impact > 9:
            raise HTTPException(status_code=400, detail="OWASP impact must be between 1 and 9")
        template.owasp_impact = owasp_impact
    if owasp_risk_rating is not None:
        template.owasp_risk_rating = owasp_risk_rating
    if default_risk_rating is not None:
        template.default_risk_rating = default_risk_rating
    if vulnerability_type is not None:
        template.vulnerability_type = vulnerability_type
    if remediation_summary is not None:
        template.remediation_summary = remediation_summary
    if remediation_steps is not None:
        template.remediation_steps = remediation_steps
    if references is not None:
        template.references = references
    if is_verified is not None:
        template.is_verified = is_verified
    
    # Update timestamp
    template.updated_at = get_utc_now()
    
    session.add(template)
    session.commit()
    session.refresh(template)
    
    logger.info(f"Updated vulnerability template: {template.id} - {template.title}")
    
    return template


@app.get("/vulnerability-templates/{template_id}/versions", response_model=List[VulnerabilityTemplateVersionRead])
def get_template_version_history(
    template_id: int,
    session: Session = Depends(get_session)
):
    """
    Get complete version history for a template.
    Returns versions in chronological order (oldest to newest).
    """
    from app.models import VulnerabilityTemplate, VulnerabilityTemplateVersion
    
    # Verify template exists
    template = session.get(VulnerabilityTemplate, template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Vulnerability template not found")
    
    # Get all versions ordered by version number
    versions = session.exec(
        select(VulnerabilityTemplateVersion)
        .where(VulnerabilityTemplateVersion.template_id == template_id)
        .order_by(VulnerabilityTemplateVersion.version_number)
    ).all()
    
    logger.info(f"Retrieved {len(versions)} versions for template {template_id}")
    
    return list(versions)


@app.post("/vulnerability-templates/{template_id}/rollback/{version_number}", response_model=VulnerabilityTemplateRead)
def rollback_template_to_version(
    template_id: int,
    version_number: int,
    changed_by: Optional[str] = Body(None),
    change_reason: Optional[str] = Body("Rolled back to previous version"),
    session: Session = Depends(get_session)
):
    """
    Rollback a template to a specific version.
    Creates a new version snapshot before rolling back.
    """
    from app.models import VulnerabilityTemplate, VulnerabilityTemplateVersion
    from sqlmodel import func
    
    # Verify template exists
    template = session.get(VulnerabilityTemplate, template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Vulnerability template not found")
    
    # Find the target version
    target_version = session.exec(
        select(VulnerabilityTemplateVersion)
        .where(
            VulnerabilityTemplateVersion.template_id == template_id,
            VulnerabilityTemplateVersion.version_number == version_number
        )
    ).first()
    
    if not target_version:
        raise HTTPException(
            status_code=404,
            detail=f"Version {version_number} not found for template {template_id}"
        )
    
    # CREATE SNAPSHOT OF CURRENT STATE BEFORE ROLLBACK
    max_version = session.exec(
        select(func.max(VulnerabilityTemplateVersion.version_number))
        .where(VulnerabilityTemplateVersion.template_id == template_id)
    ).first() or 0
    
    current_snapshot = VulnerabilityTemplateVersion(
        template_id=template.id,
        version_number=max_version + 1,
        title=template.title,
        description=template.description,
        cwe_id=template.cwe_id,
        cve_id=template.cve_id,
        cvss_vector=template.cvss_vector,
        cvss_score=template.cvss_score,
        owasp_likelihood=template.owasp_likelihood,
        owasp_impact=template.owasp_impact,
        owasp_risk_rating=template.owasp_risk_rating,
        default_risk_rating=template.default_risk_rating,
        vulnerability_type=template.vulnerability_type,
        remediation_summary=template.remediation_summary,
        remediation_steps=template.remediation_steps,
        references=template.references,
        attack_techniques=template.attack_techniques,
        source=template.source,
        is_verified=template.is_verified,
        changed_by=changed_by or "system",
        change_reason=f"Before rollback to v{version_number}",
        created_at=get_utc_now()
    )
    session.add(current_snapshot)
    
    # RESTORE FROM TARGET VERSION
    template.title = target_version.title
    template.description = target_version.description
    template.cwe_id = target_version.cwe_id
    template.cve_id = target_version.cve_id
    template.cvss_vector = target_version.cvss_vector
    template.cvss_score = target_version.cvss_score
    template.owasp_likelihood = target_version.owasp_likelihood
    template.owasp_impact = target_version.owasp_impact
    template.owasp_risk_rating = target_version.owasp_risk_rating
    template.default_risk_rating = target_version.default_risk_rating
    template.vulnerability_type = target_version.vulnerability_type
    template.remediation_summary = target_version.remediation_summary
    template.remediation_steps = target_version.remediation_steps
    template.references = target_version.references
    template.attack_techniques = target_version.attack_techniques
    template.source = target_version.source
    template.is_verified = target_version.is_verified
    template.updated_at = get_utc_now()
    
    session.add(template)
    session.commit()
    session.refresh(template)
    
    logger.info(f"Rolled back template {template_id} to version {version_number}")
    
    return template


@app.delete("/vulnerability-templates/{template_id}", status_code=204)
def delete_vulnerability_template(
    template_id: int,
    session: Session = Depends(get_session)
):
    """
    Delete a vulnerability template.
    Will fail if template is in use by findings (prevent deletion).
    """
    from app.models import VulnerabilityTemplate, Finding, VulnerabilityTemplateVersion
    
    template = session.get(VulnerabilityTemplate, template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Vulnerability template not found")
    
    # Check if template is in use
    findings_using_template = session.exec(
        select(Finding).where(Finding.template_id == template_id)
    ).first()
    
    if findings_using_template:
        raise HTTPException(
            status_code=400,
            detail=f"Cannot delete template: {template.usage_count} finding(s) are using this template. Unlink them first."
        )
    
    # Delete version history first (to avoid FK constraint issues)
    versions = session.exec(
        select(VulnerabilityTemplateVersion).where(
            VulnerabilityTemplateVersion.template_id == template_id
        )
    ).all()
    for version in versions:
        session.delete(version)
    
    session.delete(template)
    session.commit()
    
    logger.info(f"Deleted vulnerability template: {template_id}")


@app.post("/vulnerability-templates/bulk-delete")
def bulk_delete_vulnerability_templates(
    template_ids: List[int] = Body(...),
    session: Session = Depends(get_session)
):
    """
    Bulk delete multiple vulnerability templates.
    Will skip templates that are in use by findings.
    
    Returns summary with deleted count and errors for templates that couldn't be deleted.
    """
    from app.models import VulnerabilityTemplate, Finding, VulnerabilityTemplateVersion
    
    deleted = []
    errors = []
    
    for template_id in template_ids:
        template = session.get(VulnerabilityTemplate, template_id)
        
        if not template:
            errors.append({
                "id": template_id,
                "error": "Template not found"
            })
            continue
        
        # Check if template is in use
        findings_using_template = session.exec(
            select(Finding).where(Finding.template_id == template_id)
        ).first()
        
        if findings_using_template:
            errors.append({
                "id": template_id,
                "title": template.title,
                "error": f"Template in use by {template.usage_count} finding(s)"
            })
            continue
        
        # Delete version history first (to avoid FK constraint issues)
        versions = session.exec(
            select(VulnerabilityTemplateVersion).where(
                VulnerabilityTemplateVersion.template_id == template_id
            )
        ).all()
        for version in versions:
            session.delete(version)
        
        # Now safe to delete the template
        session.delete(template)
        deleted.append({
            "id": template_id,
            "title": template.title
        })
    
    session.commit()
    
    logger.info(f"Bulk delete: {len(deleted)} deleted, {len(errors)} errors")
    
    return {
        "deleted_count": len(deleted),
        "deleted": deleted,
        "error_count": len(errors),
        "errors": errors
    }


@app.post("/vulnerability-templates/bulk-update")
def bulk_update_vulnerability_templates(
    updates: List[Dict[str, Any]] = Body(...),
    changed_by: Optional[str] = Body(None),
    change_reason: Optional[str] = Body("Bulk update"),
    session: Session = Depends(get_session)
):
    """
    Bulk update multiple vulnerability templates.
    Each update object should have 'id' and the fields to update.
    
    Creates version snapshots for each template before updating.
    
    Example: [{"id": 1, "is_verified": true}, {"id": 2, "default_risk_rating": "High"}]
    
    Returns summary with updated count and errors.
    """
    from app.models import VulnerabilityTemplate, VulnerabilityTemplateVersion
    from sqlmodel import func
    
    updated = []
    errors = []
    
    for update_data in updates:
        if "id" not in update_data:
            errors.append({"error": "Missing 'id' field in update object"})
            continue
        
        template_id = update_data.pop("id")
        template = session.get(VulnerabilityTemplate, template_id)
        
        if not template:
            errors.append({
                "id": template_id,
                "error": "Template not found"
            })
            continue
        
        # CREATE VERSION SNAPSHOT BEFORE UPDATING (same logic as PATCH endpoint)
        max_version = session.exec(
            select(func.max(VulnerabilityTemplateVersion.version_number))
            .where(VulnerabilityTemplateVersion.template_id == template_id)
        ).first() or 0
        
        version_snapshot = VulnerabilityTemplateVersion(
            template_id=template.id,
            version_number=max_version + 1,
            title=template.title,
            description=template.description,
            cwe_id=template.cwe_id,
            cve_id=template.cve_id,
            cvss_vector=template.cvss_vector,
            cvss_score=template.cvss_score,
            owasp_likelihood=template.owasp_likelihood,
            owasp_impact=template.owasp_impact,
            owasp_risk_rating=template.owasp_risk_rating,
            default_risk_rating=template.default_risk_rating,
            vulnerability_type=template.vulnerability_type,
            remediation_summary=template.remediation_summary,
            remediation_steps=template.remediation_steps,
            references=template.references,
            attack_techniques=template.attack_techniques,
            source=template.source,
            is_verified=template.is_verified,
            changed_by=changed_by or "bulk_update",
            change_reason=change_reason or "Bulk update",
            created_at=get_utc_now()
        )
        session.add(version_snapshot)
        
        # Apply updates
        for key, value in update_data.items():
            if hasattr(template, key):
                setattr(template, key, value)
            else:
                errors.append({
                    "id": template_id,
                    "error": f"Invalid field: {key}"
                })
                continue
        
        template.updated_at = get_utc_now()
        session.add(template)
        updated.append({
            "id": template_id,
            "title": template.title
        })
    
    session.commit()
    
    logger.info(f"Bulk update: {len(updated)} updated, {len(errors)} errors")
    
    return {
        "updated_count": len(updated),
        "updated": updated,
        "error_count": len(errors),
        "errors": errors
    }


@app.post("/vulnerability-templates/cleanup-duplicates")
def cleanup_duplicate_templates(
    session: Session = Depends(get_session)
):
    """
    Remove duplicate vulnerability templates, keeping the most recent version
    with the highest usage count for each title+CWE/CVE combination.
    
    Returns summary of templates removed.
    """
    from app.models import VulnerabilityTemplate
    from collections import defaultdict
    
    # Get all templates
    all_templates = session.exec(select(VulnerabilityTemplate)).all()
    
    # Group by title + CWE + CVE combination
    groups = defaultdict(list)
    for template in all_templates:
        key = (template.title, template.cwe_id or '', template.cve_id or '')
        groups[key].append(template)
    
    removed_count = 0
    removed_ids = []
    kept_templates = []
    
    # Process each group
    for key, templates_in_group in groups.items():
        if len(templates_in_group) > 1:
            # Sort by: usage_count DESC, created_at DESC
            templates_in_group.sort(
                key=lambda t: (t.usage_count, t.created_at),
                reverse=True
            )
            
            # Keep the first one (highest usage, most recent)
            keep = templates_in_group[0]
            kept_templates.append(keep.id)
            
            # Remove the rest
            for template in templates_in_group[1:]:
                # Only remove if usage_count is 0 (not linked to findings)
                if template.usage_count == 0:
                    session.delete(template)
                    removed_ids.append(template.id)
                    removed_count += 1
                    logger.info(f"Removed duplicate template: {template.id} - {template.title}")
    
    session.commit()
    
    return {
        "removed_count": removed_count,
        "removed_ids": removed_ids,
        "kept_templates": kept_templates,
        "message": f"Removed {removed_count} duplicate template(s)"
    }


@app.post("/vulnerability-templates/{template_id}/enrich")
async def enrich_template_from_nvd(
    template_id: int,
    overwrite_existing: bool = Query(False, description="Overwrite existing fields with NVD data"),
    session: Session = Depends(get_session)
):
    """
    Enrich a vulnerability template with data from the National Vulnerability Database (NVD).
    
    Fetches official CVE data including:
    - Official description from NIST
    - CVSS 3.1 score and vector
    - CWE mapping
    - Severity rating
    - Official references
    
    Args:
        template_id: Template to enrich
        overwrite_existing: If False, only populate empty fields. If True, overwrite all.
    
    Returns:
        Updated template with NVD data
    
    Examples:
        POST /vulnerability-templates/5/enrich
        → Fetches CVE data and updates template
    """
    from app.nvd import enrich_template_from_nvd, NVDAPIError
    
    # Get template
    template = session.get(VulnerabilityTemplate, template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    
    # Check if template has CVE ID
    if not template.cve_id:
        raise HTTPException(
            status_code=400,
            detail="Template must have a CVE ID to enrich from NVD"
        )
    
    try:
        # Fetch enrichment data from NVD
        logger.info(f"Enriching template {template_id} from NVD: {template.cve_id}")
        enrichment_data = await enrich_template_from_nvd(
            template.cve_id,
            overwrite_existing=overwrite_existing
        )
        
        # Apply updates to template
        for field, value in enrichment_data.items():
            # Skip if field already has value and overwrite_existing=False
            if not overwrite_existing and getattr(template, field, None):
                continue
            
            # Update field
            setattr(template, field, value)
        
        session.add(template)
        session.commit()
        session.refresh(template)
        
        logger.info(f"Successfully enriched template {template_id} with NVD data")
        
        # Return updated template
        return VulnerabilityTemplateRead.model_validate(template)
    
    except NVDAPIError as e:
        logger.error(f"NVD API error: {str(e)}")
        raise HTTPException(
            status_code=502,
            detail=f"Failed to fetch data from NVD: {str(e)}"
        )
    except Exception as e:
        logger.error(f"Error enriching template: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Internal error while enriching template: {str(e)}"
        )


# =============================================================================
# CWE Database Import (v0.7.1)
# =============================================================================

@app.post("/vulnerability-templates/import-cwe-database")
async def import_cwe_database(
    file: UploadFile = File(..., description="CWE XML file (cwec_latest.xml)"),
    overwrite_existing: bool = Query(False, description="Overwrite existing CWE templates"),
    session: Session = Depends(get_session)
):
    """
    Bulk import CWE (Common Weakness Enumeration) database from MITRE XML file.
    
    Downloads: https://cwe.mitre.org/data/downloads.html
    Latest: https://cwe.mitre.org/data/xml/cwec_latest.xml.zip
    
    Process:
    1. Upload CWE XML file (cwec_latest.xml)
    2. Parse all <Weakness> elements (~900 CWEs)
    3. Create VulnerabilityTemplate for each CWE
    4. Skip existing CWEs (unless overwrite=true)
    5. Return import statistics
    
    Args:
        file: CWE XML file upload
        overwrite_existing: Whether to update existing CWE templates
        session: Database session
    
    Returns:
        Import statistics: {total_parsed, templates_created, templates_skipped, errors}
    
    Example:
        curl -X POST "http://localhost:8000/api/vulnerability-templates/import-cwe-database" \\
             -F "file=@cwec_latest.xml" \\
             -F "overwrite_existing=false"
    
    Response:
        {
          "total_parsed": 922,
          "templates_created": 845,
          "templates_skipped": 77,
          "errors": 0,
          "success_rate": 91.65,
          "imported_at": "2025-11-06T12:34:56.789012"
        }
    """
    from app.cwe import parse_cwe_xml, generate_import_statistics, CWEParseError
    from app.models import VulnerabilityTemplate
    import time
    
    # Start timing
    start_time = time.time()
    error_list = []  # Track errors for history
    
    # Validate file type
    if not file.filename.endswith('.xml'):
        raise HTTPException(
            status_code=400,
            detail="File must be an XML file (.xml extension)"
        )
    
    # Read file content
    try:
        logger.info(f"Reading CWE XML file: {file.filename}")
        xml_content = await file.read()
        
        if len(xml_content) == 0:
            raise HTTPException(status_code=400, detail="Uploaded file is empty")
        
        if len(xml_content) > 50 * 1024 * 1024:  # 50 MB limit
            raise HTTPException(
                status_code=413,
                detail="File too large. Maximum size is 50 MB."
            )
    
    except Exception as e:
        logger.error(f"Error reading uploaded file: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to read uploaded file: {str(e)}"
        )
    
    # Parse CWE XML
    try:
        logger.info("Parsing CWE XML content...")
        cwe_list = parse_cwe_xml(xml_content)
        total_parsed = len(cwe_list)
        
        if total_parsed == 0:
            raise HTTPException(
                status_code=400,
                detail="No CWE weaknesses found in XML file. Ensure you uploaded cwec_latest.xml from MITRE."
            )
        
        logger.info(f"Successfully parsed {total_parsed} CWE entries")
    
    except CWEParseError as e:
        logger.error(f"CWE parsing error: {str(e)}")
        raise HTTPException(
            status_code=400,
            detail=f"Failed to parse CWE XML: {str(e)}"
        )
    
    # Import CWEs into database
    created_count = 0
    updated_count = 0
    skipped_count = 0
    error_count = 0
    
    for cwe_data in cwe_list:
        try:
            cwe_id = cwe_data['cwe_id']
            
            # Check if CWE template already exists
            existing = session.exec(
                select(VulnerabilityTemplate).where(
                    VulnerabilityTemplate.cwe_id == cwe_id
                )
            ).first()
            
            if existing:
                if overwrite_existing:
                    # Update existing template with CWE data
                    for key, value in cwe_data.items():
                        if value is not None:  # Only update non-null values
                            setattr(existing, key, value)
                    existing.updated_at = get_utc_now()
                    session.add(existing)
                    updated_count += 1
                    logger.debug(f"Updated existing {cwe_id}")
                else:
                    # Skip existing CWE
                    skipped_count += 1
                    logger.debug(f"Skipped existing {cwe_id}")
                continue
            
            # Create new template
            template = VulnerabilityTemplate(**cwe_data)
            template.created_at = get_utc_now()
            template.updated_at = get_utc_now()
            session.add(template)
            created_count += 1
            logger.debug(f"Created template for {cwe_id}")
        
        except Exception as e:
            error_count += 1
            error_msg = f"Failed to import {cwe_data.get('cwe_id', 'unknown')}: {str(e)}"
            logger.warning(error_msg)
            error_list.append({
                "cwe_id": cwe_data.get('cwe_id', 'unknown'),
                "error": str(e)
            })
            continue
    
    # Commit all changes
    try:
        session.commit()
        logger.info(
            f"CWE import complete: {created_count} created, {updated_count} updated, "
            f"{skipped_count} skipped, {error_count} errors"
        )
    except Exception as e:
        session.rollback()
        logger.error(f"Database commit failed: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to save CWE templates to database: {str(e)}"
        )
    
    # Generate statistics
    stats = generate_import_statistics(
        total_parsed=total_parsed,
        created=created_count,
        skipped=skipped_count,
        errors=error_count
    )
    
    # Calculate duration
    duration = round(time.time() - start_time, 2)
    
    # Track import in history
    try:
        import json
        import_record = ImportHistory(
            source="cwe",
            import_type="bulk_cwe",
            file_name=file.filename,
            file_size=len(xml_content),
            templates_created=created_count,
            templates_updated=updated_count,
            templates_skipped=skipped_count,
            errors=error_count,
            total_parsed=total_parsed,
            imported_by="system",  # TODO(v1.0.0): Get from current_user.email when auth is implemented
            imported_at=get_utc_now(),
            duration_seconds=duration,
            error_details=json.dumps(error_list) if error_list else None
        )
        session.add(import_record)
        session.commit()
        logger.info(f"Import history record created: ID {import_record.id}, duration: {duration}s")
    except Exception as e:
        # Don't fail the whole import if history tracking fails
        logger.warning(f"Failed to create import history record: {str(e)}")
        session.rollback()
    
    return stats


@app.get("/cwe/{cwe_id}")
def get_cwe_details(
    cwe_id: str = Path(..., description="CWE ID (e.g., CWE-79 or just 79)"),
    session: Session = Depends(get_session)
):
    """
    Get details for a specific CWE from local database or MITRE website.
    
    Searches local vulnerability templates first. If not found locally,
    redirects to official MITRE CWE page.
    
    Args:
        cwe_id: CWE identifier (CWE-79 or 79)
        session: Database session
    
    Returns:
        CWE template details or redirect to MITRE
    
    Example:
        GET /api/cwe/79
        GET /api/cwe/CWE-79
    """
    from app.models import VulnerabilityTemplate
    
    # Normalize CWE ID format
    if not cwe_id.upper().startswith('CWE-'):
        cwe_id = f"CWE-{cwe_id}"
    else:
        cwe_id = cwe_id.upper()
    
    # Search local database
    template = session.exec(
        select(VulnerabilityTemplate).where(
            VulnerabilityTemplate.cwe_id == cwe_id
        )
    ).first()
    
    if template:
        return VulnerabilityTemplateRead.model_validate(template)
    
    # If not found locally, return MITRE URL
    cwe_number = cwe_id.replace('CWE-', '')
    mitre_url = f"https://cwe.mitre.org/data/definitions/{cwe_number}.html"
    
    return {
        "cwe_id": cwe_id,
        "found_locally": False,
        "mitre_url": mitre_url,
        "message": f"{cwe_id} not found in local database. Import CWE database or visit MITRE URL."
    }


# =============================================================================
# Import History Endpoints (v0.7.2)
# =============================================================================

@app.get("/import-history", response_model=List[ImportHistoryRead])
def list_import_history(
    skip: int = Query(0, ge=0, description="Number of records to skip"),
    limit: int = Query(50, ge=1, le=200, description="Maximum number of records to return"),
    source: Optional[str] = Query(None, description="Filter by source (cwe, nvd, manual)"),
    session: Session = Depends(get_session)
):
    """
    List all vulnerability database import history records.
    
    Returns a paginated list of import operations (CWE/CVE imports) with statistics.
    
    Args:
        skip: Pagination offset (default: 0)
        limit: Maximum records to return (default: 50, max: 200)
        source: Filter by import source (optional)
        session: Database session
    
    Returns:
        List of import history records with statistics
    
    Example:
        GET /api/import-history
        GET /api/import-history?source=cwe
        GET /api/import-history?limit=10&skip=0
    """
    # Build query
    query = select(ImportHistory).order_by(ImportHistory.imported_at.desc())
    
    # Apply source filter if provided
    if source:
        query = query.where(ImportHistory.source == source)
    
    # Apply pagination
    query = query.offset(skip).limit(limit)
    
    # Execute query
    history_records = session.exec(query).all()
    
    return [ImportHistoryRead.model_validate(record) for record in history_records]


@app.get("/import-history/{history_id}", response_model=ImportHistoryRead)
def get_import_history(
    history_id: int = Path(..., description="Import history record ID"),
    session: Session = Depends(get_session)
):
    """
    Get details for a specific import history record.
    
    Args:
        history_id: Import history record ID
        session: Database session
    
    Returns:
        Import history record with full details
    
    Example:
        GET /api/import-history/1
    """
    history = session.get(ImportHistory, history_id)
    
    if not history:
        raise HTTPException(
            status_code=404,
            detail=f"Import history record {history_id} not found"
        )
    
    return ImportHistoryRead.model_validate(history)


@app.delete("/import-history/{history_id}")
def delete_import_history(
    history_id: int = Path(..., description="Import history record ID"),
    session: Session = Depends(get_session)
):
    """
    Delete an import history record.
    
    Note: This only deletes the history record, not the imported templates.
    
    Args:
        history_id: Import history record ID
        session: Database session
    
    Returns:
        Success message
    
    Example:
        DELETE /api/import-history/1
    """
    history = session.get(ImportHistory, history_id)
    
    if not history:
        raise HTTPException(
            status_code=404,
            detail=f"Import history record {history_id} not found"
        )
    
    session.delete(history)
    session.commit()
    
    logger.info(f"Deleted import history record {history_id}")
    
    return {
        "message": f"Import history record {history_id} deleted successfully",
        "deleted_id": history_id
    }


@app.post("/vulnerability-templates/import-cve", response_model=VulnerabilityTemplateRead)
async def import_cve_by_id(
    cve_id: str = Query(..., description="CVE ID to import (e.g., CVE-2024-1234)"),
    overwrite_existing: bool = Query(False, description="Overwrite existing CVE template if found"),
    session: Session = Depends(get_session)
):
    """
    Import a single CVE directly from NVD API by CVE ID.
    
    Fetches CVE data from NIST National Vulnerability Database and creates
    a vulnerability template. If CVE already exists, optionally overwrites it.
    
    Args:
        cve_id: CVE identifier (e.g., CVE-2024-1234, CVE-2021-44228)
        overwrite_existing: Update existing template if CVE already imported
        session: Database session
    
    Returns:
        Created or updated VulnerabilityTemplate
    
    Example:
        POST /api/vulnerability-templates/import-cve?cve_id=CVE-2021-44228
        
    Response:
        {
          "id": 123,
          "title": "Log4j RCE Vulnerability",
          "cve_id": "CVE-2021-44228",
          "cvss_score": 10.0,
          ...
        }
    """
    from app.nvd import fetch_cve_data, NVDAPIError
    from app.models import VulnerabilityTemplate
    import time
    import json
    
    start_time = time.time()
    error_details = []
    
    # Normalize CVE ID
    cve_id = cve_id.upper().strip()
    if not cve_id.startswith('CVE-'):
        cve_id = f"CVE-{cve_id}"
    
    # Check if CVE already exists
    existing = session.exec(
        select(VulnerabilityTemplate).where(
            VulnerabilityTemplate.cve_id == cve_id
        )
    ).first()
    
    if existing and not overwrite_existing:
        raise HTTPException(
            status_code=409,
            detail=f"{cve_id} already exists in database (ID: {existing.id}). Use overwrite_existing=true to update."
        )
    
    # Fetch from NVD API
    try:
        logger.info(f"Fetching {cve_id} from NVD API...")
        cve_data = await fetch_cve_data(cve_id, use_cache=False)
        
        if not cve_data:
            raise HTTPException(
                status_code=404,
                detail=f"{cve_id} not found in NIST NVD database. Verify CVE ID is correct."
            )
    
    except NVDAPIError as e:
        logger.error(f"NVD API error for {cve_id}: {str(e)}")
        raise HTTPException(
            status_code=502,
            detail=f"Failed to fetch from NVD API: {str(e)}"
        )
    except Exception as e:
        logger.error(f"Unexpected error fetching {cve_id}: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to import CVE: {str(e)}"
        )
    
    # Create or update template
    try:
        if existing:
            # Update existing template
            for key, value in cve_data.items():
                if value is not None:
                    setattr(existing, key, value)
            existing.updated_at = get_utc_now()
            existing.source = "nvd"
            existing.is_verified = True
            session.add(existing)
            session.commit()
            session.refresh(existing)
            
            duration = time.time() - start_time
            logger.info(f"Updated existing template for {cve_id} (ID: {existing.id}) in {duration:.2f}s")
            
            # Track in import history
            try:
                import_record = ImportHistory(
                    source="nvd",
                    import_type="single_cve",
                    file_name=None,
                    file_size=None,
                    templates_created=0,
                    templates_updated=1,
                    templates_skipped=0,
                    errors=0,
                    total_parsed=1,
                    imported_by="system",  # TODO(v1.0.0): Get from current_user.email when auth is implemented
                    imported_at=get_utc_now(),
                    duration_seconds=round(duration, 2),
                    error_details=None
                )
                session.add(import_record)
                session.commit()
            except Exception as e:
                logger.warning(f"Failed to create import history: {str(e)}")
            
            return VulnerabilityTemplateRead.model_validate(existing)
        else:
            # Create new template
            template = VulnerabilityTemplate(**cve_data)
            template.created_at = get_utc_now()
            template.updated_at = get_utc_now()
            template.source = "nvd"
            template.is_verified = True
            session.add(template)
            session.commit()
            session.refresh(template)
            
            duration = time.time() - start_time
            logger.info(f"Created new template for {cve_id} (ID: {template.id}) in {duration:.2f}s")
            
            # Track in import history
            try:
                import_record = ImportHistory(
                    source="nvd",
                    import_type="single_cve",
                    file_name=None,
                    file_size=None,
                    templates_created=1,
                    templates_updated=0,
                    templates_skipped=0,
                    errors=0,
                    total_parsed=1,
                    imported_by="system",  # TODO(v1.0.0): Get from current_user.email when auth is implemented
                    imported_at=get_utc_now(),
                    duration_seconds=round(duration, 2),
                    error_details=None
                )
                session.add(import_record)
                session.commit()
            except Exception as e:
                logger.warning(f"Failed to create import history: {str(e)}")
            
            return VulnerabilityTemplateRead.model_validate(template)
    
    except Exception as e:
        session.rollback()
        logger.error(f"Database error creating template for {cve_id}: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to save CVE template: {str(e)}"
        )


# =============================================================================
# ATT&CK Mapping Endpoints (v0.7.0 Phase 2B)
# =============================================================================

@app.get("/attack/techniques")
def get_attack_techniques(
    query: str = Query(None, description="Search query for filtering techniques"),
    session: Session = Depends(get_session)
):
    """
    Get all available MITRE ATT&CK techniques or search for specific ones.
    
    GET /attack/techniques - Get all techniques
    GET /attack/techniques?query=injection - Search techniques
    
    Returns:
        List of ATT&CK techniques with details
    """
    from app.attack import get_all_techniques, search_techniques
    
    if query:
        techniques = search_techniques(query)
    else:
        techniques = get_all_techniques()
    
    return {
        "count": len(techniques),
        "techniques": techniques
    }


@app.post("/vulnerability-templates/{template_id}/suggest-attack")
def suggest_attack_techniques(
    template_id: int,
    session: Session = Depends(get_session)
):
    """
    Suggest MITRE ATT&CK techniques for a vulnerability template.
    
    POST /vulnerability-templates/5/suggest-attack
    
    Returns:
        List of suggested techniques with relevance scores
    """
    from app.attack import suggest_techniques
    
    # Get template
    template = session.get(VulnerabilityTemplate, template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    
    # Get suggestions
    suggestions = suggest_techniques(
        description=template.description,
        cwe_id=template.cwe_id,
        vulnerability_type=template.vulnerability_type
    )
    
    logger.info(f"Generated {len(suggestions)} ATT&CK technique suggestions for template {template_id}")
    
    return {
        "template_id": template_id,
        "template_title": template.title,
        "suggestion_count": len(suggestions),
        "suggestions": suggestions
    }


@app.patch("/vulnerability-templates/{template_id}/attack-techniques")
def update_attack_techniques(
    template_id: int,
    technique_ids: List[str],
    session: Session = Depends(get_session)
):
    """
    Update ATT&CK techniques for a vulnerability template.
    
    PATCH /vulnerability-templates/5/attack-techniques
    Body: ["T1190", "T1059", "T1505.003"]
    
    Returns:
        Updated template with ATT&CK mappings
    """
    from app.attack import get_all_techniques, format_techniques_for_storage
    
    # Get template
    template = session.get(VulnerabilityTemplate, template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    
    # Validate technique IDs
    all_techniques = get_all_techniques()
    valid_ids = {t['technique_id'] for t in all_techniques}
    
    selected_techniques = []
    for tid in technique_ids:
        if tid not in valid_ids:
            raise HTTPException(
                status_code=400,
                detail=f"Invalid technique ID: {tid}"
            )
        
        # Find technique details
        tech = next(t for t in all_techniques if t['technique_id'] == tid)
        selected_techniques.append(tech)
    
    # Store as JSON
    template.attack_techniques = format_techniques_for_storage(selected_techniques)
    template.updated_at = datetime.utcnow()
    
    session.add(template)
    session.commit()
    session.refresh(template)
    
    logger.info(f"Updated ATT&CK techniques for template {template_id}: {technique_ids}")
    
    return VulnerabilityTemplateRead.model_validate(template)


@app.post("/projects/{project_id}/auto-match")
def auto_match_project_findings(
    project_id: int,
    min_score: float = Query(0.6, ge=0.0, le=1.0, description="Minimum similarity score (0.0-1.0)"),
    auto_create: bool = Query(True, description="Automatically create matches above threshold"),
    session: Session = Depends(get_session)
):
    """
    Automatically match all findings in a project to vulnerability templates.
    
    Uses tiered matching strategy (v0.7.0 Phase 1):
    - Tier 1 (Exact): CWE/CVE exact matching → 100% confidence
    - Tier 2 (Fuzzy): Title/description fuzzy matching → 60-99% confidence
    - Tier 3 (Semantic): AI embeddings → future enhancement
    
    Args:
        project_id: Project to process
        min_score: Minimum similarity threshold (0.0-1.0). Default 0.6 (60%)
        auto_create: If True, automatically create VulnerabilityMatch records.
                     If False, return suggestions only (preview mode)
    
    Returns:
        {
            "project_id": int,
            "total_findings": int,
            "matched_count": int,
            "unmatched_count": int,
            "matches": [
                {
                    "finding_id": int,
                    "finding_title": str,
                    "template_id": int,
                    "template_title": str,
                    "similarity_score": float,
                    "match_method": str,
                    "created": bool  # True if auto_create=True
                }
            ]
        }
    
    Examples:
        # Preview mode (no changes)
        POST /projects/1/auto-match?auto_create=false
        
        # Auto-create matches with 70% minimum
        POST /projects/1/auto-match?min_score=0.7&auto_create=true
    """
    from app.matching import find_all_matches, create_vulnerability_match
    
    # Verify project exists
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Get all findings in project
    findings = session.exec(
        select(Finding).where(Finding.project_id == project_id)
    ).all()
    
    total_findings = len(findings)
    matched_count = 0
    results = []
    
    for finding in findings:
        # Find all potential matches for this finding
        matches = find_all_matches(session, finding, min_score=min_score)
        
        if matches:
            # Take the best match (first in list, highest score)
            template, score, method = matches[0]
            matched_count += 1
            
            created = False
            if auto_create:
                # Create VulnerabilityMatch record
                create_vulnerability_match(
                    session, finding, template, score, method, matched_by="auto_match_api"
                )
                created = True
            
            results.append({
                "finding_id": finding.id,
                "finding_title": finding.title,
                "template_id": template.id,
                "template_title": template.title,
                "similarity_score": round(score, 3),
                "match_method": method,
                "created": created
            })
    
    unmatched_count = total_findings - matched_count
    
    logger.info(
        f"Auto-match complete for project {project_id}: "
        f"{matched_count}/{total_findings} findings matched "
        f"(min_score={min_score}, auto_create={auto_create})"
    )
    
    return {
        "project_id": project_id,
        "total_findings": total_findings,
        "matched_count": matched_count,
        "unmatched_count": unmatched_count,
        "matches": results
    }


# =====================================================
# TAG MANAGEMENT ENDPOINTS
# =====================================================

@app.get("/tags", response_model=List[TagRead])
@limiter.limit("100/minute")
def list_tags(
    request: Request,
    search: Optional[str] = Query(None, description="Search tags by name"),
    session: Session = Depends(get_session)
):
    """
    List all tags with optional search.
    Returns tags sorted by usage count (most used first).
    """
    
    statement = select(Tag)
    
    if search:
        statement = statement.where(Tag.name.ilike(f"%{search}%"))
    
    statement = statement.order_by(Tag.usage_count.desc(), Tag.name)
    
    tags = session.exec(statement).all()
    
    return tags


@app.post("/tags", response_model=TagRead, status_code=201)
@limiter.limit("30/hour")
def create_tag(
    request: Request,
    tag: TagCreate,
    session: Session = Depends(get_session)
):
    """
    Create a new tag.
    Tag names must be unique (case-insensitive).
    """
    from app.timezone_utils import get_utc_now
    
    # Check if tag already exists (case-insensitive)
    existing = session.exec(
        select(Tag).where(func.lower(Tag.name) == tag.name.lower())
    ).first()
    
    if existing:
        raise HTTPException(
            status_code=400,
            detail=f"Tag '{tag.name}' already exists"
        )
    
    # Validate color format (hex color)
    if tag.color and not re.match(r'^#[0-9A-Fa-f]{6}$', tag.color):
        raise HTTPException(
            status_code=400,
            detail="Color must be a valid hex color code (e.g., #2196F3)"
        )
    
    new_tag = Tag(
        name=tag.name,
        color=tag.color or "#2196F3",
        description=tag.description,
        created_at=get_utc_now(),
        usage_count=0
    )
    
    session.add(new_tag)
    session.commit()
    session.refresh(new_tag)
    
    logger.info(f"Created tag: {new_tag.id} - {new_tag.name}")
    
    return new_tag


@app.get("/tags/{tag_id}", response_model=TagRead)
def get_tag(
    tag_id: int,
    session: Session = Depends(get_session)
):
    """Get a single tag by ID."""
    
    tag = session.get(Tag, tag_id)
    if not tag:
        raise HTTPException(status_code=404, detail="Tag not found")
    
    return tag


@app.patch("/tags/{tag_id}", response_model=TagRead)
@limiter.limit("60/hour")
def update_tag(
    request: Request,
    tag_id: int,
    tag_update: TagUpdate,
    session: Session = Depends(get_session)
):
    """
    Update a tag.
    Only provided fields will be updated.
    """
    from app.timezone_utils import get_utc_now
    
    tag = session.get(Tag, tag_id)
    if not tag:
        raise HTTPException(status_code=404, detail="Tag not found")
    
    # Check name uniqueness if being updated
    if tag_update.name and tag_update.name != tag.name:
        existing = session.exec(
            select(Tag).where(
                func.lower(Tag.name) == tag_update.name.lower(),
                Tag.id != tag_id
            )
        ).first()
        
        if existing:
            raise HTTPException(
                status_code=400,
                detail=f"Tag '{tag_update.name}' already exists"
            )
        tag.name = tag_update.name
    
    if tag_update.color:
        if not re.match(r'^#[0-9A-Fa-f]{6}$', tag_update.color):
            raise HTTPException(
                status_code=400,
                detail="Color must be a valid hex color code (e.g., #2196F3)"
            )
        tag.color = tag_update.color
    
    if tag_update.description is not None:
        tag.description = tag_update.description
    
    session.add(tag)
    session.commit()
    session.refresh(tag)
    
    logger.info(f"Updated tag: {tag_id} - {tag.name}")
    
    return tag


@app.delete("/tags/{tag_id}", status_code=204)
def delete_tag(
    tag_id: int,
    session: Session = Depends(get_session)
):
    """
    Delete a tag.
    This will also remove all associations with findings.
    """
    
    tag = session.get(Tag, tag_id)
    if not tag:
        raise HTTPException(status_code=404, detail="Tag not found")
    
    # Delete all finding-tag associations
    session.exec(delete(FindingTag).where(FindingTag.tag_id == tag_id))
    
    # Delete the tag
    session.delete(tag)
    session.commit()
    
    logger.info(f"Deleted tag: {tag_id} - {tag.name}")


@app.post("/findings/{finding_id}/tags/{tag_id}", status_code=201)
@limiter.limit("100/minute")
def add_tag_to_finding(
    request: Request,
    finding_id: int,
    tag_id: int,
    session: Session = Depends(get_session)
):
    """Add a tag to a finding."""
    from app.models import Finding
    from app.timezone_utils import get_utc_now
    
    # Verify finding exists
    finding = session.get(Finding, finding_id)
    if not finding:
        raise HTTPException(status_code=404, detail="Finding not found")
    
    # Verify tag exists
    tag = session.get(Tag, tag_id)
    if not tag:
        raise HTTPException(status_code=404, detail="Tag not found")
    
    # Check if already associated
    existing = session.exec(
        select(FindingTag).where(
            FindingTag.finding_id == finding_id,
            FindingTag.tag_id == tag_id
        )
    ).first()
    
    if existing:
        return {"message": "Tag already associated with finding"}
    
    # Create association
    finding_tag = FindingTag(
        finding_id=finding_id,
        tag_id=tag_id,
        created_at=get_utc_now()
    )
    
    session.add(finding_tag)
    
    # Update tag usage count
    tag.usage_count += 1
    session.add(tag)
    
    session.commit()
    
    logger.info(f"Added tag {tag_id} to finding {finding_id}")
    
    return {"message": "Tag added to finding"}


@app.delete("/findings/{finding_id}/tags/{tag_id}", status_code=204)
def remove_tag_from_finding(
    finding_id: int,
    tag_id: int,
    session: Session = Depends(get_session)
):
    """Remove a tag from a finding."""
    
    # Find the association
    finding_tag = session.exec(
        select(FindingTag).where(
            FindingTag.finding_id == finding_id,
            FindingTag.tag_id == tag_id
        )
    ).first()
    
    if not finding_tag:
        raise HTTPException(status_code=404, detail="Tag association not found")
    
    # Delete association
    session.delete(finding_tag)
    
    # Update tag usage count
    tag = session.get(Tag, tag_id)
    if tag and tag.usage_count > 0:
        tag.usage_count -= 1
        session.add(tag)
    
    session.commit()
    
    logger.info(f"Removed tag {tag_id} from finding {finding_id}")


@app.get("/findings/{finding_id}/tags", response_model=List[TagRead])
def get_finding_tags(
    finding_id: int,
    session: Session = Depends(get_session)
):
    """Get all tags for a finding."""
    from app.models import Finding
    
    # Verify finding exists
    finding = session.get(Finding, finding_id)
    if not finding:
        raise HTTPException(status_code=404, detail="Finding not found")
    
    # Get all tags for this finding
    tags = session.exec(
        select(Tag)
        .join(FindingTag, Tag.id == FindingTag.tag_id)
        .where(FindingTag.finding_id == finding_id)
        .order_by(Tag.name)
    ).all()
    
    return tags


# =====================================================
# REPORT TEMPLATE ENDPOINTS (v1.1.0)
# =====================================================

@app.get("/templates", response_model=List[ReportTemplateRead])
def get_templates(
    template_type: Optional[str] = Query(None, description="Filter by template type"),
    session: Session = Depends(get_session)
):
    """
    Get all report templates.
    Optionally filter by template_type.
    """
    from app.models import ReportTemplate
    
    statement = select(ReportTemplate)
    
    if template_type:
        statement = statement.where(ReportTemplate.template_type == template_type)
    
    statement = statement.order_by(ReportTemplate.is_system_template.desc(), ReportTemplate.name)
    
    templates = session.exec(statement).all()
    
    return templates


@app.post("/templates", response_model=ReportTemplateRead, status_code=201)
def create_template(
    template: ReportTemplateCreate,
    session: Session = Depends(get_session)
):
    """
    Create a new unified report template.
    Supports both simple parameterized AND complex widget-based templates.
    """
    from app.models import ReportTemplate
    from app.timezone_utils import get_utc_now
    
    # Validate JSON fields
    try:
        json.loads(template.sections)
    except json.JSONDecodeError:
        raise HTTPException(status_code=400, detail="Invalid JSON in sections field")
    
    try:
        json.loads(template.variables)
    except json.JSONDecodeError:
        raise HTTPException(status_code=400, detail="Invalid JSON in variables field")
    
    if template.layout_config:
        try:
            json.loads(template.layout_config)
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="Invalid JSON in layout_config field")
    
    # Create template
    now = get_utc_now()
    db_template = ReportTemplate(
        name=template.name,
        description=template.description,
        template_type=template.template_type,
        sections=template.sections,
        variables=template.variables,
        layout_config=template.layout_config,
        is_system_template=False,
        is_public=template.is_public,
        usage_count=0,
        last_used_at=None,
        created_at=now,
        updated_at=now,
        created_by_user_id=None  # No auth required yet
    )
    
    session.add(db_template)
    session.commit()
    session.refresh(db_template)
    
    logger.info(f"Created report template: {db_template.id} - {db_template.name}")
    
    return db_template


@app.get("/templates/{template_id}", response_model=ReportTemplateRead)
def get_template(
    template_id: int,
    session: Session = Depends(get_session)
):
    """Get a specific report template by ID."""
    from app.models import ReportTemplate
    
    template = session.get(ReportTemplate, template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    
    return template


@app.put("/templates/{template_id}", response_model=ReportTemplateRead)
def update_template(
    template_id: int,
    template_update: ReportTemplateUpdate,
    session: Session = Depends(get_session)
):
    """
    Update a report template.
    System templates cannot be modified.
    """
    from app.models import ReportTemplate
    from app.timezone_utils import get_utc_now
    
    db_template = session.get(ReportTemplate, template_id)
    if not db_template:
        raise HTTPException(status_code=404, detail="Template not found")
    
    # Prevent modification of system templates
    if db_template.is_system_template:
        raise HTTPException(status_code=403, detail="Cannot modify system templates")
    
    # Update fields
    if template_update.name is not None:
        db_template.name = template_update.name
    if template_update.description is not None:
        db_template.description = template_update.description
    if template_update.template_type is not None:
        db_template.template_type = template_update.template_type
    if template_update.sections is not None:
        try:
            json.loads(template_update.sections)
            db_template.sections = template_update.sections
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="Invalid JSON in sections field")
    if template_update.variables is not None:
        try:
            json.loads(template_update.variables)
            db_template.variables = template_update.variables
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="Invalid JSON in variables field")
    
    db_template.updated_at = get_utc_now()
    
    session.add(db_template)
    session.commit()
    session.refresh(db_template)
    
    logger.info(f"Updated report template: {template_id}")
    
    return db_template


@app.delete("/templates/{template_id}", status_code=204)
def delete_template(
    template_id: int,
    session: Session = Depends(get_session)
):
    """
    Delete a report template.
    System templates cannot be deleted.
    """
    from app.models import ReportTemplate
    
    template = session.get(ReportTemplate, template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    
    # Prevent deletion of system templates
    if template.is_system_template:
        raise HTTPException(status_code=403, detail="Cannot delete system templates")
    
    session.delete(template)
    session.commit()
    
    logger.info(f"Deleted report template: {template_id} - {template.name}")


# =====================================================
# QUICK ADD / TEMPLATE SEARCH ENDPOINTS
# =====================================================

@app.get("/repository/search", response_model=List[VulnerabilityTemplateRead])
def search_repository_templates(
    q: str = Query(..., min_length=1, description="Search query for title, description, CWE, CVE"),
    limit: int = Query(20, ge=1, le=50, description="Maximum results to return"),
    verified_only: bool = Query(False, description="Only return verified templates"),
    session: Session = Depends(get_session)
):
    """
    Quick search for vulnerability templates (optimized for Quick Add feature).
    
    Searches across:
    - Title (fuzzy match)
    - Description (fuzzy match)
    - CWE ID (exact match)
    - CVE ID (exact match)
    - Vulnerability type (exact match)
    
    Returns templates ordered by:
    1. Exact title matches first
    2. Usage count (most used first)
    3. Creation date (newest first)
    """
    from app.models import VulnerabilityTemplate
    
    # Build search pattern for fuzzy matching
    search_pattern = f"%{q}%"
    
    # Base query
    statement = select(VulnerabilityTemplate)
    
    # Apply verified filter
    if verified_only:
        statement = statement.where(VulnerabilityTemplate.is_verified == True)
    
    # Search across multiple fields
    statement = statement.where(
        (VulnerabilityTemplate.title.ilike(search_pattern)) |
        (VulnerabilityTemplate.description.ilike(search_pattern)) |
        (VulnerabilityTemplate.cwe_id.ilike(search_pattern)) |
        (VulnerabilityTemplate.cve_id.ilike(search_pattern)) |
        (VulnerabilityTemplate.vulnerability_type.ilike(search_pattern))
    )
    
    # Order by relevance:
    # 1. Exact title match first (case-insensitive)
    # 2. Usage count (most used first)
    # 3. Creation date (newest first)
    
    # Create a score for exact matches
    exact_match_score = case(
        (func.lower(VulnerabilityTemplate.title) == q.lower(), 1),
        else_=0
    )
    
    statement = statement.order_by(
        exact_match_score.desc(),
        VulnerabilityTemplate.usage_count.desc(),
        VulnerabilityTemplate.created_at.desc()
    ).limit(limit)
    
    templates = session.exec(statement).all()
    
    logger.info(f"Repository search for '{q}' returned {len(templates)} results")
    
    return templates


@app.get("/projects/{project_id}/template-suggestions", response_model=List[VulnerabilityTemplateRead])
def get_project_template_suggestions(
    project_id: int,
    limit: int = Query(10, ge=1, le=50, description="Maximum suggestions to return"),
    session: Session = Depends(get_session)
):
    """
    Get template suggestions for Quick Add based on project's existing findings.
    
    Returns templates that:
    1. Are already used in this project (sorted by frequency)
    2. Have high overall usage counts (popular templates)
    3. Are verified templates
    
    Useful for "Add Similar" and suggesting common vulnerabilities.
    """
    from app.models import VulnerabilityTemplate, Finding, Project
    
    # Verify project exists
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Get templates already used in this project
    project_templates_statement = (
        select(VulnerabilityTemplate.id, func.count(Finding.id).label("usage_in_project"))
        .join(Finding, Finding.template_id == VulnerabilityTemplate.id)
        .where(Finding.project_id == project_id)
        .group_by(VulnerabilityTemplate.id)
        .order_by(func.count(Finding.id).desc())
    )
    
    project_template_usage = session.exec(project_templates_statement).all()
    project_template_ids = [t[0] for t in project_template_usage]
    
    # Combine:
    # 1. Templates used in this project (by frequency)
    # 2. Most popular verified templates not yet used
    suggestions = []
    
    # Add templates from this project first
    if project_template_ids:
        project_templates = session.exec(
            select(VulnerabilityTemplate)
            .where(VulnerabilityTemplate.id.in_(project_template_ids))
            .order_by(VulnerabilityTemplate.usage_count.desc())
        ).all()
        suggestions.extend(project_templates)
    
    # Fill remaining slots with popular verified templates
    remaining = limit - len(suggestions)
    if remaining > 0:
        popular_templates = session.exec(
            select(VulnerabilityTemplate)
            .where(
                VulnerabilityTemplate.is_verified == True,
                ~VulnerabilityTemplate.id.in_(project_template_ids) if project_template_ids else True
            )
            .order_by(VulnerabilityTemplate.usage_count.desc())
            .limit(remaining)
        ).all()
        suggestions.extend(popular_templates)
    
    logger.info(f"Generated {len(suggestions)} template suggestions for project {project_id}")
    
    return suggestions[:limit]


# ============================================================================
# EXECUTIVE DASHBOARD ENDPOINTS
# ============================================================================

from app.executive import ExecutiveMetrics


@app.get("/executive/summary")
def get_executive_summary(session: Session = Depends(get_session)):
    """
    Executive Summary Dashboard - High-level KPIs for C-level stakeholders.
    
    Returns:
    - Total projects (active)
    - Total findings with breakdown by severity
    - MTTR (Mean Time To Remediation) in days
    - Trend direction (improving/worsening/stable)
    - Compliance coverage (OWASP, CWE, ATT&CK)
    - Open critical/high findings count
    - Top 5 risky projects
    """
    try:
        summary = ExecutiveMetrics.get_executive_summary(session)
        logger.info("Generated executive summary")
        return summary
    except Exception as e:
        logger.error(f"Error generating executive summary: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate executive summary: {str(e)}")


@app.get("/executive/risk-heatmap")
def get_risk_heatmap(session: Session = Depends(get_session)):
    """
    Risk Heat Map - Visual grid showing risk scores across all projects.
    
    Returns list of projects with:
    - project_id, project_name
    - risk_score (weighted by severity)
    - severity_counts (critical, high, medium, low, informational)
    - color coding (red/orange/yellow/green)
    - total_findings, open_critical_high count
    
    Projects are sorted by risk_score descending.
    """
    try:
        heat_map = ExecutiveMetrics.get_risk_heat_map(session)
        logger.info(f"Generated risk heat map for {len(heat_map)} projects")
        return heat_map
    except Exception as e:
        logger.error(f"Error generating risk heat map: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate risk heat map: {str(e)}")


# ============================================================================
# ADVANCED REPORTING ENDPOINTS
# ============================================================================

from app.report_templates import ReportTemplateEngine
from app.email_service import EmailService
from app.models import (
    EmailSettings,
    EmailSettingsRead,
    EmailSettingsCreate,
    EmailSettingsUpdate,
    ReportBranding,
    ReportBrandingRead,
    ReportBrandingUpdate,
    ReportGenerationRequest,
    ReportTemplateType,
    ReportFormat,
    CustomReportTemplate,
    CustomReportTemplateRead,
    CustomReportTemplateCreate,
    CustomReportTemplateUpdate
)


# --- Custom Report Template CRUD ---

@app.post("/custom-templates", response_model=CustomReportTemplateRead)
@limiter.limit("30/minute")
def create_custom_template(
    request: Request,
    template: CustomReportTemplateCreate,
    session: Session = Depends(get_session)
):
    """
    Create a new custom report template.
    
    Template JSON structure should follow this schema:
    {
        "sections": [
            {
                "type": "text|table|chart|metrics|findings",
                "title": "Section Title",
                "content": "...",  # For text sections
                "widget": "...",   # For chart/metrics widgets
                "filters": {...},  # For data sections
                "layout": {...}    # Layout options
            }
        ],
        "layout": {
            "page_size": "letter|a4",
            "orientation": "portrait|landscape",
            "margins": {...}
        }
    }
    """
    try:
        # Validate JSON structure
        import json
        template_data = json.loads(template.template_json)
        
        if "sections" not in template_data:
            raise HTTPException(
                status_code=400,
                detail="Template JSON must contain 'sections' array"
            )
        
        # Create template
        db_template = CustomReportTemplate(
            name=template.name,
            description=template.description,
            template_json=template.template_json,
            is_public=template.is_public,
            created_by=template.created_by,
            created_at=get_utc_now(),
            updated_at=get_utc_now(),
            usage_count=0
        )
        
        session.add(db_template)
        session.commit()
        session.refresh(db_template)
        
        logger.info(f"Created custom template: {db_template.name} (ID: {db_template.id})")
        return db_template
        
    except json.JSONDecodeError as e:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid JSON in template_json: {str(e)}"
        )
    except Exception as e:
        logger.error(f"Error creating custom template: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to create template: {str(e)}")


@app.get("/custom-templates", response_model=List[CustomReportTemplateRead])
@limiter.limit("60/minute")
def list_custom_templates(
    request: Request,
    skip: int = Query(0, ge=0),
    limit: int = Query(50, ge=1, le=100),
    search: Optional[str] = Query(None),
    is_public: Optional[bool] = Query(None),
    created_by: Optional[str] = Query(None),
    session: Session = Depends(get_session)
):
    """
    List all custom report templates with optional filtering.
    """
    try:
        query = select(CustomReportTemplate)
        
        # Apply filters
        if search:
            query = query.where(
                CustomReportTemplate.name.ilike(f"%{search}%") |
                CustomReportTemplate.description.ilike(f"%{search}%")
            )
        
        if is_public is not None:
            query = query.where(CustomReportTemplate.is_public == is_public)
        
        if created_by:
            query = query.where(CustomReportTemplate.created_by == created_by)
        
        # Order by most recently used, then by name
        query = query.order_by(
            CustomReportTemplate.last_used_at.desc().nullslast(),
            CustomReportTemplate.name
        )
        
        # Apply pagination
        query = query.offset(skip).limit(limit)
        
        templates = session.exec(query).all()
        return templates
        
    except Exception as e:
        logger.error(f"Error listing custom templates: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to list templates: {str(e)}")


@app.get("/custom-templates/{template_id}", response_model=CustomReportTemplateRead)
@limiter.limit("60/minute")
def get_custom_template(
    request: Request,
    template_id: int,
    session: Session = Depends(get_session)
):
    """Get a specific custom report template by ID."""
    template = session.get(CustomReportTemplate, template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    return template


@app.patch("/custom-templates/{template_id}", response_model=CustomReportTemplateRead)
@limiter.limit("30/minute")
def update_custom_template(
    request: Request,
    template_id: int,
    updates: CustomReportTemplateUpdate,
    session: Session = Depends(get_session)
):
    """Update a custom report template."""
    try:
        template = session.get(CustomReportTemplate, template_id)
        if not template:
            raise HTTPException(status_code=404, detail="Template not found")
        
        # Validate JSON if provided
        if updates.template_json:
            import json
            template_data = json.loads(updates.template_json)
            if "sections" not in template_data:
                raise HTTPException(
                    status_code=400,
                    detail="Template JSON must contain 'sections' array"
                )
        
        # Apply updates
        update_data = updates.model_dump(exclude_unset=True)
        for field, value in update_data.items():
            setattr(template, field, value)
        
        template.updated_at = get_utc_now()
        
        session.add(template)
        session.commit()
        session.refresh(template)
        
        logger.info(f"Updated custom template: {template.name} (ID: {template_id})")
        return template
        
    except json.JSONDecodeError as e:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid JSON in template_json: {str(e)}"
        )
    except Exception as e:
        logger.error(f"Error updating custom template: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to update template: {str(e)}")


@app.delete("/custom-templates/{template_id}", status_code=204)
@limiter.limit("30/minute")
def delete_custom_template(
    request: Request,
    template_id: int,
    session: Session = Depends(get_session)
):
    """Delete a custom report template."""
    try:
        template = session.get(CustomReportTemplate, template_id)
        if not template:
            raise HTTPException(status_code=404, detail="Template not found")
        
        session.delete(template)
        session.commit()
        
        logger.info(f"Deleted custom template: {template.name} (ID: {template_id})")
        return None
        
    except Exception as e:
        logger.error(f"Error deleting custom template: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to delete template: {str(e)}")


@app.post("/custom-templates/{template_id}/duplicate", response_model=CustomReportTemplateRead)
@limiter.limit("30/minute")
def duplicate_custom_template(
    request: Request,
    template_id: int,
    new_name: Optional[str] = Query(None),
    session: Session = Depends(get_session)
):
    """Duplicate an existing custom report template."""
    try:
        original = session.get(CustomReportTemplate, template_id)
        if not original:
            raise HTTPException(status_code=404, detail="Template not found")
        
        # Create duplicate
        duplicate_name = new_name or f"{original.name} (Copy)"
        
        duplicate = CustomReportTemplate(
            name=duplicate_name,
            description=original.description,
            template_json=original.template_json,
            is_public=False,  # Always private by default
            created_by=original.created_by,
            created_at=get_utc_now(),
            updated_at=get_utc_now(),
            usage_count=0
        )
        
        session.add(duplicate)
        session.commit()
        session.refresh(duplicate)
        
        logger.info(f"Duplicated template {template_id} as {duplicate.id}: {duplicate.name}")
        return duplicate
        
    except Exception as e:
        logger.error(f"Error duplicating custom template: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to duplicate template: {str(e)}")


@app.post("/reports/generate")
def generate_advanced_report(
    request: ReportGenerationRequest,
    session: Session = Depends(get_session)
):
    """
    Generate advanced report with template selection and multi-format support.
    
    Supports:
    - Multiple template types (Executive Summary, Technical Findings, etc.)
    - Multiple formats (DOCX, PDF, HTML)
    - Project filtering
    - Date range filtering
    - Optional email delivery
    - Branding customization
    """
    try:
        logger.info(f"Generating report: {request.template_type} in {request.format}")
        
        # Get branding settings
        branding = session.exec(select(ReportBranding)).first()
        
        # Create template engine
        engine = ReportTemplateEngine(session, branding)
        
        # Generate report
        file_path = engine.generate_report(
            template_type=request.template_type,
            format=request.format,
            project_ids=request.project_ids,
            start_date=request.start_date,
            end_date=request.end_date,
            include_sections=request.include_sections,
            custom_template_id=request.custom_template_id
        )
        
        # Handle email delivery if requested
        if request.send_email and request.email_to:
            email_service = EmailService(session)
            
            # Get project names for email
            if request.project_ids:
                projects = session.exec(
                    select(Project).where(Project.id.in_(request.project_ids))
                ).all()
                project_names = [p.name for p in projects]
            else:
                project_names = ["All Projects"]
            
            # Generate email body
            plain_text, html_body = email_service.generate_report_email_body(
                report_name=f"{request.template_type.value} Report",
                project_names=project_names
            )
            
            # Send email
            email_sent = email_service.send_report(
                to_emails=request.email_to,
                subject=request.email_subject or f"{request.template_type.value} Report - {datetime.now().strftime('%Y-%m-%d')}",
                body_text=plain_text,
                body_html=html_body,
                attachment_paths=[file_path],
                cc_emails=request.email_cc,
                bcc_emails=request.email_bcc
            )
            
            if email_sent:
                logger.info(f"Report emailed successfully to {', '.join(request.email_to)}")
                return {
                    "success": True,
                    "message": "Report generated and emailed successfully",
                    "file_path": file_path,
                    "email_sent": True
                }
            else:
                logger.warning("Report generated but email delivery failed")
                return {
                    "success": True,
                    "message": "Report generated but email delivery failed",
                    "file_path": file_path,
                    "email_sent": False
                }
        
        # Return file download if no email
        logger.info(f"Report generated: {file_path}")
        return FileResponse(
            path=file_path,
            media_type="application/octet-stream",
            filename=os.path.basename(file_path)
        )
    
    except Exception as e:
        logger.error(f"Error generating report: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate report: {str(e)}")


# Email Settings Endpoints

@app.get("/settings/email", response_model=Optional[EmailSettingsRead])
def get_email_settings(session: Session = Depends(get_session)):
    """Get active email settings."""
    settings = session.exec(select(EmailSettings).where(EmailSettings.is_active == True)).first()
    return settings


@app.post("/settings/email", response_model=EmailSettingsRead)
def create_email_settings(
    settings: EmailSettingsCreate,
    session: Session = Depends(get_session)
):
    """Create new email settings (deactivates others)."""
    # Deactivate all existing settings
    existing = session.exec(select(EmailSettings)).all()
    for s in existing:
        s.is_active = False
        session.add(s)
    
    # Create new settings
    new_settings = EmailSettings(
        **settings.model_dump(),
        created_at=get_utc_now(),
        updated_at=get_utc_now()
    )
    session.add(new_settings)
    session.commit()
    session.refresh(new_settings)
    
    logger.info("Email settings created successfully")
    return new_settings


@app.put("/settings/email/{settings_id}", response_model=EmailSettingsRead)
def update_email_settings(
    settings_id: int,
    settings: EmailSettingsUpdate,
    session: Session = Depends(get_session)
):
    """Update existing email settings."""
    db_settings = session.get(EmailSettings, settings_id)
    if not db_settings:
        raise HTTPException(status_code=404, detail="Email settings not found")
    
    update_data = settings.model_dump(exclude_unset=True)
    for key, value in update_data.items():
        setattr(db_settings, key, value)
    
    db_settings.updated_at = get_utc_now()
    session.add(db_settings)
    session.commit()
    session.refresh(db_settings)
    
    logger.info(f"Email settings {settings_id} updated")
    return db_settings


@app.post("/settings/email/test")
def test_email_connection(session: Session = Depends(get_session)):
    """Test email SMTP connection."""
    email_service = EmailService(session)
    result = email_service.test_connection()
    
    if not result["success"]:
        raise HTTPException(status_code=400, detail=result["message"])
    
    return result


# Report Branding Endpoints

@app.get("/settings/branding", response_model=Optional[ReportBrandingRead])
def get_branding_settings(session: Session = Depends(get_session)):
    """Get report branding settings."""
    branding = session.exec(select(ReportBranding)).first()
    return branding


@app.post("/settings/branding", response_model=ReportBrandingRead)
def create_branding_settings(
    branding: ReportBrandingUpdate,
    session: Session = Depends(get_session)
):
    """Create or update branding settings."""
    existing = session.exec(select(ReportBranding)).first()
    
    if existing:
        # Update existing
        update_data = branding.model_dump(exclude_unset=True)
        for key, value in update_data.items():
            setattr(existing, key, value)
        existing.updated_at = get_utc_now()
        session.add(existing)
        session.commit()
        session.refresh(existing)
        logger.info("Branding settings updated")
        return existing
    else:
        # Create new
        new_branding = ReportBranding(
            **branding.model_dump(),
            created_at=get_utc_now(),
            updated_at=get_utc_now()
        )
        session.add(new_branding)
        session.commit()
        session.refresh(new_branding)
        logger.info("Branding settings created")
        return new_branding


# ============================================================================
# SCORING CALCULATORS ENDPOINTS
# ============================================================================

class CVSSCalculateRequest(BaseModel):
    """Request model for CVSS calculation."""
    vector: str


class CVSSCalculateResponse(BaseModel):
    """Response model for CVSS calculation."""
    vector: str
    base_score: float
    severity: str
    is_valid: bool
    error: Optional[str] = None


@app.post("/cvss/calculate", response_model=CVSSCalculateResponse)
def calculate_cvss(
    request: CVSSCalculateRequest,
    session: Session = Depends(get_session)
):
    """
    Calculate CVSS 3.1 Base Score from a vector string.
    
    Uses the official CVSS 3.1 formula:
    - Impact Sub-Score (ISS) calculation
    - Scope-adjusted Impact
    - Exploitability calculation
    - Final Base Score with proper rounding
    
    Returns score, severity rating, and validation status.
    """
    vector = request.vector.strip()
    
    # Validate vector
    is_valid, error_msg = scoring.validate_cvss_vector(vector)
    if not is_valid:
        return CVSSCalculateResponse(
            vector=vector,
            base_score=0.0,
            severity="None",
            is_valid=False,
            error=error_msg
        )
    
    # Calculate score
    result = scoring.calculate_cvss_score(vector)
    if result is None:
        return CVSSCalculateResponse(
            vector=vector,
            base_score=0.0,
            severity="None",
            is_valid=False,
            error="Failed to calculate CVSS score"
        )
    
    base_score, severity = result
    
    logger.info(f"CVSS calculation: {vector} -> {base_score} ({severity})")
    
    return CVSSCalculateResponse(
        vector=vector,
        base_score=base_score,
        severity=severity,
        is_valid=True
    )


class OWASPCalculateRequest(BaseModel):
    """Request model for OWASP risk calculation."""
    likelihood: int
    impact: int


class OWASPCalculateResponse(BaseModel):
    """Response model for OWASP risk calculation."""
    likelihood: int
    impact: int
    risk_score: int
    risk_rating: str
    is_valid: bool
    error: Optional[str] = None


@app.post("/owasp/calculate", response_model=OWASPCalculateResponse)
def calculate_owasp(
    request: OWASPCalculateRequest,
    session: Session = Depends(get_session)
):
    """
    Calculate OWASP Risk Rating using Likelihood × Impact methodology.
    
    Inputs:
    - likelihood: 1-9 scale (1=rare, 9=almost certain)
    - impact: 1-9 scale (1=minimal, 9=catastrophic)
    
    Risk Rating Thresholds:
    - Critical: >= 18
    - High: 12-17
    - Medium: 6-11
    - Low: < 6
    """
    try:
        risk_score, risk_rating = scoring.calculate_owasp_risk(
            request.likelihood,
            request.impact
        )
        
        logger.info(f"OWASP calculation: L={request.likelihood} × I={request.impact} -> {risk_score} ({risk_rating})")
        
        return OWASPCalculateResponse(
            likelihood=request.likelihood,
            impact=request.impact,
            risk_score=risk_score,
            risk_rating=risk_rating,
            is_valid=True
        )
    
    except ValueError as e:
        return OWASPCalculateResponse(
            likelihood=request.likelihood,
            impact=request.impact,
            risk_score=0,
            risk_rating="Low",
            is_valid=False,
            error=str(e)
        )

    
    return None

